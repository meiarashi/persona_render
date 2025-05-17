from flask import Flask, request, jsonify, send_from_directory, send_file, Response
from flask_cors import CORS
import random
import os
import traceback
# --- Import DB and AI Libraries ---
# from flask_sqlalchemy import SQLAlchemy # Removed SQLAlchemy import
from openai import OpenAI
import google.generativeai as genai
from anthropic import Anthropic
from fpdf import FPDF
import io
from urllib.parse import quote
from pptx import Presentation
from pptx.util import Inches, Pt # Inches for image size, Pt for font size
from pptx.enum.text import PP_ALIGN # For text alignment
from pptx.dml.color import RGBColor

# --- Flask App Setup ---
app = Flask(__name__, static_url_path='', static_folder='.')
# app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///settings.db' # Removed DB config
# app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False # Removed DB config
# db = SQLAlchemy(app) # Removed SQLAlchemy initialization
CORS(app)

# --- Database Model Definition --- (Removed)
# class Setting(db.Model):
#     id = db.Column(db.Integer, primary_key=True)
#     key = db.Column(db.String(80), unique=True, nullable=False)
#     value = db.Column(db.String(500), nullable=True)
#
#     def __repr__(self):
#         return f'<Setting {self.key}: {self.value}>'

# --- Helper Functions --- (Removed)
# def get_setting(key, default=None):
#     """Gets a setting value from the database."""
#     setting = Setting.query.filter_by(key=key).first()
#     return setting.value if setting else default
#
# def set_setting(key, value):
#     """Sets or updates a setting value in the database."""
#     setting = Setting.query.filter_by(key=key).first()
#     if setting:
#         setting.value = value
#     else:
#         setting = Setting(key=key, value=value)
#         db.session.add(setting)
#     db.session.commit()

# --- Function to initialize DB and load initial settings --- (Removed)
# def initialize_database():
#     with app.app_context():
#         db.create_all() # Create table if it doesn't exist
#
#         # Load initial DEFAULTS for non-sensitive settings if not already in DB
#         if not get_setting("selected_text_model"):
#              set_setting("selected_text_model", "gpt-4.1")
#         if not get_setting("selected_image_model"):
#              set_setting("selected_image_model", "dall-e-3")
#         # Load initial char limits if not already in DB
#         default_limits = {
#             "personality": "100", "reason": "100", "behavior": "100",
#             "reviews": "100", "values": "100", "demands": "100"
#         }
#         for key, default_value in default_limits.items():
#             db_key = f"limit_{key}"
#             if not get_setting(db_key):
#                  set_setting(db_key, default_value)
#
#         # Add default output settings
#         if get_setting("output_pdf_enabled") is None: # Check if None explicitly
#              set_setting("output_pdf_enabled", "true") # Store as string 'true'/'false'
#         if get_setting("output_ppt_enabled") is None:
#              set_setting("output_ppt_enabled", "true")
#         if get_setting("output_gslide_enabled") is None:
#              set_setting("output_gslide_enabled", "false") # Default GSlide to disabled

# --- AI Client Initialization Helper --- 
def get_ai_client(model_name, api_key):
    """Initializes and returns the correct AI client based on model name."""
    if model_name.startswith("gpt"):
        if not api_key:
            raise ValueError("OpenAI APIキーが設定されていません。")
        try:
            return OpenAI(api_key=api_key)
        except Exception as e:
            raise ValueError(f"OpenAI Client Error: {e}")
    elif model_name.startswith("claude"):
        if not api_key:
            raise ValueError("Anthropic APIキーが設定されていません。")
        try:
            return Anthropic(api_key=api_key)
        except Exception as e:
            raise ValueError(f"Anthropic Client Error: {e}")
    elif model_name.startswith("gemini"):
        if not api_key:
             raise ValueError("Google APIキーが設定されていません。")
        try:
             genai.configure(api_key=api_key) # Configure library
             # Return the configured model object directly for gemini
             return genai.GenerativeModel('gemini-1.5-pro-latest') 
        except Exception as e:
             raise ValueError(f"Google Client Error: {e}")
    else:
        raise ValueError(f"未対応のモデルです: {model_name}")

# --- Prompt and Response Parsing Helpers --- 
def build_prompt(data):
    limit_personality = os.environ.get("LIMIT_PERSONALITY", "100")
    limit_reason = os.environ.get("LIMIT_REASON", "100")
    limit_behavior = os.environ.get("LIMIT_BEHAVIOR", "100")
    limit_reviews = os.environ.get("LIMIT_REVIEWS", "100")
    limit_values = os.environ.get("LIMIT_VALUES", "100")
    limit_demands = os.environ.get("LIMIT_DEMANDS", "100")

    # --- Patient Type Descriptions (Copied from script.js) ---
    patient_type_details_map = {
        '利便性重視型': { "description": 'アクセスの良さ、待ち時間の短さ、診療時間の柔軟性など、便利さを最優先', "example": '忙しいビジネスパーソン、オンライン診療を好む患者' },
        '専門医療追求型': { "description": '専門医や高度専門医療機関での治療を希望し、医師の経歴や実績を重視', "example": '難病患者、複雑な症状を持つ患者' },
        '予防健康管理型': { "description": '病気になる前の予防や早期発見、健康維持に関心が高い', "example": '定期健診を欠かさない人、予防接種に積極的な人' },
        '代替医療志向型': { "description": '漢方、鍼灸、ホメオパシーなど、西洋医学以外の選択肢を積極的に取り入れる', "example": '自然療法愛好者、慢性疾患の患者' },
        '経済合理型': { "description": '自己負担額、保険適用の有無、費用対効果を重視', "example": '経済的制約のある患者、医療費控除を意識する人' },
        '情報探求型': { "description": '徹底的な情報収集、セカンドオピニオン取得、比較検討を行う', "example": '高学歴層、慎重な意思決定を好む患者' },
        '革新技術指向型': { "description": '最先端の医療技術、新薬、臨床試験などに積極的に関心を持つ', "example": '既存治療で効果が出なかった患者、医療イノベーションに関心がある人' },
        '対話重視型': { "description": '医師からの丁寧な説明や対話を求め、質問が多い', "example": '不安を感じやすい患者、医療従事者' },
        '信頼基盤型': { "description": 'かかりつけ医との長期的な関係や医療機関の評判を重視', "example": '地域密着型の患者、同じ医師に長期通院する患者' },
        '緊急解決型': { "description": '症状の即時改善を求め、緊急性を重視', "example": '急性疾患患者、痛みに耐性が低い患者' },
        '受動依存型': { "description": '医師の判断に全面的に依存し、自らの決定より医師の指示を優先', "example": '高齢者、医療知識が少ない患者' },
        '自律決定型': { "description": '自分の治療に主体的に関わり、最終決定権を持ちたいと考える', "example": '医療リテラシーが高い患者、自己管理を好む慢性疾患患者' }
    }

    prompt_parts = [
        "以下の情報に基づいて、医療系のペルソナを作成してください。",
        "各項目は自然な文章で記述し、**日本語で、指定されたおおよその文字数制限に従ってください**。",
        "",
        "# 利用者からの入力情報"
    ]

    # 基本情報
    basic_info = {
        "診療科": data.get('department'),
        "ペルソナ作成目的": data.get('purpose'),
        "名前": data.get('name'),
        "性別": data.get('gender'),
        "年齢": data.get('age'),
        "都道府県": data.get('prefecture'),
        "市区町村": data.get('municipality'),
        "家族構成": data.get('family'),
        "職業": data.get('occupation'),
        "年収": data.get('income'),
        "趣味": data.get('hobby'),
        "ライフイベント": data.get('life_events'),
        "患者タイプ": data.get('patient_type')
    }
    # setting_type と patient_type も追加
    if data.get('setting_type'):
        basic_info["基本情報設定タイプ"] = data.get('setting_type')
        # 以下の重複する可能性のある「選択された患者タイプ」の追加はコメントアウトまたは削除
        # if data.get('setting_type') == 'patient_type' and data.get('patient_type'):
        #     basic_info["選択された患者タイプ"] = data.get('patient_type')


    for key, value in basic_info.items():
        if value: # 値が存在する場合のみプロンプトに追加
            prompt_parts.append(f"- {key}: {value}")
            if key == "患者タイプ" and value in patient_type_details_map:
                details = patient_type_details_map[value]
                prompt_parts.append(f"  - 患者タイプの特徴: {details['description']}")
                prompt_parts.append(f"  - 患者タイプの例: {details['example']}")
        # 「選択された患者タイプ」の(自動生成)条件から patient_type を除外 (「患者タイプ」で処理されるため)
        elif key in ["名前", "性別", "年齢", "都道府県", "市区町村", "家族構成", "職業", "年収", "趣味", "ライフイベント"] and key != "患者タイプ": 
             prompt_parts.append(f"- {key}: (自動生成)")
        elif key == "患者タイプ" and not value and data.get('setting_type') == 'patient_type': # setting_typeがpatient_typeで患者タイプが未選択の場合
             prompt_parts.append(f"- {key}: (指定なし/自動生成)")
        # '診療科', '作成目的', '基本情報設定タイプ' は値がなければ「指定なし」とするか、キー自体含めないか検討。ここでは値がある場合のみ表示。

    # Step 4の固定追加項目
    fixed_additional_info = {
        "座右の銘": data.get('motto'),
        "最近の悩み/関心": data.get('concerns'),
        "好きな有名人/尊敬する人物": data.get('favorite_person'),
        "よく見るメディア/SNS": data.get('media_sns'),
        "性格キーワード": data.get('personality_keywords'),
        "最近した健康に関する行動": data.get('health_actions'),
        "休日の過ごし方": data.get('holiday_activities'),
        "キャッチコピー": data.get('catchphrase'),
    }
    has_fixed_additional_info = any(fixed_additional_info.values())
    if has_fixed_additional_info:
        prompt_parts.append("\n## 追加情報（固定項目）")
        for key, value in fixed_additional_info.items():
            if value:
                prompt_parts.append(f"- {key}: {value}")

    # Step 4の動的追加項目
    additional_field_names = data.get('additional_field_name', [])
    additional_field_values = data.get('additional_field_value', [])
    
    dynamic_additional_info = []
    if isinstance(additional_field_names, list) and isinstance(additional_field_values, list):
        for i in range(len(additional_field_names)):
            field_name = additional_field_names[i]
            field_value = additional_field_values[i] if i < len(additional_field_values) else None
            if field_name and field_value: # 項目名と内容の両方がある場合のみ
                dynamic_additional_info.append(f"- {field_name}: {field_value}")
    
    if dynamic_additional_info:
        prompt_parts.append("\n## 追加情報（自由入力項目）")
        prompt_parts.extend(dynamic_additional_info)

    prompt_parts.append("\n# 生成項目 (日本語で、指定文字数程度で)")
    prompt_parts.append("以下の項目について、上記情報に基づいた自然な文章を生成してください。")
    prompt_parts.append(f"\n1.  **性格（価値観・人生観）**: (日本語で{limit_personality}文字程度)")
    prompt_parts.append(f"2.  **通院理由**: (日本語で{limit_reason}文字程度)")
    prompt_parts.append(f"3.  **症状通院頻度・行動パターン**: (日本語で{limit_behavior}文字程度)")
    prompt_parts.append(f"4.  **口コミの重視ポイント**: (日本語で{limit_reviews}文字程度)")
    prompt_parts.append(f"5.  **医療機関への価値観・行動傾向**: (日本語で{limit_values}文字程度)")
    prompt_parts.append(f"6.  **医療機関に求めるもの**: (日本語で{limit_demands}文字程度)")
    
    prompt_parts.append("\n回答は、各生成項目の見出しを含めず、項目に対応するテキストのみを以下の形式で記述してください。")
    prompt_parts.append("---性格---")
    prompt_parts.append("[ここに性格のテキスト]")
    prompt_parts.append("---通院理由---")
    prompt_parts.append("[ここに通院理由のテキスト]")
    prompt_parts.append("---行動パターン---")
    prompt_parts.append("[ここに症状通院頻度・行動パターンのテキスト]")
    prompt_parts.append("---口コミ重視点---")
    prompt_parts.append("[ここに口コミの重視ポイントのテキスト]")
    prompt_parts.append("---医療価値観---")
    prompt_parts.append("[ここに医療機関への価値観・行動傾向のテキスト]")
    prompt_parts.append("---医療求めるもの---")
    prompt_parts.append("[ここに医療機関に求めるもののテキスト]")

    return "\n".join(prompt_parts)

def parse_ai_response(text):
    details = {}
    sections = {
        "---性格---": "personality",
        "---通院理由---": "reason",
        "---行動パターン---": "behavior",
        "---口コミ重視点---": "reviews",
        "---医療価値観---": "values",
        "---医療求めるもの---": "demands"
    }
    current_key = None
    current_text = []

    # AIの応答がNoneの場合のフォールバック
    if text is None:
        print("DEBUG parse_ai_response: Input text is None.") # ログ追加
        print("Warning: AI response text is None.")
        for marker, key in sections.items():
            details[key] = "(AIからの応答がありませんでした)"
        return details

    for line in text.splitlines():
        found_key = None
        for marker, key in sections.items():
            if line.strip() == marker:
                found_key = key
                break
        if found_key:
            if current_key and current_text:
                details[sections[current_key]] = "\n".join(current_text).strip()
            current_key = marker # ここを sections[found_key] ではなく marker にする
            current_text = []
        elif current_key: # current_key が None でないことを確認
            current_text.append(line)
    
    if current_key and current_text: # ループ終了後、最後のセクションを処理
        details[sections[current_key]] = "\n".join(current_text).strip()
    
    # 解析できなかったセクションにデフォルト値を設定
    for marker, key in sections.items():
        if key not in details:
            details[key] = "(AIからの応答が解析できませんでした)"
    return details

# --- Static File Serving --- 
@app.route('/')
def serve_index():
    return send_from_directory('.', 'index.html')

@app.route('/admin')
def serve_admin():
     return send_from_directory('.', 'admin.html')

# --- API Endpoints --- 
@app.route('/api/generate', methods=['POST'])
def generate_persona():
    data = request.get_json()
    if not data:
        return jsonify({"error": "リクエストが空です。"}), 400

    try:
        # --- Get API keys and selected models from environment variables ---
        openai_api_key = os.environ.get("OPENAI_API_KEY")
        print(f"DEBUG: Attempting to use OpenAI API Key from env: '{openai_api_key}'") # <--- ログ出力追加

        anthropic_api_key = os.environ.get("ANTHROPIC_API_KEY")
        google_api_key = os.environ.get("GOOGLE_API_KEY")

        selected_text_model = os.environ.get("SELECTED_TEXT_MODEL", "gpt-4.1") # Default if not set
        print(f"DEBUG: Selected text model: {selected_text_model}") # <--- モデルもログ出力
        selected_image_model = os.environ.get("SELECTED_IMAGE_MODEL", "dall-e-3") # Default

        # Determine which API key to use based on the selected text model
        api_key_to_use = None
        if selected_text_model.startswith("gpt"):
            api_key_to_use = openai_api_key
        elif selected_text_model.startswith("claude"):
            api_key_to_use = anthropic_api_key
        elif selected_text_model.startswith("gemini"):
            api_key_to_use = google_api_key
        else: # Fallback or if model doesn't clearly indicate provider
            api_key_to_use = openai_api_key # Default to OpenAI key if unsure

        # APIキーがNoneの場合のエラーチェックを追加
        if api_key_to_use is None and selected_text_model.startswith("gpt"):
             print("ERROR: OpenAI API key is None after selection logic!")
             raise ValueError("OpenAI APIキーが環境変数から取得できませんでした。適切なAPIキーが設定されているか確認してください。")

        # --- Initialize AI Client for Text Generation ---
        text_generation_client = get_ai_client(selected_text_model, api_key_to_use)
        
        prompt = build_prompt(data)
        # print(f"DEBUG: Generated Prompt:\n{prompt}") # 必要に応じてプロンプトもログ出力

        generated_text_str = None # 初期化
        image_url = None # 初期化

        # --- Text Generation ---
        if selected_text_model.startswith("gpt"):
            try:
                print(f"DEBUG: Calling OpenAI API ({selected_text_model}) for text generation...")
                completion = text_generation_client.chat.completions.create(
                    model=selected_text_model,
                    messages=[{"role": "user", "content": prompt}]
                )
                generated_text_str = completion.choices[0].message.content
                print(f"DEBUG OpenAI Response content: >>>\n{generated_text_str}\n<<<END OpenAI Response")
                print("DEBUG: OpenAI API call for text generation successful.")
            except Exception as e:
                print(f"ERROR calling OpenAI API for text generation ({selected_text_model}): {str(e)}")
                traceback.print_exc()
                generated_text_str = None # Ensure it's None on error
        elif selected_text_model.startswith("claude"):
            try:
                print(f"DEBUG: Calling Anthropic API ({selected_text_model}) for text generation...")
                if not isinstance(text_generation_client, Anthropic):
                    raise TypeError("Expected Anthropic client for Claude model")
                response = text_generation_client.messages.create(
                    model=selected_text_model,
                    max_tokens=2000,
                    messages=[{"role": "user", "content": prompt}]
                 )
                generated_text_str = response.content[0].text if response.content else None
                print("DEBUG: Anthropic API call for text generation successful.")
            except Exception as e:
                print(f"ERROR calling Anthropic API for text generation ({selected_text_model}): {str(e)}")
                traceback.print_exc()
                generated_text_str = None # Ensure it's None on error
        elif selected_text_model.startswith("gemini"):
            try:
                print(f"DEBUG: Calling Google Gemini API ({selected_text_model}) for text generation...")
                response = text_generation_client.generate_content(prompt)
                generated_text_str = response.text
                print("DEBUG: Google Gemini API call for text generation successful.")
            except Exception as e:
                print(f"ERROR calling Google Gemini API for text generation ({selected_text_model}): {str(e)}")
                traceback.print_exc()
                generated_text_str = None # Ensure it's None on error
        else:
            print(f"WARNING: Text generation skipped, unsupported model type: {selected_text_model}")
            generated_text_str = None # Ensure it's None if skipped

        generated_details = parse_ai_response(generated_text_str)
        print(f"DEBUG parse_ai_response output: {generated_details}")

        # --- Image Generation (常に試行) ---
        # selected_image_model = os.environ.get("SELECTED_IMAGE_MODEL", "dall-e-3")
        # image_api_key_for_dalle = os.environ.get("OPENAI_API_KEY")
        # 
        # if image_api_key_for_dalle and selected_image_model and selected_image_model != "none":
        #     try:
        #         image_client = OpenAI(api_key=image_api_key_for_dalle)
        #         img_name = profile_data.get('name', 'person')
        #         img_age = profile_data.get('age', 'age unknown')
        #         img_gender = profile_data.get('gender', 'gender unknown')
        #         img_occupation = profile_data.get('occupation')
        #
        #         img_prompt_parts = [
        #             f"Create a profile picture for a persona named {img_name}",
        #             f"who is {img_age}, {img_gender}.",
        #             "Style: realistic photo."
        #         ]
        #         if img_occupation and img_occupation.strip():
        #             img_prompt_parts.append(f"Occupation: {img_occupation}.")
        #         
        #         img_prompt = " ".join(img_prompt_parts)
        #         print(f"Attempting image generation with prompt: {img_prompt}")
        #
        #         image_response = image_client.images.generate(
        #             model=selected_image_model,
        #             prompt=img_prompt,
        #             size="1024x1024", 
        #             quality="standard", 
        #             n=1,
        #         )
        #         image_url = image_response.data[0].url
        #         print(f"Image generated successfully: {image_url}")
        #     except Exception as img_e:
        #         print(f"Image generation failed: {img_e}")
        #         traceback.print_exc()
        # else:
        #     print(f"Image generation skipped: OpenAI API key for DALL-E not found or image model is '{selected_image_model}'.")

        # Set a dummy image URL
        image_url = "https://placehold.jp/150x150.png" # 以前のプレースホルダーURLを変更

        response_data = {
            "profile": data,
            "details": generated_details,
            "image_url": image_url
        }
        return jsonify(response_data)

    except ValueError as ve:
        print(f"Value Error: {ve}")
        return jsonify({"error": str(ve)}), 400
    except Exception as e:
        print(f"Error during AI generation: {e}")
        traceback.print_exc()
        return jsonify({"error": "ペルソナ生成中にエラーが発生しました。"}), 500

# --- Admin API Endpoints (Simplified for Environment Variables) --- 

@app.route('/api/admin/settings', methods=['GET'])
def get_all_settings():
    """Gets settings from environment variables."""
    settings = {
        # API Keys (Return placeholders or indicate if set, DO NOT RETURN ACTUAL KEYS)
        "openai_api_key": "Set" if os.environ.get("OPENAI_API_KEY") else "Not Set",
        "google_api_key": "Set" if os.environ.get("GOOGLE_API_KEY") else "Not Set",
        "anthropic_api_key": "Set" if os.environ.get("ANTHROPIC_API_KEY") else "Not Set",
        # Selected Models
        "selected_text_model": os.environ.get("SELECTED_TEXT_MODEL", "gpt-4.1"),
        "selected_image_model": os.environ.get("SELECTED_IMAGE_MODEL", "dall-e-3"),
        # Limits
        "limit_personality": os.environ.get("LIMIT_PERSONALITY", "100"),
        "limit_reason": os.environ.get("LIMIT_REASON", "100"),
        "limit_behavior": os.environ.get("LIMIT_BEHAVIOR", "100"),
        "limit_reviews": os.environ.get("LIMIT_REVIEWS", "100"),
        "limit_values": os.environ.get("LIMIT_VALUES", "100"),
        "limit_demands": os.environ.get("LIMIT_DEMANDS", "100"),
        # Output Formats Enabled (Read from Env Vars, default to 'true')
        "output_pdf_enabled": os.environ.get("OUTPUT_PDF_ENABLED", "true"),
        "output_ppt_enabled": os.environ.get("OUTPUT_PPT_ENABLED", "true"),
        "output_gslide_enabled": os.environ.get("OUTPUT_GSLIDE_ENABLED", "false") # Default false
    }
    return jsonify(settings)

# --- Settings Saving Endpoints (Removed/Placeholder) --- 
# Saving settings directly via API is not suitable for environment variables.
# These need to be set in the Render dashboard.
@app.route('/api/admin/settings/api', methods=['POST'])
def save_api_settings():
    # This endpoint is now effectively read-only for env vars
    return jsonify({"message": "APIキーはRenderの環境変数で設定してください。"}), 405 # Method Not Allowed

@app.route('/api/admin/settings/limits', methods=['POST'])
def save_limit_settings():
    # This endpoint is now effectively read-only for env vars
    return jsonify({"message": "文字数制限はRenderの環境変数で設定してください。"}), 405 # Method Not Allowed

@app.route('/api/admin/settings/output', methods=['POST'])
def save_output_settings():
     # This endpoint is now effectively read-only for env vars
     return jsonify({"message": "出力設定はRenderの環境変数で設定してください。"}), 405 # Method Not Allowed


# --- User-facing Output Settings Endpoint --- 
@app.route('/api/settings/output', methods=['GET'])
def get_output_settings_for_user():
    """Gets output format enabled status from environment variables."""
    settings = {
        "output_pdf_enabled": os.environ.get("OUTPUT_PDF_ENABLED", "true").lower() == "true",
        "output_ppt_enabled": os.environ.get("OUTPUT_PPT_ENABLED", "true").lower() == "true",
        "output_gslide_enabled": os.environ.get("OUTPUT_GSLIDE_ENABLED", "false").lower() == "true"
    }
    return jsonify(settings)


# --- PDF Generation --- 
def generate_pdf(data):
    # A4サイズに設定、レイアウト最適化
    pdf = FPDF(orientation='P', unit='mm', format='A4')
    
    # 余白を小さく設定 (左右上下の余白を10mmに)
    pdf.set_margins(10, 10, 10)
    pdf.set_auto_page_break(True, margin=10)
    
    pdf.add_page()
    
    # フォント設定
    try:
        # Register Regular font
        pdf.add_font("ipa", "", "ipaexg.ttf", uni=True)
        # Register Bold style using the same file
        pdf.add_font("ipa", "B", "ipaexg.ttf", uni=True)
        pdf.set_font("ipa", size=10) # 全体的に小さいフォントサイズをデフォルトに
    except RuntimeError as e:
        print(f"WARNING: Could not load/register font 'ipaexg.ttf'. Error: {e}. Using default font.")
        pdf.set_font("Arial", size=10) # フォールバックフォントも小さく

    profile = data.get('profile', {})
    details = data.get('details', {})

    # --- ヘッダー部分（タイトルと診療科・目的） ---
    # 名前とタイトル
    pdf.set_font_size(14)
    pdf.cell(0, 8, profile.get('name', 'Unknown Persona'), ln=True, align='C')
    
    # 診療科と目的を横並びで表示
    pdf.set_font_size(9)
    pdf.cell(30, 5, "診療科:", align='R')
    pdf.cell(70, 5, profile.get('department', '-'), ln=0)
    pdf.cell(30, 5, "作成目的:", align='R')
    pdf.cell(50, 5, profile.get('purpose', '-'), ln=1)
    pdf.ln(2)
    
    # ページを2段組にレイアウト - 左側に基本情報とアイコン、右側に詳細情報
    
    # --- 左カラム (80mm幅) ---
    left_column_width = 80
    left_margin = pdf.l_margin
    current_y = pdf.get_y()
    
    # プロフィール画像用のスペース (仮のテキスト)
    pdf.set_xy(left_margin, current_y)
    pdf.set_font_size(8)
    pdf.cell(left_column_width, 5, "ペルソナ画像", border=0, ln=1, align='C')
    pdf.ln(15) # 画像用スペース
    
    # 基本情報セクション
    current_y = pdf.get_y()
    pdf.set_xy(left_margin, current_y)
    pdf.set_font("ipa", 'B', 11)
    pdf.cell(left_column_width, 5, "基本情報", ln=1, fill=False)
    pdf.ln(1)
    
    # 基本情報の項目を2列に表示するための設定
    info_col_width = left_column_width / 2
    
    # 左列の情報項目
    pdf.set_font("ipa", 'B', 8)
    info_items_left = [
        ("性別", profile.get('gender', '-')),
        ("年齢", profile.get('age', '-')),
        ("都道府県", profile.get('prefecture', '-')),
        ("市区町村", profile.get('municipality', '-')),
        ("職業", profile.get('occupation', '-'))
    ]
    
    # 右列の情報項目
    info_items_right = [
        ("年収", profile.get('income', '-')),
        ("家族構成", profile.get('family', '-')),
        ("趣味", profile.get('hobby', '-')),
        ("ライフイベント", profile.get('life_events', '-')),
        ("患者タイプ", profile.get('patient_type', '-'))
    ]
    
    # 左列の描画
    current_y = pdf.get_y()
    for key, value in info_items_left:
        pdf.set_xy(left_margin, current_y)
        pdf.set_font("ipa", 'B', 8)
        pdf.cell(25, 4, f"{key}:", 0)
        pdf.set_font("ipa", '', 8)
        pdf.cell(info_col_width - 25, 4, str(value) if value else '-', 0)
        current_y += 4
    
    # 右列の描画
    current_y = pdf.get_y() - 20 # 上に戻る
    for key, value in info_items_right:
        pdf.set_xy(left_margin + info_col_width, current_y)
        pdf.set_font("ipa", 'B', 8)
        pdf.cell(25, 4, f"{key}:", 0)
        pdf.set_font("ipa", '', 8)
        pdf.cell(info_col_width - 25, 4, str(value) if value else '-', 0)
        current_y += 4
    
    # 追加情報セクション（一列表示）
    pdf.ln(5)
    current_y = pdf.get_y()
    pdf.set_xy(left_margin, current_y)
    pdf.set_font("ipa", 'B', 11)
    pdf.cell(left_column_width, 5, "その他の特徴", ln=1, fill=False)
    pdf.ln(1)
    
    # 追加情報の項目
    additional_items = [
        ("座右の銘", profile.get('motto', '-')),
        ("最近の悩み/関心", profile.get('concerns', '-')),
        ("好きな有名人", profile.get('favorite_person', '-')),
        ("よく見るメディア", profile.get('media_sns', '-')),
        ("性格キーワード", profile.get('personality_keywords', '-')),
        ("健康に関する行動", profile.get('health_actions', '-')),
        ("休日の過ごし方", profile.get('holiday_activities', '-')),
        ("キャッチコピー", profile.get('catchphrase', '-'))
    ]
    
    # 追加項目の描画（幅を短くして説明文を追加）
    for key, value in additional_items:
        pdf.set_font("ipa", 'B', 7)
        pdf.cell(30, 3.5, f"{key}:", 0)
        pdf.set_font("ipa", '', 7)
        # 長いテキストは折り返して表示するが、高さは固定
        pdf.cell(left_column_width - 30, 3.5, str(value) if value else '-', 0, ln=1)
    
    # 動的に追加された項目があれば表示
    if profile.get('additional_field_name') and profile.get('additional_field_value'):
        additional_fields = zip(profile.get('additional_field_name'), profile.get('additional_field_value'))
        for field_name, field_value in additional_fields:
            if field_name or field_value:
                pdf.set_font("ipa", 'B', 7)
                pdf.cell(30, 3.5, f"{field_name}:", 0)
                pdf.set_font("ipa", '', 7)
                pdf.cell(left_column_width - 30, 3.5, str(field_value) if field_value else '-', 0, ln=1)
    
    # --- 右カラム (詳細情報) ---
    right_column_width = pdf.w - pdf.r_margin - left_margin - left_column_width
    right_margin = left_margin + left_column_width
    
    # Map internal keys to Japanese headers
    header_map = {
        "personality": "性格（価値観・人生観）",
        "reason": "通院理由",
        "behavior": "症状通院頻度・行動パターン",
        "reviews": "口コミの重視ポイント",
        "values": "医療機関への価値観・行動傾向",
        "demands": "医療機関に求めるもの"
    }

    # 右カラムの開始位置を設定
    current_y = current_y - (len(additional_items) * 3.5) # 左カラムの追加情報の高さ分戻る
    if current_y < 30: # あまりに上すぎる場合は下げる
        current_y = 30
        
    # 詳細情報の描画
    pdf.set_xy(right_margin, current_y)
    for key, japanese_header in header_map.items():
        value = details.get(key) # Get value using the internal key
        if value: # Only add if value exists
            # ヘッダー
            pdf.set_xy(right_margin, current_y)
            pdf.set_font("ipa", 'B', 10)
            pdf.set_fill_color(240, 240, 240) # より薄いグレー
            pdf.cell(right_column_width, 5, japanese_header, ln=1, fill=True)
            current_y += 5
            
            # 内容
            pdf.set_xy(right_margin, current_y)
            pdf.set_font("ipa", '', 9)
            
            # テキストを複数行に分割してレンダリング（MultiCellの高さを固定）
            text_height = pdf.font_size * 1.5
            lines = str(value).split('\n')
            for line in lines:
                while len(line) > 0:
                    # 約60文字で改行
                    if len(line) > 60:
                        chunk = line[:60]
                        line = line[60:]
                    else:
                        chunk = line
                        line = ""
                    
                    pdf.set_xy(right_margin, current_y)
                    pdf.cell(right_column_width, text_height, chunk, ln=1)
                    current_y += text_height
            
            current_y += 2 # 各セクション間の余白

    # Generate PDF in memory
    pdf_output = pdf.output() # Get output as bytes directly
    buffer = io.BytesIO(pdf_output)
    buffer.seek(0)
    return buffer

@app.route('/api/download/pdf', methods=['POST'])
def download_pdf():
    data = request.json
    if not data:
        return jsonify({"error": "No data provided"}), 400

    try:
        pdf_buffer = generate_pdf(data) # generate_pdf returns a BytesIO buffer
        # Construct filename, handle potential missing name
        profile_name = data.get('profile', {}).get('name', 'persona')
        safe_profile_name = "".join(c if c.isalnum() or c in [' ', '(', ')'] else '_' for c in profile_name)
        filename_utf8 = f"{safe_profile_name}_persona.pdf"
        filename_encoded = quote(filename_utf8) # Use urllib.parse.quote

        # Create Response object manually
        response = Response(pdf_buffer.getvalue(), mimetype='application/pdf')
        # Set headers on the response object
        response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{filename_encoded}"

        return response # Return the Response object

        # --- Previous code using send_file (commented out) ---
        # headers = {
        #     'Content-Disposition': f"attachment; filename*=UTF-8''{filename_encoded}"
        # }
        # return send_file(
        #     pdf_buffer,
        #     mimetype='application/pdf',
        #     as_attachment=True,
        #     headers=headers # Pass the custom header
        # )
    except Exception as e:
        print(f"Error generating PDF: {e}")
        import traceback
        traceback.print_exc() # Print full traceback for debugging
        return jsonify({"error": "Failed to generate PDF"}), 500

# Helper function to generate PPT (basic implementation)
def generate_ppt(data):
    prs = Presentation()
    profile = data.get('profile', {})
    details = data.get('details', {})

    # --- 1枚のスライドに全ての情報を収める ---
    # ワイドスクリーン(16:9)のきれいなレイアウトを使用
    slide_layout = prs.slide_layouts[5]  # 白紙のレイアウト
    slide = prs.slides.add_slide(slide_layout)
    
    # スライドの寸法を取得（インチ単位）
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    
    # 余白の設定
    margin = Inches(0.3)

    # --- ヘッダー部分（タイトルと診療科・目的） ---
    # 名前とタイトル
    title_shape = slide.shapes.add_textbox(margin, margin, Inches(slide_width - 0.6), Inches(0.6))
    title_frame = title_shape.text_frame
    title_frame.text = profile.get('name', 'Unknown Persona')
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True

    # 診療科と目的を横並びで表示
    header_shape = slide.shapes.add_textbox(margin, Inches(0.9), Inches(slide_width - 0.6), Inches(0.4))
    header_frame = header_shape.text_frame
    header_para = header_frame.add_paragraph()
    header_para.text = f"診療科: {profile.get('department', '-')}    作成目的: {profile.get('purpose', '-')}"
    header_para.font.size = Pt(12)
    
    # 2段組レイアウトの寸法計算
    column_height = slide_height - 1.7  # ヘッダー部分を引く
    left_column_width = slide_width * 0.4  # 左側40%
    right_column_width = slide_width * 0.55  # 右側55%
    
    # --- 左カラム (基本情報) ---
    left_column_y = Inches(1.3)
    
    # 「基本情報」ヘッダー
    basic_header = slide.shapes.add_textbox(margin, left_column_y, Inches(left_column_width), Inches(0.3))
    basic_header_frame = basic_header.text_frame
    basic_para = basic_header_frame.add_paragraph()
    basic_para.text = "基本情報"
    basic_para.font.size = Pt(14)
    basic_para.font.bold = True
    
    # 基本情報の項目を2列に表示するための設定
    left_column_y += Inches(0.4)
    info_col_width = left_column_width / 2
    
    # 左列の情報項目
    info_items_left = [
        ("性別", profile.get('gender', '-')),
        ("年齢", profile.get('age', '-')),
        ("都道府県", profile.get('prefecture', '-')),
        ("市区町村", profile.get('municipality', '-')),
        ("職業", profile.get('occupation', '-'))
    ]
    
    # 右列の情報項目
    info_items_right = [
        ("年収", profile.get('income', '-')),
        ("家族構成", profile.get('family', '-')),
        ("趣味", profile.get('hobby', '-')),
        ("ライフイベント", profile.get('life_events', '-')),
        ("患者タイプ", profile.get('patient_type', '-'))
    ]
    
    # 左列の描画
    current_y = left_column_y
    item_height = Inches(0.25)
    for key, value in info_items_left:
        # キー
        key_box = slide.shapes.add_textbox(margin, current_y, Inches(info_col_width * 0.4), item_height)
        key_frame = key_box.text_frame
        key_para = key_frame.add_paragraph()
        key_para.text = f"{key}:"
        key_para.font.bold = True
        key_para.font.size = Pt(10)
        
        # 値
        val_box = slide.shapes.add_textbox(margin + Inches(info_col_width * 0.4), current_y, Inches(info_col_width * 0.6), item_height)
        val_frame = val_box.text_frame
        val_para = val_frame.add_paragraph()
        val_para.text = str(value) if value else '-'
        val_para.font.size = Pt(10)
        
        current_y += item_height
    
    # 右列の描画
    current_y = left_column_y
    for key, value in info_items_right:
        # キー
        key_box = slide.shapes.add_textbox(margin + Inches(info_col_width), current_y, Inches(info_col_width * 0.4), item_height)
        key_frame = key_box.text_frame
        key_para = key_frame.add_paragraph()
        key_para.text = f"{key}:"
        key_para.font.bold = True
        key_para.font.size = Pt(10)

        # 値
        val_box = slide.shapes.add_textbox(margin + Inches(info_col_width * 1.4), current_y, Inches(info_col_width * 0.6), item_height)
        val_frame = val_box.text_frame
        val_para = val_frame.add_paragraph()
        val_para.text = str(value) if value else '-'
        val_para.font.size = Pt(10)
        
        current_y += item_height
    
    # その他の特徴
    additional_header_y = current_y + Inches(0.2)
    additional_header = slide.shapes.add_textbox(margin, additional_header_y, Inches(left_column_width), Inches(0.3))
    additional_header_frame = additional_header.text_frame
    additional_para = additional_header_frame.add_paragraph()
    additional_para.text = "その他の特徴"
    additional_para.font.size = Pt(14)
    additional_para.font.bold = True
    
    # 追加情報の項目
    additional_items = [
        ("座右の銘", profile.get('motto', '-')),
        ("最近の悩み/関心", profile.get('concerns', '-')),
        ("好きな有名人", profile.get('favorite_person', '-')),
        ("よく見るメディア", profile.get('media_sns', '-')),
        ("性格キーワード", profile.get('personality_keywords', '-')),
        ("健康に関する行動", profile.get('health_actions', '-')),
        ("休日の過ごし方", profile.get('holiday_activities', '-')),
        ("キャッチコピー", profile.get('catchphrase', '-'))
    ]
    
    # 追加項目の描画
    current_y = additional_header_y + Inches(0.3)
    additional_item_height = Inches(0.20)
    for key, value in additional_items:
        # 1行に収める形式で描画
        key_box = slide.shapes.add_textbox(margin, current_y, Inches(left_column_width * 0.35), additional_item_height)
        key_frame = key_box.text_frame
        key_para = key_frame.add_paragraph()
        key_para.text = f"{key}:"
        key_para.font.bold = True
        key_para.font.size = Pt(9)
        
        val_box = slide.shapes.add_textbox(margin + Inches(left_column_width * 0.35), current_y, Inches(left_column_width * 0.65), additional_item_height)
        val_frame = val_box.text_frame
        val_para = val_frame.add_paragraph()
        val_para.text = str(value) if value else '-'
        val_para.font.size = Pt(9)

        current_y += additional_item_height
    
    # 動的に追加された項目があれば表示
    if profile.get('additional_field_name') and profile.get('additional_field_value'):
        additional_fields = zip(profile.get('additional_field_name'), profile.get('additional_field_value'))
        for field_name, field_value in additional_fields:
            if field_name or field_value:
                key_box = slide.shapes.add_textbox(margin, current_y, Inches(left_column_width * 0.35), additional_item_height)
                key_frame = key_box.text_frame
                key_para = key_frame.add_paragraph()
                key_para.text = f"{field_name}:"
                key_para.font.bold = True
                key_para.font.size = Pt(9)
                
                val_box = slide.shapes.add_textbox(margin + Inches(left_column_width * 0.35), current_y, Inches(left_column_width * 0.65), additional_item_height)
                val_frame = val_box.text_frame
                val_para = val_frame.add_paragraph()
                val_para.text = str(field_value) if field_value else '-'
                val_para.font.size = Pt(9)
                
                current_y += additional_item_height
    
    # --- 右カラム (詳細情報) ---
    right_column_x = margin + Inches(left_column_width + 0.2)
    right_column_y = Inches(1.3)
    
    # ヘッダーマップ
    header_map = {
        "personality": "性格（価値観・人生観）",
        "reason": "通院理由",
        "behavior": "症状通院頻度・行動パターン",
        "reviews": "口コミの重視ポイント",
        "values": "医療機関への価値観・行動傾向",
        "demands": "医療機関に求めるもの"
    }

    # 詳細情報の描画
    for key, japanese_header in header_map.items():
        value = details.get(key)
        if value:
            # ヘッダー
            header_box = slide.shapes.add_textbox(right_column_x, right_column_y, Inches(right_column_width), Inches(0.3))
            header_frame = header_box.text_frame
            header_para = header_frame.add_paragraph()
            header_para.text = japanese_header
            header_para.font.bold = True
            header_para.font.size = Pt(12)
            
            # 背景色付けのための長方形
            header_rect = slide.shapes.add_shape(
                1, # 長方形
                right_column_x, 
                right_column_y, 
                Inches(right_column_width), 
                Inches(0.3)
            )
            header_rect.fill.solid()
            header_rect.fill.fore_color.rgb = RGBColor(240, 240, 240)  # 薄いグレー
            header_rect.line.fill.background()  # 枠線を消す
            header_rect.zorder = 1  # 下に配置
            header_box.zorder = 2  # テキストを上に

            # 内容
            content_box = slide.shapes.add_textbox(
                right_column_x, 
                right_column_y + Inches(0.3), 
                Inches(right_column_width), 
                Inches(0.5)  # 高さは仮設定、後で調整
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_para = content_frame.add_paragraph()
            content_para.text = str(value)
            content_para.font.size = Pt(10)
            
            right_column_y += Inches(0.9)  # 次のセクションへの間隔

    # Save presentation to a BytesIO buffer
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

@app.route('/api/download/ppt', methods=['POST'])
def download_ppt():
    data = request.json
    if not data:
        return jsonify({"error": "No data provided"}), 400

    try:
        ppt_buffer = generate_ppt(data)
        profile_name = data.get('profile', {}).get('name', 'persona')
        safe_profile_name = "".join(c if c.isalnum() or c in [' ', '(', ')'] else '_' for c in profile_name)
        filename_utf8 = f"{safe_profile_name}_persona.pptx" # Use .pptx extension
        filename_encoded = quote(filename_utf8)

        response = Response(
            ppt_buffer.getvalue(),
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation' # PPTX mimetype
        )
        response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{filename_encoded}"

        return response

    except Exception as e:
        print(f"Error generating PPT: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": "Failed to generate PPT"}), 500

# --- Run App --- (Removed)
# if __name__ == '__main__':
#     # initialize_database() # Removed DB initialization call
#     app.run(debug=True, port=5000) 
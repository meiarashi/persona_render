from fastapi import FastAPI, HTTPException, Request
from fastapi.staticfiles import StaticFiles
from fastapi.responses import JSONResponse, Response, FileResponse
from pathlib import Path
import json
import os
import random
import traceback
import io
import re
from urllib.parse import quote
from urllib.request import urlopen
import base64

# For PDF/PPT generation
from PIL import Image
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import MSO_VERTICAL_ANCHOR

# For AI clients
try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

# New Gemini SDK (preferred)
try:
    from google import genai as google_genai_sdk # Alias to avoid conflict
    from google.genai import types as google_genai_types
except ImportError:
    google_genai_sdk = None
    google_genai_types = None

# Old Gemini SDK (fallback)
try:
    import google.generativeai as old_gemini_sdk
except ImportError:
    old_gemini_sdk = None
    
try:
    from anthropic import Anthropic
except ImportError:
    Anthropic = None

# Import the router from the routers module
from .routers import admin_settings
from . import crud
from . import models
from . import rag_processor

# Create a FastAPI app instance
app = FastAPI(
    title="Persona Render Admin API",
    description="API for managing Persona Render application settings.",
    version="0.1.0"
)

# Mount the admin_settings router
app.include_router(admin_settings.router)

# --- Static files hosting ---
current_file_dir = Path(__file__).resolve().parent
backend_dir = current_file_dir.parent
project_root_dir = backend_dir.parent
frontend_dir = project_root_dir / "frontend"

if not frontend_dir.exists() or not frontend_dir.is_dir():
    alt_frontend_dir = project_root_dir / "frontend"
    if (project_root_dir / "backend").exists() and (project_root_dir / "frontend").exists():
        frontend_dir = project_root_dir / "frontend"
    elif (backend_dir.parent / "frontend").exists():
         frontend_dir = backend_dir.parent / "frontend"
    else:
        frontend_dir = backend_dir / "frontend"

if frontend_dir.exists() and frontend_dir.is_dir():
    print(f"Serving static files from: {frontend_dir}")
    app.mount("/static/user", StaticFiles(directory=frontend_dir / "user"), name="user_static_assets")
    app.mount("/static/admin", StaticFiles(directory=frontend_dir / "admin"), name="admin_static_assets")
    images_dir = project_root_dir / "images"
    if images_dir.exists() and images_dir.is_dir():
        app.mount("/images", StaticFiles(directory=images_dir), name="image_assets")
        print(f"Serving image files from: {images_dir}")

    @app.get("/admin", include_in_schema=False)
    async def serve_admin_html():
        admin_html_path = frontend_dir / "admin/admin.html"
        if admin_html_path.exists(): return FileResponse(admin_html_path)
        fallback_html_path = project_root_dir / "frontend/admin/admin.html"
        if fallback_html_path.exists(): return FileResponse(fallback_html_path)
        raise HTTPException(status_code=404, detail="admin.html not found")

    @app.get("/", include_in_schema=False)
    async def serve_user_html():
        user_html_path = frontend_dir / "user/index.html"
        if user_html_path.exists(): return FileResponse(user_html_path)
        fallback_html_path = project_root_dir / "frontend/user/index.html"
        if fallback_html_path.exists(): return FileResponse(fallback_html_path)
        raise HTTPException(status_code=404, detail="index.html not found")
    print("User UI and Admin UI routes are set up.")
else:
    print(f"Frontend directory not found at {frontend_dir}. Static file serving skipped.")

# --- Settings Migration ---
def migrate_image_model_settings():
    """Migrates old Gemini image model name in settings to the new one."""
    try:
        settings = crud.read_settings()
        if settings and settings.models and \
           settings.models.image_api_model == "gemini-2.0-flash-preview-image-generation":  # 古い名前をチェック
            print("Migrating image_api_model name in settings...")
            settings.models.image_api_model = "gemini-2.0-flash-exp-image-generation"  # 新しい名前に設定
            crud.write_settings(settings)
            print("Settings migrated: Updated Gemini image model name to gemini-2.0-flash-exp-image-generation")
    except Exception as e:
        print(f"Settings migration for image model failed: {e}")

@app.on_event("startup")
async def startup_event():
    print("Running startup migration...")
    migrate_image_model_settings()

# --- AI Client Initialization Helper --- 
def get_ai_client(model_name, api_key):
    """Initializes and returns the correct AI client based on model name."""
    if model_name.startswith("gpt"):
        if not OpenAI: raise ValueError("OpenAI library not loaded.")
        if not api_key: raise ValueError("OpenAI APIキーが設定されていません。")
        try: return OpenAI(api_key=api_key)
        except Exception as e: raise ValueError(f"OpenAI Client Error: {e}")
    elif model_name.startswith("claude"):
        if not Anthropic: raise ValueError("Anthropic library not loaded.")
        if not api_key: raise ValueError("Anthropic APIキーが設定されていません。")
        try: return Anthropic(api_key=api_key)
        except Exception as e: raise ValueError(f"Anthropic Client Error: {e}")
    elif model_name.startswith("gemini"):
        if not api_key: raise ValueError("Google APIキーが設定されていません。")
        # Prefer new SDK
        if google_genai_sdk:
            try:
                print("[DEBUG] Using new google-genai SDK for client initialization.")
                return google_genai_sdk.Client(api_key=api_key)
            except Exception as e:
                print(f"[WARNING] New google-genai SDK client init failed: {e}. Trying old SDK.")
        # Fallback to old SDK
        if old_gemini_sdk:
            try:
                print("[DEBUG] Using old google.generativeai SDK for client initialization (fallback).")
                old_gemini_sdk.configure(api_key=api_key)
                # For the old SDK, GenerativeModel is returned, not a 'Client' instance
                return old_gemini_sdk.GenerativeModel(model_name) 
            except Exception as e_old:
                raise ValueError(f"Google Client Error (old SDK): {e_old}")
        raise ValueError("Google AI SDK (new or old) not available or failed to initialize.")
    else:
        raise ValueError(f"未対応のモデルです: {model_name}")

# --- Function to build prompts ---
def build_prompt(data, limit_personality="100", limit_reason="100", limit_behavior="100", 
                 limit_reviews="100", limit_values="100", limit_demands="100"):

    # --- Patient Type Descriptions ---
    patient_type_details_map = {
        '利便性重視型': { "description": 'アクセスの良さ、待ち時間の短さ、診療時間の柔軟性など、便利さを最優先', "example": '忙しいビジネスパーソン、オンライン診療を好む患者' },
        '専門医療追求型': { "description": '専門医や高度専門医療機関での治療を希望し、医師の経歴や実績を重視', "example": '難病患者、複雑な症状を持つ患者' },
        '予防健康管理型': { "description": '病気になる前の予防や早期発見、健康維持に関心が高い', "example": '定期健診を欠かさない人、予防接種に積極的な人' },
        '代替医療志向型': { "description": '漢方、鍼灸、ホメオパシーなど、西洋医学以外の選択肢を積極的に取り入れる', "example": '自然療法愛好者、慢性疾患の患者' },
        '経済合理型': { "description": '自己負担額、保険適用の有無、費用対効果を重視', "example": '経済的制約のある患者、医療費控除を意識する人' },
        '情報探求型': { "description": '徹底的な情報収集、セカンドオピニオン取得、比較検討を行う', "example": '高学歴層、慎重な意思決定を好む患者' },
        '革新技術指向型': { "description": '最先端の医療技術、新薬、臨床試験などに積極的に関心を持つ', "example": '既存治療で効果が出なかった患者、医療イノベーションに関心がある人' },
        '対話重視型': { "description": '医師からの丁寧な説明や対話を求め、質問が多い', "example": '不安を感じやすい患者、医療従事者' },
        '信頼基盤型': { "description": 'かかりつけ医との長期的な関係や医療機関の評判を重視', "example": '地域密着型の患者、同じ医師に長期通院する患者' },
        '緊急解決型': { "description": '症状の即時改善を求め、緊急性を重視', "example": '急性疾患患者、痛みに耐性が低い患者' },
        '受動依存型': { "description": '医師の判断に全面的に依存し、自らの決定より医師の指示を優先', "example": '高齢者、医療知識が少ない患者' },
        '自律決定型': { "description": '自分の治療に主体的に関わり、最終決定権を持ちたいと考える', "example": '医療リテラシーが高い患者、自己管理を好む慢性疾患患者' }
    }

    prompt_parts = [
        "以下の情報に基づいて、医療系のペルソナを作成してください。",
        "各項目は自然な文章で記述し、**日本語で、指定されたおおよその文字数制限に従ってください**。",
        "",
        "# 利用者からの入力情報"
    ]

    # 基本情報
    basic_info = {
        "診療科": data.get('department'),
        "ペルソナ作成目的": data.get('purpose'),
        "名前": data.get('name'),
        "性別": data.get('gender'),
        "年齢": data.get('age'),
        "都道府県": data.get('prefecture'),
        "市区町村": data.get('municipality'),
        "家族構成": data.get('family'),
        "職業": data.get('occupation'),
        "年収": data.get('income'),
        "趣味": data.get('hobby'),
        "ライフイベント": data.get('life_events'),
        "患者タイプ": data.get('patient_type')
    }
    # setting_type と patient_type も追加
    if data.get('setting_type'):
        basic_info["基本情報設定タイプ"] = data.get('setting_type')

    for key, value in basic_info.items():
        if value: # 値が存在する場合のみプロンプトに追加
            prompt_parts.append(f"- {key}: {value}")
            if key == "患者タイプ" and value in patient_type_details_map:
                details = patient_type_details_map[value]
                prompt_parts.append(f"  - 患者タイプの特徴: {details['description']}")
                prompt_parts.append(f"  - 患者タイプの例: {details['example']}")
        # 「選択された患者タイプ」の(自動生成)条件から patient_type を除外 (「患者タイプ」で処理されるため)
        elif key in ["名前", "性別", "年齢", "都道府県", "市区町村", "家族構成", "職業", "年収", "趣味", "ライフイベント"] and key != "患者タイプ": 
             prompt_parts.append(f"- {key}: (自動生成)")
        elif key == "患者タイプ" and not value and data.get('setting_type') == 'patient_type': # setting_typeがpatient_typeで患者タイプが未選択の場合
             prompt_parts.append(f"- {key}: (指定なし/自動生成)")

    # Step 4の固定追加項目
    fixed_additional_info = {
        "座右の銘": data.get('motto'),
        "最近の悩み/関心": data.get('concerns'),
        "好きな有名人/尊敬する人物": data.get('favorite_person'),
        "よく見るメディア/SNS": data.get('media_sns'),
        "性格キーワード": data.get('personality_keywords'),
        "最近した健康に関する行動": data.get('health_actions'),
        "休日の過ごし方": data.get('holiday_activities'),
        "キャッチコピー": data.get('catchphrase'),
    }
    has_fixed_additional_info = any(fixed_additional_info.values())
    if has_fixed_additional_info:
        prompt_parts.append("\n## 追加情報（固定項目）")
        for key, value in fixed_additional_info.items():
            if value:
                prompt_parts.append(f"- {key}: {value}")

    # Step 4の動的追加項目
    additional_field_names = data.get('additional_field_name', [])
    additional_field_values = data.get('additional_field_value', [])
    
    dynamic_additional_info = []
    if isinstance(additional_field_names, list) and isinstance(additional_field_values, list):
        for i in range(len(additional_field_names)):
            field_name = additional_field_names[i]
            field_value = additional_field_values[i] if i < len(additional_field_values) else None
            if field_name and field_value: # 項目名と内容の両方がある場合のみ
                dynamic_additional_info.append(f"- {field_name}: {field_value}")
    
    if dynamic_additional_info:
        prompt_parts.append("\n## 追加情報（自由入力項目）")
        prompt_parts.extend(dynamic_additional_info)

    prompt_parts.append("\n# 生成項目")
    prompt_parts.append("以下の項目について、上記情報に基づいた自然な文章を生成してください。各項目は指定された文字数の目安で記述してください。")
    prompt_parts.append("\n重要: 出力には項目名と内容のみを含め、文字数の指定（例：「(100文字程度)」）は出力に含めないでください。")
    prompt_parts.append(f"\n1. **性格（価値観・人生観）**: {limit_personality}文字程度で記述")
    prompt_parts.append(f"2. **通院理由**: {limit_reason}文字程度で記述")
    prompt_parts.append(f"3. **症状通院頻度・行動パターン**: {limit_behavior}文字程度で記述")
    prompt_parts.append(f"4. **口コミの重視ポイント**: {limit_reviews}文字程度で記述")
    prompt_parts.append(f"5. **医療機関への価値観・行動傾向**: {limit_values}文字程度で記述")
    prompt_parts.append(f"6. **医療機関に求めるもの**: {limit_demands}文字程度で記述")
    
    return "\n".join(prompt_parts)

# --- Function to parse AI response ---
def parse_ai_response(text):
    """Extract sections from AI-generated text response."""
    sections = {
        "personality": "",
        "reason": "",
        "behavior": "",
        "reviews": "",
        "values": "",
        "demands": ""
    }
    
    # Try to parse structured output
    try:
        lines = text.strip().split('\n')
        current_section = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Check if this is a section header for personality section
            if line.startswith("1.") and ("性格" in line or "価値観" in line or "人生観" in line):
                current_section = "personality"
                # Remove the header and keep only content
                content = line.split(':', 1)[1].strip() if ':' in line else ""
                if content:
                    # 「(100文字程度)」のような文字数指定を削除
                    content = re.sub(r'[（(]\d+文字程度[）)]\s*', '', content)
                    content = re.sub(r'\d+文字程度\s*', '', content)
                    sections[current_section] = content
                continue
                
            # Check for reason section
            elif line.startswith("2.") and "通院理由" in line:
                current_section = "reason"
                content = line.split(':', 1)[1].strip() if ':' in line else ""
                if content:
                    content = re.sub(r'\(\d+文字程度\)\s*', '', content)
                    sections[current_section] = content
                continue
                
            # Check for behavior section
            elif line.startswith("3.") and ("症状" in line or "通院頻度" in line or "行動パターン" in line):
                current_section = "behavior"
                content = line.split(':', 1)[1].strip() if ':' in line else ""
                if content:
                    content = re.sub(r'\(\d+文字程度\)\s*', '', content)
                    sections[current_section] = content
                continue
                
            # Check for reviews section
            elif line.startswith("4.") and ("口コミ" in line or "重視ポイント" in line):
                current_section = "reviews"
                content = line.split(':', 1)[1].strip() if ':' in line else ""
                if content:
                    content = re.sub(r'\(\d+文字程度\)\s*', '', content)
                    sections[current_section] = content
                continue
                
            # Check for values section
            elif line.startswith("5.") and ("医療機関" in line or "価値観" in line or "行動傾向" in line):
                current_section = "values"
                content = line.split(':', 1)[1].strip() if ':' in line else ""
                if content:
                    content = re.sub(r'\(\d+文字程度\)\s*', '', content)
                    sections[current_section] = content
                continue
                
            # Check for demands section
            elif line.startswith("6.") and ("医療機関" in line or "求めるもの" in line):
                current_section = "demands"
                content = line.split(':', 1)[1].strip() if ':' in line else ""
                if content:
                    content = re.sub(r'\(\d+文字程度\)\s*', '', content)
                    sections[current_section] = content
                continue
                
            # If we're in a section, append text
            if current_section and not line.startswith(("#", "##")):
                # 「(100文字程度)」のような文字数指定を削除
                cleaned_line = re.sub(r'[（(]\d+文字程度[）)]\s*', '', line)
                cleaned_line = re.sub(r'\d+文字程度\s*', '', cleaned_line)
                if sections[current_section]:
                    sections[current_section] += " " + cleaned_line
                else:
                    sections[current_section] = cleaned_line
    
    except Exception as e:
        print(f"Error parsing AI response: {e}")
        traceback.print_exc()
    
    return sections

@app.post("/api/generate")
async def generate_persona(request: Request):
    try:
        data = await request.json()
        
        # --- 設定をcrudから読み込む ---
        app_settings = crud.read_settings() # AdminSettings インスタンスが返る
        
        selected_text_model = app_settings.models.text_api_model if app_settings.models else "gpt-3.5-turbo" # デフォルト値
        selected_image_model = app_settings.models.image_api_model if app_settings.models else "none" # デフォルト値

        # --- APIキーの取得 ---
        openai_api_key = os.environ.get("OPENAI_API_KEY")
        anthropic_api_key = os.environ.get("ANTHROPIC_API_KEY")
        google_api_key = os.environ.get("GOOGLE_API_KEY")

        # テキスト生成用APIキーの選択
        text_api_key_to_use = None
        if selected_text_model.startswith("gpt"):
            text_api_key_to_use = openai_api_key
        elif selected_text_model.startswith("claude"):
            text_api_key_to_use = anthropic_api_key
        elif selected_text_model.startswith("gemini"):
            text_api_key_to_use = google_api_key
        
        # テスト用ダミーデータ (テキスト生成) - APIキーがない場合
        if not text_api_key_to_use and not (selected_text_model.startswith("gemini") and google_api_key): # Geminiはキーがなくてもgenai.configureされるため別途判定
             if selected_text_model != "dummy": # "dummy" モデル選択時以外でキーがない場合
                print(f"Warning: API key for {selected_text_model} not found. Text generation might fail or use defaults.")
                # ここでエラーを返すか、ダミーデータを返すかは設計次第
                # 今回はget_ai_clientやAPI呼び出しでエラーになることを許容


        # 文字数制限を取得（管理画面の設定を優先）
        if app_settings.limits:
            limit_personality = app_settings.limits.get("personality", "100")
            limit_reason = app_settings.limits.get("reason", "100")
            limit_behavior = app_settings.limits.get("behavior", "100")
            limit_reviews = app_settings.limits.get("reviews", "100")
            limit_values = app_settings.limits.get("values", "100")
            limit_demands = app_settings.limits.get("demands", "100")
        else:
            # フォールバック：環境変数から取得
            limit_personality = os.environ.get("LIMIT_PERSONALITY", "100")
            limit_reason = os.environ.get("LIMIT_REASON", "100")
            limit_behavior = os.environ.get("LIMIT_BEHAVIOR", "100")
            limit_reviews = os.environ.get("LIMIT_REVIEWS", "100")
            limit_values = os.environ.get("LIMIT_VALUES", "100")
            limit_demands = os.environ.get("LIMIT_DEMANDS", "100")
        
        # RAGデータベースの初期化
        rag_processor.init_rag_database()
        
        # RAGデータの検索
        rag_context = ""
        rag_results = []
        department = data.get('department')
        age = data.get('age')
        gender = data.get('gender')
        
        if department:
            # 年齢グループの判定
            age_group = None
            if age:
                try:
                    age_str = str(age)
                    # 年齢を数値に変換
                    if 'y' in age_str:
                        # "45y3m" -> 45
                        age_num = int(age_str.split('y')[0])
                    elif age_str.isdigit():
                        # "45" -> 45
                        age_num = int(age_str)
                    else:
                        age_num = None
                    
                    if age_num is not None:
                        if age_num < 20:
                            age_group = "10s"
                        elif age_num < 30:
                            age_group = "20s"
                        elif age_num < 40:
                            age_group = "30s"
                        elif age_num < 50:
                            age_group = "40s"
                        elif age_num < 60:
                            age_group = "50s"
                        elif age_num < 70:
                            age_group = "60s"
                        else:
                            age_group = "70s"
                except (ValueError, AttributeError):
                    print(f"Warning: Could not parse age '{age}' for RAG search")
            
            # RAGデータの検索
            print(f"[DEBUG] Searching RAG data for department={department}, age_group={age_group}, gender={gender}")
            rag_results = rag_processor.search_rag_data(
                specialty=department,
                age_group=age_group,
                gender=gender,
                limit=5
            )
            print(f"[DEBUG] RAG search returned {len(rag_results)} results")
            
            if rag_results:
                rag_context = "\n\n# 参考情報（この診療科の患者が検索するキーワード）\n"
                rag_context += "以下は、同じ診療科・年代・性別の患者がよく検索するキーワードです。ペルソナ作成の参考にしてください：\n"
                for i, result in enumerate(rag_results, 1):
                    rag_context += f"{i}. {result['keyword']} (検索数: {result['search_volume']}人)\n"
                    # カテゴリーフィールドは削除（CSVに存在しないため）
        
        # プロンプト構築（RAGコンテキストを含む）
        prompt_text = build_prompt(data, limit_personality, limit_reason, limit_behavior, 
                                  limit_reviews, limit_values, limit_demands) + rag_context
        
        # 生成情報をログ出力
        print(f"[INFO] Generating persona with model: {selected_text_model}")
        print(f"[INFO] Department: {department}, Age: {age}, Gender: {gender}")
        print(f"[INFO] Character limits - Personality: {limit_personality}文字, "
              f"Reason: {limit_reason}文字, "
              f"Behavior: {limit_behavior}文字, "
              f"Reviews: {limit_reviews}文字, "
              f"Values: {limit_values}文字, "
              f"Demands: {limit_demands}文字")
        if rag_results:
            print(f"[INFO] RAG data included: {len(rag_results)} keywords from {department}")
            for i, result in enumerate(rag_results[:3], 1):
                print(f"  - Keyword {i}: {result['keyword']}")
        else:
            print(f"[INFO] No RAG data available for {department}")
        
        # AIクライアント初期化 (テキスト生成用)
        text_generation_client = None
        client_init_error = None
        if text_api_key_to_use or (selected_text_model.startswith("gemini") and google_api_key):
            try:
                text_generation_client = get_ai_client(selected_text_model, text_api_key_to_use)
            except Exception as e:
                client_init_error = str(e)
                print(f"AI Client initialization error for text model {selected_text_model}: {client_init_error}")
        
        generated_text_str = None
        # テキスト生成実行
        if text_generation_client:
            try:
                if selected_text_model.startswith("gpt"):
                    completion = text_generation_client.chat.completions.create(
                        model=selected_text_model,
                        messages=[{"role": "user", "content": prompt_text}],
                        temperature=0.7,
                        max_tokens=2000
                    )
                    generated_text_str = completion.choices[0].message.content
                elif selected_text_model.startswith("claude"):
                    response = text_generation_client.messages.create(
                        model=selected_text_model,
                        max_tokens=2000,
                        messages=[{"role": "user", "content": prompt_text}]
                    )
                    generated_text_str = response.content[0].text if response.content else None
                elif selected_text_model.startswith("gemini"):
                    # 新しいSDKの場合
                    if hasattr(text_generation_client, 'models'):
                        print(f"[DEBUG] Using new SDK for text generation with model: {selected_text_model}")
                        response = text_generation_client.models.generate_content(
                            model=selected_text_model,
                            contents=prompt_text
                        )
                        if response.candidates and response.candidates[0].content.parts:
                            generated_text_str = response.candidates[0].content.parts[0].text
                    # 古いSDKの場合
                    else:
                        print(f"[DEBUG] Using old SDK for text generation with model: {selected_text_model}")
                        response = text_generation_client.generate_content(prompt_text)
                        generated_text_str = response.text
                        # Gemini特有の文字数指定を追加で削除
                        generated_text_str = re.sub(r'[（(]\d+文字程度[）)]\s*', '', generated_text_str)
                        generated_text_str = re.sub(r'\d+文字程度[\s、。]', '', generated_text_str)
            except Exception as e:
                print(f"Error during text generation with {selected_text_model}: {e}")
                traceback.print_exc()
        
        if generated_text_str is None: # AI生成失敗またはスキップ時のフォールバック
            error_msg = "Text generation failed or was skipped."
            if client_init_error:
                error_msg += f" Client initialization error: {client_init_error}"
            print(error_msg)
            generated_details = {
                "personality": "真面目で責任感が強く、家族を大切にする。健康意識が高く、予防医療に関心がある。",
                "reason": "定期的な健康診断と、軽度の高血圧の管理のため。",
                "behavior": "3ヶ月に一度の定期検診に欠かさず通院。処方された降圧剤を規則正しく服用している。",
                "reviews": "医師の説明がわかりやすいこと、待ち時間が短いこと、スタッフの対応が丁寧であることを重視する。",
                "values": "信頼できる医師との長期的な関係を望む。予防医療に前向きで、医師のアドバイスを真摯に受け止める。",
                "demands": "わかりやすい説明と、必要に応じて専門医への適切な紹介。予防医療のアドバイスも欲しい。"
            }
        else:
            print(f"[DEBUG] Raw AI response preview (first 200 chars): {generated_text_str[:200] if generated_text_str else 'None'}")
            generated_details = parse_ai_response(generated_text_str)
            # 生成後のチェック
            for key, value in generated_details.items():
                if "文字程度" in str(value):
                    print(f"[WARNING] Character count text found in {key}: {value[:50]}...")
            

        # --- 画像生成 ---
        image_url = "https://placehold.jp/150x150.png" # デフォルトプレースホルダー

        img_prompt_parts = [
            f"Create a profile picture for a persona named {data.get('name', 'person')}",
            f"who is {data.get('age', 'age unknown')}, {data.get('gender', 'gender unknown')}.",
            "Style: realistic photo."
        ]
        if data.get('occupation'):
            img_prompt_parts.append(f"Occupation: {data.get('occupation')}.")
        img_prompt = " ".join(img_prompt_parts)

        if selected_image_model == "dall-e-3":
            if openai_api_key:
                try:
                    image_client = OpenAI(api_key=openai_api_key)
                    print(f"Attempting DALL-E 3 image generation with prompt: {img_prompt}")
                    image_response = image_client.images.generate(
                        model="dall-e-3",
                        prompt=img_prompt,
                        size="1024x1024", 
                        quality="standard", 
                        n=1,
                    )
                    image_url = image_response.data[0].url
                    print(f"DALL-E 3 Image generated successfully: {image_url}")
                except Exception as img_e:
                    print(f"DALL-E 3 Image generation failed: {img_e}")
                    traceback.print_exc()
                    image_url = "https://placehold.jp/300x200/EEE/777?text=DALL-E+Error"
            else:
                print("OpenAI API key not found for DALL-E 3. Skipping image generation.")
                image_url = "https://placehold.jp/300x200/CCC/555?text=No+OpenAI+Key"

        elif selected_image_model == "gemini-2.0-flash-exp-image-generation":
            if google_api_key:
                try:
                    print(f"[INFO] Attempting Gemini image generation with prompt: {img_prompt}")
                    
                    # 新しいSDKが利用可能かチェック
                    if google_genai_sdk and google_genai_types:
                        # 新しいSDKでクライアントを作成
                        client = google_genai_sdk.Client(api_key=google_api_key)
                        
                        print(f"[DEBUG] Using new SDK for image generation with model: {selected_image_model}")
                        
                        # 画像生成API呼び出し
                        response = client.models.generate_content(
                            model=selected_image_model,
                            contents=img_prompt,
                            config=google_genai_types.GenerateContentConfig(
                                response_modalities=['TEXT', 'IMAGE']
                            )
                        )
                        
                        print(f"[DEBUG] Gemini Raw Response Object: {type(response)}")
                        
                        # レスポンスの検査 (より詳細に)
                        if hasattr(response, 'candidates') and response.candidates:
                            candidate = response.candidates[0]
                            if candidate.content and candidate.content.parts:
                                # 画像データを探す
                                image_found = False
                                for part in candidate.content.parts:
                                    if hasattr(part, 'inline_data') and part.inline_data:
                                        if hasattr(part.inline_data, 'data'):
                                            print(f"[DEBUG] Gemini image part found. Mime type: {getattr(part.inline_data, 'mime_type', 'image/png')}")
                                            image_data = part.inline_data.data # bytes
                                            mime_type = getattr(part.inline_data, 'mime_type', 'image/png')
                                            base64_image = base64.b64encode(image_data).decode('utf-8')
                                            image_url = f"data:{mime_type};base64,{base64_image}"
                                            print("[INFO] Gemini Image generated successfully as Base64 Data URI.")
                                            image_found = True
                                            break
                                
                                if not image_found:
                                    print("[ERROR] No image data found in response parts")
                                    image_url = "https://placehold.jp/300x200/EEE/777?text=No+Image+Data"
                            else:
                                print("[ERROR] Gemini response candidate found, but no content or parts.")
                                if hasattr(candidate, 'finish_reason'):
                                    print(f"[ERROR] Gemini Finish Reason: {candidate.finish_reason}")
                                if hasattr(candidate, 'safety_ratings'):
                                    print(f"[ERROR] Gemini Safety Ratings: {candidate.safety_ratings}")
                                image_url = "https://placehold.jp/300x200/EEE/777?text=No+Content"
                        elif hasattr(response, 'text') and response.text:
                            print(f"[ERROR] Gemini image generation returned text instead of image: {response.text}")
                            image_url = "https://placehold.jp/300x200/EEE/777?text=Text+Response"
                        else:
                            print("[ERROR] No candidates in Gemini response")
                            image_url = "https://placehold.jp/300x200/EEE/777?text=No+Candidates"
                            
                    else:
                        print("[ERROR] New Gemini SDK not available for image generation")
                        image_url = "https://placehold.jp/300x200/EEE/777?text=No+SDK"

                except Exception as img_e:
                    print(f"[ERROR] Gemini Image generation failed: {img_e}")
                    traceback.print_exc()
                    # エラーメッセージに詳細を含める
                    error_message = str(img_e).replace("\\n", " ") # 改行をスペースに
                    image_url = f"https://placehold.jp/300x200/EEE/777?text=Gemini+Error"
            else:
                print("Google API key not found for Gemini. Skipping image generation.")
                image_url = "https://placehold.jp/300x200/CCC/555?text=No+Google+Key"
        
        elif selected_image_model == "none":
            print("Image generation set to 'none'.")
            image_url = "https://placehold.jp/150x150.png?text=No+Image" # No Image
        
        # レスポンスデータ作成
        response_data = {
            "profile": data, # フロントから送られてきた入力データをそのまま返す
            "details": generated_details,
            "image_url": image_url # DALL-EならURL、GeminiならBase64 Data URI
        }
        return response_data

    except Exception as e:
        print(f"Error in generate_persona: {e}")
        traceback.print_exc()
        return JSONResponse(
            status_code=500,
            content={"error": f"サーバーエラー: {str(e)}"}
        )

@app.post("/api/download/pdf")
async def download_pdf(request: Request):
    """ペルソナデータをPDFとしてダウンロードするエンドポイント"""
    try:
        data = await request.json()
        if not data:
            return JSONResponse(
                status_code=400,
                content={"error": "No data provided"}
            )

        # PDF生成
        pdf_buffer = generate_pdf(data)
        
        # ファイル名の設定
        profile_name = data.get('profile', {}).get('name', 'persona')
        safe_profile_name = "".join(c if c.isalnum() or c in [' ', '(', ')'] else '_' for c in profile_name)
        filename_utf8 = f"{safe_profile_name}_persona.pdf"
        filename_encoded = quote(filename_utf8)

        # レスポンス作成
        response = Response(pdf_buffer.getvalue(), media_type='application/pdf')
        response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{filename_encoded}"
        
        return response
        
    except Exception as e:
        print(f"Error generating PDF: {e}")
        traceback.print_exc()
        return JSONResponse(
            status_code=500,
            content={"error": f"Failed to generate PDF: {str(e)}"}
        )

@app.post("/api/download/ppt")
async def download_ppt(request: Request):
    """ペルソナデータをPPTXとしてダウンロードするエンドポイント"""
    try:
        data = await request.json()
        if not data:
            return JSONResponse(
                status_code=400,
                content={"error": "No data provided"}
            )
        
        # PPTXデータ作成のためのデータ変換
        persona_data = {}
        for key, value in data.get('profile', {}).items():
            persona_data[key] = value
        
        # details からの転送
        for key, value in data.get('details', {}).items():
            persona_data[key] = value
        
        # 画像URL
        image_url = data.get('image_url')
        image_path = None
        
        # 診療科と目的の取得
        department_val = persona_data.get('department', '-')
        department_text = DEPARTMENT_MAP.get(department_val.lower(), department_val) if department_val else '-'
        
        purpose_val = persona_data.get('purpose', '-')
        purpose_text = PURPOSE_MAP.get(purpose_val.lower(), purpose_val) if purpose_val else '-'
        
        # 画像があれば一時ファイルに保存
        if image_url:
            try:
                # 一時画像ファイルの処理ここでは省略
                pass
            except Exception as e:
                print(f"Error downloading image: {e}")
        
        # PPTX生成
        pptx_buffer = generate_ppt(persona_data, image_path, department_text, purpose_text)
        
        # ファイル名の設定
        profile_name = data.get('profile', {}).get('name', 'persona')
        safe_profile_name = "".join(c if c.isalnum() or c in [' ', '(', ')'] else '_' for c in profile_name)
        filename_utf8 = f"{safe_profile_name}_persona.pptx"
        filename_encoded = quote(filename_utf8)
        
        # レスポンス作成
        response = Response(
            pptx_buffer.getvalue(),
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{filename_encoded}"
        
        return response
        
    except Exception as e:
        print(f"Error generating PPTX: {e}")
        traceback.print_exc()
        return JSONResponse(
            status_code=500,
            content={"error": f"Failed to generate PPTX: {str(e)}"}
        )

@app.get("/health", summary="Health check endpoint", tags=["Health"])
async def health_check():
    """
    Simple health check endpoint to confirm the API is running.
    """
    return {"status": "ok"}

# To run this app locally (from the 'backend' directory, assuming main.py is in 'app'):
# Ensure 'app' is a package (has __init__.py if needed, though often not for FastAPI structure like this)
# cd backend
# python -m uvicorn app.main:app --reload --port 8000

# If your structure is flatter, e.g. main.py in 'backend' root:
# cd backend
# uvicorn main:app --reload --port 8000 

# --- Global Maps and Helper Functions for PDF/PPT Generation ---
GENDER_MAP = {
    "male": "男性", "female": "女性", "other": "その他",
}
INCOME_MAP = {
    "<100": "100万円未満", "100-200": "100-200万円", "200-300": "200-300万円",
    "300-400": "300-400万円", "400-500": "400-500万円", "500-600": "500-600万円",
    "600-700": "600-700万円", "700-800": "700-800万円", "800-900": "800-900万円",
    "900-1000": "900-1000万円", "1000-1100": "1000-1100万円", "1100-1200": "1100-1200万円",
    "1200-1300": "1200-1300万円", "1300-1400": "1300-1400万円", "1400-1500": "1400-1500万円",
    "1500-1600": "1500-1600万円", "1600-1700": "1600-1700万円", "1700-1800": "1700-1800万円",
    "1800-1900": "1800-1900万円", "1900-2000": "1900-2000万円", "2000-2100": "2000-2100万円",
    "2100-2200": "2100-2200万円", "2200-2300": "2200-2300万円", "2300-2400": "2300-2400万円",
    "2400-2500": "2400-2500万円", "2500-2600": "2500-2600万円", "2600-2700": "2600-2700万円",
    "2700-2800": "2700-2800万円", "2800-2900": "2800-2900万円", "2900-3000": "2900-3000万円",
    "3000-3100": "3000-3100万円", "3100-3200": "3100-3200万円", "3200-3300": "3200-3300万円",
    "3300-3400": "3300-3400万円", "3400-3500": "3400-3500万円", "3500-3600": "3500-3600万円",
    "3600-3700": "3600-3700万円", "3700-3800": "3700-3800万円", "3800-3900": "3800-3900万円",
    "3900-4000": "3900-4000万円", "4000-4100": "4000-4100万円", "4100-4200": "4100-4200万円",
    "4200-4300": "4200-4300万円", "4300-4400": "4300-4400万円", "4400-4500": "4400-4500万円",
    "4500-4600": "4500-4600万円", "4600-4700": "4600-4700万円", "4700-4800": "4700-4800万円",
}
DEPARTMENT_MAP = {
    "internal_medicine": "内科", "surgery": "外科", "pediatrics": "小児科",
    "orthopedics": "整形外科", "dermatology": "皮膚科", "ophthalmology": "眼科",
    "cardiology": "循環器内科", "psychiatry": "精神科", "dentistry": "歯科",
    "pediatric_dentistry": "小児歯科", "otorhinolaryngology": "耳鼻咽喉科",
    "ent": "耳鼻咽喉科",
    "gynecology": "婦人科",
    "urology": "泌尿器科",
    "neurosurgery": "脳神経外科",
    "general_dentistry": "一般歯科",
    "orthodontics": "矯正歯科",
    "cosmetic_dentistry": "審美歯科",
    "oral_surgery": "口腔外科",
    "anesthesiology": "麻酔科",
    "radiology": "放射線科",
    "rehabilitation": "リハビリテーション科",
    "allergy": "アレルギー科",
    "gastroenterology": "消化器内科",
    "respiratory_medicine": "呼吸器内科",
    "diabetes_medicine": "糖尿病内科",
    "nephrology": "腎臓内科",
    "neurology": "神経内科",
    "hematology": "血液内科",
    "plastic_surgery": "形成外科",
    "beauty_surgery": "美容外科",
}
PURPOSE_MAP = {
    "increase_patients": "患者数を増やす",
    "increase_frequency": "来院頻度を増やす",
    "increase_spend": "客単価を増やす",
}
HEADER_MAP = {
    "personality": "性格（価値観・人生観）",
    "reason": "通院理由",
    "behavior": "症状通院頻度・行動パターン",
    "reviews": "口コミの重視ポイント",
    "values": "医療機関への価値観・行動傾向",
    "demands": "医療機関に求めるもの"
}

def format_age_for_pdf_ppt(age_value):
    if not age_value: return '-'
    try:
        age_value_str = str(age_value) # Ensure it's a string
        if 'm' in age_value_str and 'y' in age_value_str:
            parts = age_value_str.split('y')
            years = parts[0]
            months = parts[1].replace('m', '')
            if years == '0' and months == '0':
                return "0歳"
            return f"{years}歳{months}ヶ月"
        elif 'y' in age_value_str:
            return f"{age_value_str.replace('y', '')}歳"
        elif 'm' in age_value_str:
            months = age_value_str.replace('m', '')
            if months == '0':
                return "0歳"
            return f"0歳{months}ヶ月"
        elif age_value_str.isdigit():
            return f"{age_value_str}歳"
        return age_value_str
    except Exception:
        return str(age_value) 

# Helper function to generate PPT
def add_text_to_shape(shape, text, font_size=Pt(9), is_bold=False, alignment=PP_ALIGN.LEFT, font_name='Meiryo UI'):
    text_frame = shape.text_frame
    text_frame.clear() # Clear existing text and formatting
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    p.clear() # Clear existing runs in the paragraph

    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = font_size
    font.bold = is_bold
    font.color.rgb = RGBColor(0, 0, 0) # Black text
    p.alignment = alignment

    # Estimate height based on text length and font size
    num_lines = len(text.split('\\n'))
    estimated_height = num_lines * font_size.pt * 1.5 # Approximate line height factor
    return Cm(estimated_height / 28.3465 / 2.54) # Convert points to cm (rough estimate)

def generate_pdf(data):
    # A4サイズ横長に設定、レイアウト最適化
    pdf = FPDF(orientation='L', unit='mm', format='A4')
    
    # 余白を小さく設定 (左右上下の余白を8mmに)
    pdf.set_margins(8, 8, 8)
    pdf.set_auto_page_break(True, margin=8)
    
    pdf.add_page()
    
    # フォント設定
    try:
        # Register Regular font
        pdf.add_font("ipa", "", "ipaexg.ttf", uni=True)
        # Register Bold style using the same file
        pdf.add_font("ipa", "B", "ipaexg.ttf", uni=True)
        pdf.set_font("ipa", size=10) # 全体的に小さいフォントサイズをデフォルトに
    except RuntimeError as e:
        print(f"WARNING: Could not load/register font 'ipaexg.ttf'. Error: {e}. Using default font.")
        pdf.set_font("Arial", size=10) # フォールバックフォントも小さく

    profile = data.get('profile', {})
    details = data.get('details', {})
    image_url = data.get('image_url')

    # --- 定数定義 ---
    name_line_height = 6  # ペルソナ名の行の高さ (mm)
    section_title_height = 7 # セクションタイトルのセルの高さ (mm)

    # --- ページとカラムの基本設定 ---
    left_column_content_x = pdf.l_margin
    page_width = pdf.w - pdf.l_margin - pdf.r_margin
    left_column_width = page_width * 0.35
    right_column_width = page_width * 0.65
    column_gap = 5
    right_column_x = left_column_content_x + left_column_width + column_gap
    
    content_start_y = pdf.get_y() # 全体のコンテンツ開始Y座標（一番上の要素に合わせる）
    header_end_y = content_start_y # header_end_y を定義 (アイコンなどの開始Y座標の基準)

    # --- 左カラムの描画開始 ---
    pdf.set_xy(left_column_content_x, header_end_y + 5) # ヘッダーの下に少しスペース

    # --- ペルソナアイコンと名前 ---
    icon_size = 30 # アイコンのサイズ (mm)
    icon_y_position = pdf.get_y()

    if image_url:
        try:
            print(f"Fetching image from URL: {image_url}")
            # Data URLの場合の処理
            if image_url.startswith('data:'):
                # data:image/png;base64,... の形式から画像データを抽出
                try:
                    header, encoded = image_url.split(',', 1)
                    image_data = base64.b64decode(encoded)
                except (ValueError, base64.binascii.Error) as e:
                    print(f"Error decoding data URL: {e}")
                    raise
            else:
                # 通常のURLの場合
                image_data = urlopen(image_url).read()
            
            # PIL Imageで画像を開く
            img_file_obj = io.BytesIO(image_data)
            pil_image = Image.open(img_file_obj)
            
            # アスペクト比を保ちながら正方形にリサイズ（中央クロップ）
            original_width, original_height = pil_image.size
            
            # 正方形にクロップ
            min_dimension = min(original_width, original_height)
            left = (original_width - min_dimension) // 2
            top = (original_height - min_dimension) // 2
            right = left + min_dimension
            bottom = top + min_dimension
            
            # クロップして正方形にする
            cropped_image = pil_image.crop((left, top, right, bottom))
            
            # リサイズして最終サイズに調整
            final_size = int(icon_size * 10)  # mmをpixelに変換（おおよそ）
            resized_image = cropped_image.resize((final_size, final_size), Image.Resampling.LANCZOS)
            
            # RGB形式に変換（PDFで確実に表示されるように）
            if resized_image.mode != 'RGB':
                resized_image = resized_image.convert('RGB')
            
            # バイトストリームに保存
            processed_img_obj = io.BytesIO()
            resized_image.save(processed_img_obj, format='JPEG', quality=85)
            processed_img_obj.seek(0)
            
            # PDFに画像を追加
            pdf.image(processed_img_obj, x=left_column_content_x, y=icon_y_position, w=icon_size, h=icon_size)
            
        except Exception as e:
            print(f"Error loading image: {e}")
            # アイコン失敗時は代替テキストや枠線などを表示
            pdf.rect(left_column_content_x, icon_y_position, icon_size, icon_size, style='D')
            pdf.set_xy(left_column_content_x + 1, icon_y_position + icon_size / 2 - 2)
            pdf.multi_cell(icon_size - 2, 4, "No Img", 0, 'C')

    # 名前の描画開始位置をアイコンの右隣に設定
    name_x_position = left_column_content_x + icon_size + 3 # アイコンの右に少しスペース
    pdf.set_xy(name_x_position, icon_y_position + (icon_size / 2) - (name_line_height / 2) - 2) # 上下中央揃えっぽく調整
    
    pdf.set_font("ipa", 'B', 14) # 名前は少し大きく太字に
    # 名前の最大幅を左カラムの残り幅に制限
    name_max_width = left_column_width - (icon_size + 3) # アイコンとスペース分を引く
    pdf.multi_cell(name_max_width, name_line_height, profile.get('name', '-'), 0, 'L')
    
    # アイコンと名前の下に基本情報を配置するためのY座標を設定
    # アイコンの高さ、または名前の高さのうち、大きい方を基準にする
    current_y_after_icon_name = icon_y_position + icon_size + 5 # アイコンの下端 + 余白

    # --- 基本情報セクション ---
    pdf.set_xy(left_column_content_x, current_y_after_icon_name)
    pdf.set_font("ipa", 'BU', 11) # セクションタイトル (太字・下線)
    pdf.cell(left_column_width, section_title_height, "基本情報", 0, 1, 'L')
    current_y_after_icon_name = pdf.get_y() # 「基本情報」タイトルの後にY座標を更新
    pdf.set_font("ipa", '', 10) # 内容のフォントに戻す

    # 1. 診療科 (左カラム)
    pdf.set_xy(left_column_content_x, current_y_after_icon_name) # 更新されたY座標を使用
    pdf.set_font("ipa", '', 10)
    department_val = profile.get('department', '-')
    department_display = DEPARTMENT_MAP.get(department_val.lower(), department_val) # Lowercase for map key
    pdf.multi_cell(left_column_width, 7, f"診療科: {department_display}", 0, 'L')
    current_y_after_icon_name = pdf.get_y()

    # 2. 作成目的 (左カラム)
    pdf.set_xy(left_column_content_x, current_y_after_icon_name)
    purpose_val = profile.get('purpose', '-')
    purpose_display = PURPOSE_MAP.get(purpose_val.lower(), purpose_val) # Lowercase for map key
    pdf.multi_cell(left_column_width, 7, f"作成目的: {purpose_display}", 0, 'L')
    current_y_after_icon_name = pdf.get_y()
    pdf.ln(3) # 少しスペース
    current_y_after_icon_name = pdf.get_y()

    # 3. 基本情報セクション (左カラム)
    pdf.set_xy(left_column_content_x, current_y_after_icon_name) 
    pdf.set_font("ipa", '', 9)
    info_items = [
        ("性別", GENDER_MAP.get(profile.get('gender', '-'), profile.get('gender', '-'))),
        ("年齢", format_age_for_pdf_ppt(profile.get('age', '-'))),
        ("都道府県", profile.get('prefecture', '-')),
        ("市区町村", profile.get('municipality', '-')),
        ("職業", profile.get('occupation', '-')),
        ("年収", INCOME_MAP.get(profile.get('income', '-'), profile.get('income', '-'))),
        ("家族構成", profile.get('family', '-')),
        ("趣味", profile.get('hobby', '-')),
        ("ライフイベント", profile.get('life_events', '-')),
        ("患者タイプ", profile.get('patient_type', '-'))
    ]
    
    item_height = 4 
    key_width = 25 
    value_width = left_column_width - key_width

    for key, value_display in info_items:
        pdf.set_xy(left_column_content_x, current_y_after_icon_name)
        pdf.set_font("ipa", '', 9)
        pdf.cell(key_width, item_height, f"{key}:", 0, 0, 'L')
        pdf.set_font("ipa", '', 9)
        pdf.set_xy(left_column_content_x + key_width, current_y_after_icon_name)
        pdf.multi_cell(value_width, item_height, str(value_display), 0, 'L')
        current_y_after_icon_name = pdf.get_y() 

    current_y_after_icon_name += 3 
    
    # 4. その他の特徴セクション (左カラム)
    pdf.set_xy(left_column_content_x, current_y_after_icon_name)
    pdf.set_font("ipa", '', 11)
    pdf.set_fill_color(240, 240, 240) 
    pdf.cell(left_column_width, 6, "その他の特徴", 0, 1, 'L', fill=True)
    current_y_after_icon_name = pdf.get_y() 
    pdf.ln(1)
    current_y_after_icon_name = pdf.get_y() 
    
    additional_items = [
        ("座右の銘", profile.get('motto', '-')),
        ("最近の悩み/関心", profile.get('concerns', '-')),
        ("好きな有名人", profile.get('favorite_person', '-')),
        ("よく見るメディア", profile.get('media_sns', '-')),
        ("性格キーワード", profile.get('personality_keywords', '-')),
        ("健康に関する行動", profile.get('health_actions', '-')),
        ("休日の過ごし方", profile.get('holiday_activities', '-')),
        ("キャッチコピー", profile.get('catchphrase', '-'))
    ]
    
    additional_key_width = 30
    additional_value_width = left_column_width - additional_key_width

    for key, value in additional_items:
        pdf.set_xy(left_column_content_x, current_y_after_icon_name)
        pdf.set_font("ipa", '', 9)
        pdf.cell(additional_key_width, item_height, f"{key}:", 0, 0, 'L')
        pdf.set_font("ipa", '', 9)
        pdf.set_xy(left_column_content_x + additional_key_width, current_y_after_icon_name)
        pdf.multi_cell(additional_value_width, item_height, str(value), 0, 'L')
        current_y_after_icon_name = pdf.get_y()

    if profile.get('additional_field_name') and profile.get('additional_field_value'):
        additional_fields = zip(profile.get('additional_field_name'), profile.get('additional_field_value'))
        current_y_after_icon_name +=1 
        for field_name, field_value in additional_fields:
            if field_name or field_value: 
                pdf.set_xy(left_column_content_x, current_y_after_icon_name)
                pdf.set_font("ipa", '', 9)
                pdf.cell(additional_key_width, item_height, f"{field_name if field_name else ''}:", 0, 0, 'L')
                pdf.set_font("ipa", '', 9)
                pdf.set_xy(left_column_content_x + additional_key_width, current_y_after_icon_name)
                pdf.multi_cell(additional_value_width, item_height, str(field_value) if field_value else '-', 0, 'L')
                current_y_after_icon_name = pdf.get_y()
    
    max_left_y = current_y_after_icon_name # 左カラムの最終Y座標

    # --- 右カラムの描画開始 ---
    # 右カラムの開始Y座標を、左カラムのアイコンの開始高さに合わせる
    right_column_current_y = icon_y_position 
    
    for detail_key, japanese_header_text in HEADER_MAP.items():
        value = details.get(detail_key)
        if value and str(value).strip(): # 値が存在し、空でない場合のみ描画
            pdf.set_xy(right_column_x, right_column_current_y)
            
            # セクションヘッダー
            pdf.set_font("ipa", 'BU', 11) # 太字・下線付き、サイズ11
            pdf.set_fill_color(240, 240, 240) # 薄いグレーの背景
            pdf.cell(right_column_width, 6, str(japanese_header_text), 0, 1, 'L', fill=True) # ln=1 で改行
            # pdf.get_y() はこのセルの後に自動更新される

            # コンテンツ
            pdf.set_x(right_column_x) # multi_cell のためにX座標を右カラムの開始位置にリセット
            pdf.set_font("ipa", '', 9) # 通常フォント、サイズ9
            pdf.multi_cell(right_column_width, 5, str(value), 0, 'L') # 行の高さ5mm
            right_column_current_y = pdf.get_y() # multi_cell 後のY座標を更新
            
            pdf.ln(3) # セクション間のスペース
            right_column_current_y = pdf.get_y() # スペース後のY座標を更新

    # Generate PDF in memory
    pdf_output = pdf.output() # Get output as bytes directly
    buffer = io.BytesIO(pdf_output)
    buffer.seek(0)
    return buffer

def generate_ppt(persona_data, image_path=None, department_text=None, purpose_text=None):
    prs = Presentation()
    prs.slide_width = Inches(11.69)  # A4 Landscape width
    prs.slide_height = Inches(8.27) # A4 Landscape height
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Margins (approximated from PDF's 8mm)
    left_margin_ppt = Cm(0.8)
    right_margin_ppt = Cm(0.8)
    top_margin_ppt = Cm(0.8)
    bottom_margin_ppt = Cm(0.8)
    content_width = prs.slide_width - left_margin_ppt - right_margin_ppt
    content_height = prs.slide_height - top_margin_ppt - bottom_margin_ppt
    
    item_spacing_ppt = Cm(0.15) # Spacing between items

    # Title
    title_shape = slide.shapes.add_textbox(left=Cm(0.5), top=Cm(0.2), width=prs.slide_width - Cm(1.0), height=Cm(1.0))
    text_frame = title_shape.text_frame
    p = text_frame.add_paragraph()
    p.text = "生成されたペルソナ"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.name = 'Meiryo UI'

    # Icon
    icon_left = left_margin_ppt
    icon_top = top_margin_ppt + Cm(1.0) # Below title
    icon_size = Cm(3.0)
    
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, icon_left, icon_top, height=icon_size)
        except Exception as e:
            print(f"Error adding image to PPT: {e}")
            # Add a placeholder if image fails
            icon_placeholder = slide.shapes.add_textbox(icon_left, icon_top, icon_size, icon_size)
            add_text_to_shape(icon_placeholder, "画像エラー", font_size=Pt(8))
    else:
        icon_placeholder = slide.shapes.add_textbox(icon_left, icon_top, icon_size, icon_size)
        add_text_to_shape(icon_placeholder, "画像なし", font_size=Pt(8))

    # Persona Name (right of icon)
    name_left = icon_left + icon_size + Cm(0.3)
    name_top = icon_top 
    name_width = content_width * 0.35 - icon_size - Cm(0.3) # Available width in the 35% column part
    name_height = icon_size # Align height with icon
    
    name_text_box = slide.shapes.add_textbox(name_left, name_top, name_width, name_height)
    # Vertical centering
    tf_name = name_text_box.text_frame
    tf_name.word_wrap = True
    tf_name.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE 
    p_name = tf_name.paragraphs[0] if tf_name.paragraphs else tf_name.add_paragraph()
    p_name.clear()
    run_name = p_name.add_run()
    run_name.text = persona_data.get('name', '')
    font_name_style = run_name.font
    font_name_style.name = 'Meiryo UI'
    font_name_style.size = Pt(12)
    font_name_style.bold = True
    p_name.alignment = PP_ALIGN.LEFT

    # Department and Purpose below icon and name
    if not department_text:
        department_text = DEPARTMENT_MAP.get(persona_data.get('department', '').lower(), persona_data.get('department', '-'))
    if not purpose_text:
        purpose_text = PURPOSE_MAP.get(persona_data.get('purpose', '').lower(), persona_data.get('purpose', '-'))

    header_info_top = icon_top + icon_size + Cm(0.2)
    header_info_shape = slide.shapes.add_textbox(icon_left, header_info_top, content_width * 0.33, Cm(1.0))
    text_frame = header_info_shape.text_frame
    text_frame.clear()
    
    header_items = [
        ("診療科", department_text),
        ("作成目的", purpose_text)
    ]
    for i, (label, value) in enumerate(header_items):
        item_p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        if i == 0 and not text_frame.paragraphs: item_p = text_frame.add_paragraph()
        
        run_label = item_p.add_run()
        run_label.text = f"{label}: "
        font_label = run_label.font
        font_label.name = 'Meiryo UI'
        font_label.size = Pt(9)
        font_label.bold = True

        run_value = item_p.add_run()
        run_value.text = value
        font_value = run_value.font
        font_value.name = 'Meiryo UI'
        font_value.size = Pt(9)
        item_p.alignment = PP_ALIGN.LEFT
        if i < len(header_items) -1: # Add line break except for last item
             item_p.add_run().text = '\\n'

    y_position_left = header_info_top + Cm(1.0) + Cm(0.3) # Start Y for left column content
    y_position_right = icon_top # Start Y for right column content (aligned with icon top)

    # Left Column (Basic Information)
    left_column_x = left_margin_ppt
    left_width = content_width * 0.33
    
    current_y_left = y_position_left

    # Section Title: 基本情報
    shape_title_basic = slide.shapes.add_textbox(left_column_x, current_y_left, left_width, Cm(0.6))
    add_text_to_shape(shape_title_basic, "基本情報", font_size=Pt(11), is_bold=True, font_name='Meiryo UI')
    current_y_left += Cm(0.6) + item_spacing_ppt # Update Y position

    basic_info_keys = ["gender", "age", "prefecture", "municipality", "family", "occupation", "income", "hobby", "life_events", "patient_type"]
    basic_info_labels = {
        "gender": "性別", "age": "年齢", "prefecture": "都道府県", "municipality": "市区町村",
        "family": "家族構成", "occupation": "職業", "income": "年収", "hobby": "趣味",
        "life_events": "ライフイベント", "patient_type": "患者タイプ"
    }

    for key in basic_info_keys:
        value = persona_data.get(key, "-")
        if key == "gender": value = GENDER_MAP.get(value, value)
        elif key == "age": value = format_age_for_pdf_ppt(value)
        elif key == "income": value = INCOME_MAP.get(value, value)
        
        item_text = f"{basic_info_labels.get(key, key)}: {value}"
        item_shape = slide.shapes.add_textbox(left_column_x, current_y_left, left_width, Cm(0.5))
        add_text_to_shape(item_shape, item_text, font_size=Pt(9), font_name='Meiryo UI')
        current_y_left += Cm(0.5) + item_spacing_ppt

    # Additional Fixed Fields (Left Column)
    additional_fixed_keys_labels = {
        "motto": "座右の銘", "concerns": "最近の悩み/関心", "favorite_person": "好きな有名人/尊敬する人物",
        "media_sns": "よく見るメディア/SNS", "personality_keywords": "性格キーワード",
        "health_actions": "最近した健康に関する行動", "holiday_activities": "休日の過ごし方",
        "catchphrase": "キャッチコピー"
    }
    # Section Title: 追加情報 (Fixed)
    current_y_left += item_spacing_ppt # Extra space before next section
    shape_title_add_fixed = slide.shapes.add_textbox(left_column_x, current_y_left, left_width, Cm(0.6))
    add_text_to_shape(shape_title_add_fixed, "追加情報", font_size=Pt(11), is_bold=True, font_name='Meiryo UI')
    current_y_left += Cm(0.6) + item_spacing_ppt

    for key, label in additional_fixed_keys_labels.items():
        value = persona_data.get(key, persona_data.get(key.replace('_input', ''), "-"))
        item_text = f"{label}: {value}"
        item_shape = slide.shapes.add_textbox(left_column_x, current_y_left, left_width, Cm(0.5))
        add_text_to_shape(item_shape, item_text, font_size=Pt(9), font_name='Meiryo UI')
        current_y_left += Cm(0.5) + item_spacing_ppt
        
    # Dynamic Additional Fields (Left Column)
    if persona_data.get("additional_field_name") and persona_data.get("additional_field_value"):
        fields = zip(persona_data.get("additional_field_name"), persona_data.get("additional_field_value"))
        has_fields = False
        for field_name, field_value in fields:
            if field_name and field_value:
                if not has_fields:
                    # Section title on first valid field
                    current_y_left += item_spacing_ppt
                    shape_title_add_dyn = slide.shapes.add_textbox(left_column_x, current_y_left, left_width, Cm(0.6))
                    add_text_to_shape(shape_title_add_dyn, "自由記述項目", font_size=Pt(11), is_bold=True, font_name='Meiryo UI')
                    current_y_left += Cm(0.6) + item_spacing_ppt
                    has_fields = True
                
                item_text = f"{field_name}: {field_value}"
                item_shape = slide.shapes.add_textbox(left_column_x, current_y_left, left_width, Cm(0.5))
                add_text_to_shape(item_shape, item_text, font_size=Pt(9), font_name='Meiryo UI')
                current_y_left += Cm(0.5) + item_spacing_ppt

    # Right Column (Detailed Information)
    right_column_x = left_margin_ppt + content_width * 0.35 + Cm(0.3) # Start after left column + gap
    right_width = content_width * 0.65 - Cm(0.3) # Remaining width
    
    current_y_right = y_position_right

    detail_key_map = {
        "personality": "性格（価値観・人生観）",
        "reason": "通院理由",
        "behavior": "症状通院頻度・行動パターン",
        "reviews": "口コミの重視ポイント",
        "values": "医療機関への価値観・行動傾向",
        "demands": "医療機関に求めるもの"
    }
    
    # Draw "性格（価値観・人生観）" first as it's a primary section
    if persona_data.get('personality'):
        section_title = detail_key_map['personality']
        value = persona_data['personality']
        
        title_shape = slide.shapes.add_textbox(right_column_x, current_y_right, right_width, Cm(0.6))
        add_text_to_shape(title_shape, section_title, font_size=Pt(11), is_bold=True, font_name='Meiryo UI')
        current_y_right += Cm(0.6)

        content_shape = slide.shapes.add_textbox(right_column_x, current_y_right, right_width, Cm(2.5))
        add_text_to_shape(content_shape, value, font_size=Pt(9), font_name='Meiryo UI')
        current_y_right += Cm(2.5) + item_spacing_ppt * 2
        
    # Other detailed sections
    for key in ["reason", "behavior", "reviews", "values", "demands"]:
        if key in persona_data and persona_data[key]:
            section_title = detail_key_map.get(key, key)
            value = persona_data[key]
            
            title_shape = slide.shapes.add_textbox(right_column_x, current_y_right, right_width, Cm(0.6))
            add_text_to_shape(title_shape, section_title, font_size=Pt(11), is_bold=True, font_name='Meiryo UI')
            current_y_right += Cm(0.6)

            content_shape = slide.shapes.add_textbox(right_column_x, current_y_right, right_width, Cm(2.0))
            add_text_to_shape(content_shape, value, font_size=Pt(9), font_name='Meiryo UI')
            current_y_right += Cm(2.0) + item_spacing_ppt * 2
    
    # Save to memory stream
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer
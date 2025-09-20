from fastapi import FastAPI, HTTPException, Request, Depends
from fastapi.staticfiles import StaticFiles
from fastapi.responses import JSONResponse, Response, FileResponse, RedirectResponse
from pathlib import Path
import json
import os
import random
import traceback
import io
import re
import tempfile
import requests
from urllib.parse import quote
from urllib.request import urlopen
import base64
import asyncio
import logging

# For PDF/PPT generation
from PIL import Image
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor

# For graph generation
try:
    import matplotlib
    matplotlib.use('Agg')  # Use non-interactive backend
    # Set font configuration before importing pyplot
    matplotlib.rcParams['font.family'] = 'DejaVu Sans'
    matplotlib.rcParams['axes.unicode_minus'] = False
    import matplotlib.pyplot as plt
    try:
        import japanize_matplotlib  # Try to import but don't fail if not available
    except ImportError:
        pass
    GRAPH_ENABLED = True
except ImportError:
    print("Warning: matplotlib not installed. Graph generation disabled.")
    GRAPH_ENABLED = False

# For AI clients
try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

# New Gemini SDK (preferred)
try:
    from google import genai as google_genai_sdk # Alias to avoid conflict
    from google.genai import types as google_genai_types
except ImportError:
    google_genai_sdk = None
    google_genai_types = None

# Old Gemini SDK (fallback)
try:
    import google.generativeai as old_gemini_sdk
except ImportError:
    old_gemini_sdk = None
    
try:
    from anthropic import Anthropic
except ImportError:
    Anthropic = None

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Import from new structure
from api import admin_settings, config
from services import timeline_analyzer
from services import crud, rag_processor
from services.async_image_generator import generate_image_async
from services.cache_manager import get_chief_complaints, preload_cache, load_chief_complaints_data
from .services.competitive_analysis_service import CompetitiveAnalysisService
from .services.google_maps_service import GoogleMapsService
from .middleware.auth import verify_admin_credentials, verify_department_credentials
from .models import schemas as models
from .utils import config_loader, prompt_builder

# Create a FastAPI app instance
# アプリケーション起動時にキャッシュをプリロード
try:
    preload_cache()
except Exception as e:
    print(f"[Warning] Failed to preload cache: {e}")

app = FastAPI(
    title="Persona Render Admin API",
    description="API for managing Persona Render application settings.",
    version="0.1.0"
)

# Log API keys status on startup
@app.on_event("startup")
async def startup_event():
    logger.info("Starting application...")
    
    # 環境変数の包括的チェック
    try:
        from backend.services.env_validator import EnvironmentValidator
        is_valid, report = EnvironmentValidator.check_and_report()
        
        if not is_valid:
            logger.error("必須環境変数が不足しています。設定を確認してください。")
            # 必須変数が不足していても起動は続行（エラーメッセージは表示）
    except Exception as e:
        logger.error(f"環境変数チェック中にエラーが発生: {e}")
        
        # フォールバック: 従来の簡易チェック
        api_keys = {
            "OPENAI_API_KEY": bool(os.getenv("OPENAI_API_KEY")),
            "GOOGLE_MAPS_API_KEY": bool(os.getenv("GOOGLE_MAPS_API_KEY")),
            "SERPAPI_KEY": bool(os.getenv("SERPAPI_KEY")),
            "ESTAT_API_KEY": bool(os.getenv("ESTAT_API_KEY")),
            "ANTHROPIC_API_KEY": bool(os.getenv("ANTHROPIC_API_KEY")),
            "GOOGLE_API_KEY": bool(os.getenv("GOOGLE_API_KEY"))
        }
        
        for key_name, is_set in api_keys.items():
            if is_set:
                logger.info(f"✓ {key_name} is set")
            else:
                logger.warning(f"✗ {key_name} is NOT set")
    
    logger.info("Application startup complete")

# Mount the admin_settings router
app.include_router(admin_settings.router)
# Mount the config router
app.include_router(config.router)

# --- Static files hosting ---
current_file_dir = Path(__file__).resolve().parent  # backend/
project_root_dir = current_file_dir.parent  # persona_render/
frontend_dir = project_root_dir / "frontend"

if not frontend_dir.exists() or not frontend_dir.is_dir():
    alt_frontend_dir = project_root_dir / "frontend"
    if (project_root_dir / "backend").exists() and (project_root_dir / "frontend").exists():
        frontend_dir = project_root_dir / "frontend"
    else:
        frontend_dir = current_file_dir / "frontend"

if frontend_dir.exists() and frontend_dir.is_dir():
    print(f"Serving static files from: {frontend_dir}")
    app.mount("/static/user", StaticFiles(directory=frontend_dir / "user"), name="user_static_assets")
    app.mount("/static/admin", StaticFiles(directory=frontend_dir / "admin"), name="admin_static_assets")
    app.mount("/static/shared", StaticFiles(directory=frontend_dir / "shared"), name="shared_static_assets")
    app.mount("/static/medical", StaticFiles(directory=frontend_dir / "medical"), name="medical_static_assets")
    app.mount("/static/dental", StaticFiles(directory=frontend_dir / "dental"), name="dental_static_assets")
    app.mount("/static/others", StaticFiles(directory=frontend_dir / "others"), name="others_static_assets")
    
    # Mount competitive analysis directories
    app.mount("/user/competitive", StaticFiles(directory=frontend_dir / "user/competitive"), name="user_competitive")
    app.mount("/medical/competitive", StaticFiles(directory=frontend_dir / "medical/competitive"), name="medical_competitive")
    app.mount("/dental/competitive", StaticFiles(directory=frontend_dir / "dental/competitive"), name="dental_competitive")
    app.mount("/others/competitive", StaticFiles(directory=frontend_dir / "others/competitive"), name="others_competitive")
    
    # Mount images from frontend directory
    frontend_images_dir = frontend_dir / "images"
    if frontend_images_dir.exists() and frontend_images_dir.is_dir():
        app.mount("/assets/images", StaticFiles(directory=frontend_images_dir), name="frontend_images")
        print(f"Serving images from: {frontend_images_dir}")
    
    # Mount assets directory (new structure) as fallback
    assets_dir = project_root_dir / "assets"
    if assets_dir.exists() and assets_dir.is_dir():
        app.mount("/assets", StaticFiles(directory=assets_dir), name="assets")
        print(f"Serving assets from: {assets_dir}")
    
    # Legacy images directory support
    images_dir = project_root_dir / "images"
    if images_dir.exists() and images_dir.is_dir():
        app.mount("/images", StaticFiles(directory=images_dir), name="image_assets")
        print(f"Serving image files from: {images_dir}")
    
    # Serve Excel template file from root directory
    @app.get("/クリニック情報.xlsx", include_in_schema=False)
    async def serve_clinic_template():
        excel_path = project_root_dir / "クリニック情報.xlsx"
        if excel_path.exists():
            return FileResponse(
                excel_path,
                media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                filename="クリニック情報.xlsx"
            )
        raise HTTPException(status_code=404, detail="Excel template not found")

    @app.get("/admin", include_in_schema=False)
    async def serve_admin_html(username: str = Depends(verify_admin_credentials)):
        admin_html_path = frontend_dir / "admin/admin.html"
        if admin_html_path.exists(): return FileResponse(admin_html_path)
        fallback_html_path = project_root_dir / "frontend/admin/admin.html"
        if fallback_html_path.exists(): return FileResponse(fallback_html_path)
        raise HTTPException(status_code=404, detail="admin.html not found")

    @app.get("/", include_in_schema=False)
    async def serve_user_dashboard(username: str = Depends(verify_admin_credentials)):
        dashboard_path = frontend_dir / "user/dashboard.html"
        if dashboard_path.exists(): 
            return FileResponse(dashboard_path)
        # フォールバック: 旧index.htmlを使用
        user_html_path = frontend_dir / "user/index.html"
        if user_html_path.exists(): 
            return FileResponse(user_html_path)
        raise HTTPException(status_code=404, detail="user/dashboard.html not found")
    
    @app.get("/user/persona", include_in_schema=False)
    async def serve_user_persona(username: str = Depends(verify_admin_credentials)):
        persona_path = frontend_dir / "user/persona/index.html"
        if persona_path.exists(): 
            return FileResponse(persona_path)
        # フォールバック: 旧index.htmlを使用
        user_html_path = frontend_dir / "user/index.html"
        if user_html_path.exists(): 
            return FileResponse(user_html_path)
        raise HTTPException(status_code=404, detail="user/persona/index.html not found")
    
    @app.get("/user/competitive", include_in_schema=False)
    async def serve_user_competitive(username: str = Depends(verify_admin_credentials)):
        competitive_path = frontend_dir / "user/competitive/index.html"
        if competitive_path.exists(): 
            return FileResponse(competitive_path)
        raise HTTPException(status_code=404, detail="user/competitive/index.html not found")
    
    # 診療科別エンドポイント
    @app.get("/medical", include_in_schema=False)
    async def serve_medical_dashboard(username: str = Depends(verify_department_credentials("medical"))):
        dashboard_path = frontend_dir / "medical/dashboard.html"
        if dashboard_path.exists(): 
            return FileResponse(dashboard_path)
        # フォールバック: 旧index.htmlを使用
        medical_html_path = frontend_dir / "medical/index.html"
        if medical_html_path.exists(): 
            return FileResponse(medical_html_path)
        raise HTTPException(status_code=404, detail="medical/dashboard.html not found")
    
    @app.get("/medical/persona", include_in_schema=False)
    async def serve_medical_persona(username: str = Depends(verify_department_credentials("medical"))):
        persona_path = frontend_dir / "medical/persona/index.html"
        if persona_path.exists(): 
            return FileResponse(persona_path)
        # フォールバック: 旧index.htmlを使用
        medical_html_path = frontend_dir / "medical/index.html"
        if medical_html_path.exists(): 
            return FileResponse(medical_html_path)
        raise HTTPException(status_code=404, detail="medical/persona/index.html not found")
    
    @app.get("/medical/competitive", include_in_schema=False)
    async def serve_medical_competitive(username: str = Depends(verify_department_credentials("medical"))):
        competitive_path = frontend_dir / "medical/competitive/index.html"
        if competitive_path.exists(): 
            return FileResponse(competitive_path)
        raise HTTPException(status_code=404, detail="medical/competitive/index.html not found")
    
    @app.get("/dental", include_in_schema=False)
    async def serve_dental_dashboard(username: str = Depends(verify_department_credentials("dental"))):
        dashboard_path = frontend_dir / "dental/dashboard.html"
        if dashboard_path.exists(): 
            return FileResponse(dashboard_path)
        # フォールバック: 旧index.htmlを使用
        dental_html_path = frontend_dir / "dental/index.html"
        if dental_html_path.exists(): 
            return FileResponse(dental_html_path)
        raise HTTPException(status_code=404, detail="dental/dashboard.html not found")
    
    @app.get("/dental/persona", include_in_schema=False)
    async def serve_dental_persona(username: str = Depends(verify_department_credentials("dental"))):
        persona_path = frontend_dir / "dental/persona/index.html"
        if persona_path.exists(): 
            return FileResponse(persona_path)
        # フォールバック: 旧index.htmlを使用
        dental_html_path = frontend_dir / "dental/index.html"
        if dental_html_path.exists(): 
            return FileResponse(dental_html_path)
        raise HTTPException(status_code=404, detail="dental/persona/index.html not found")
    
    @app.get("/dental/competitive", include_in_schema=False)
    async def serve_dental_competitive(username: str = Depends(verify_department_credentials("dental"))):
        competitive_path = frontend_dir / "dental/competitive/index.html"
        if competitive_path.exists(): 
            return FileResponse(competitive_path)
        raise HTTPException(status_code=404, detail="dental/competitive/index.html not found")
    
    @app.get("/others", include_in_schema=False)
    async def serve_others_dashboard(username: str = Depends(verify_department_credentials("others"))):
        dashboard_path = frontend_dir / "others/dashboard.html"
        if dashboard_path.exists(): 
            return FileResponse(dashboard_path)
        # フォールバック: 旧index.htmlを使用
        others_html_path = frontend_dir / "others/index.html"
        if others_html_path.exists(): 
            return FileResponse(others_html_path)
        raise HTTPException(status_code=404, detail="others/dashboard.html not found")
    
    @app.get("/others/persona", include_in_schema=False)
    async def serve_others_persona(username: str = Depends(verify_department_credentials("others"))):
        persona_path = frontend_dir / "others/persona/index.html"
        if persona_path.exists(): 
            return FileResponse(persona_path)
        # フォールバック: 旧index.htmlを使用
        others_html_path = frontend_dir / "others/index.html"
        if others_html_path.exists(): 
            return FileResponse(others_html_path)
        raise HTTPException(status_code=404, detail="others/persona/index.html not found")
    
    @app.get("/others/competitive", include_in_schema=False)
    async def serve_others_competitive(username: str = Depends(verify_department_credentials("others"))):
        competitive_path = frontend_dir / "others/competitive/index.html"
        if competitive_path.exists(): 
            return FileResponse(competitive_path)
        raise HTTPException(status_code=404, detail="others/competitive/index.html not found")
    
    # ルートアクセスをユーザーダッシュボードへリダイレクト
    @app.get("/", include_in_schema=False)
    async def redirect_root_to_user_dashboard(username: str = Depends(verify_admin_credentials)):
        return RedirectResponse(url="/user/dashboard", status_code=302)
    
    # カテゴリールートからダッシュボードへのリダイレクト
    @app.get("/user", include_in_schema=False)
    async def redirect_user_to_dashboard(username: str = Depends(verify_admin_credentials)):
        return RedirectResponse(url="/user/dashboard", status_code=302)
    
    @app.get("/medical", include_in_schema=False)
    async def redirect_medical_to_dashboard(username: str = Depends(verify_department_credentials("medical"))):
        return RedirectResponse(url="/medical/dashboard", status_code=302)
    
    @app.get("/dental", include_in_schema=False)
    async def redirect_dental_to_dashboard(username: str = Depends(verify_department_credentials("dental"))):
        return RedirectResponse(url="/dental/dashboard", status_code=302)
    
    @app.get("/others", include_in_schema=False)
    async def redirect_others_to_dashboard(username: str = Depends(verify_department_credentials("others"))):
        return RedirectResponse(url="/others/dashboard", status_code=302)
    
    # 診療科別の設定API
    @app.get("/api/departments/by-category/{category}")
    async def get_departments_by_category(category: str, username: str = Depends(verify_admin_credentials)):
        """診療科カテゴリーに応じた診療科リストを返す"""
        departments_by_category_path = project_root_dir / "config" / "departments_by_category.json"
        if not departments_by_category_path.exists():
            raise HTTPException(status_code=404, detail="departments_by_category.json not found")
        
        with open(departments_by_category_path, 'r', encoding='utf-8') as f:
            departments_by_category = json.load(f)
        
        if category not in departments_by_category:
            raise HTTPException(status_code=404, detail=f"Category '{category}' not found")
        
        # departments.jsonから詳細情報を取得
        departments_path = project_root_dir / "config" / "departments.json"
        if not departments_path.exists():
            raise HTTPException(status_code=404, detail="departments.json not found")
        
        with open(departments_path, 'r', encoding='utf-8') as f:
            all_departments = json.load(f)
        
        # カテゴリーに属する診療科のみをフィルタリング
        category_dept_ids = departments_by_category[category]
        filtered_departments = [
            dept for dept in all_departments["departments"]
            if dept["id"] in category_dept_ids and dept.get("enabled", True)
        ]
        
        return {"departments": filtered_departments}
    
    # 主訴リスト取得API（キャッシュ使用版）
    @app.get("/api/chief-complaints/{category}/{department}")
    async def get_chief_complaints_api(category: str, department: str, username: str = Depends(verify_admin_credentials)):
        """指定された診療科の主訴リストを返す（キャッシュから高速取得）"""
        # キャッシュから主訴を取得
        chief_complaints_list = get_chief_complaints(category, department)
        
        if not chief_complaints_list:
            raise HTTPException(
                status_code=404, 
                detail=f"Chief complaints not found for {category}/{department}"
            )
        
        # 主訴リストを返す
        return {
            "category": category,
            "department": department,
            "chief_complaints": chief_complaints_list
        }
    
    print("User UI, Admin UI, and Department routes are set up.")
else:
    print(f"Frontend directory not found at {frontend_dir}. Static file serving skipped.")

# Helper function to find font file
def find_font_file():
    """Search for the Japanese font file in various locations"""
    possible_paths = [
        project_root_dir / "assets" / "fonts" / "ipaexg.ttf",  # New structure
        project_root_dir / "ipaexg.ttf",  # Legacy location
        current_file_dir / "ipaexg.ttf",
        Path("/usr/share/fonts/truetype/fonts-japanese-gothic.ttf"),
        Path("/usr/share/fonts/truetype/fonts-japanese-mincho.ttf"),
    ]
    
    for path in possible_paths:
        if path.exists():
            return str(path)
    
    # Default fallback
    return str(project_root_dir / "assets" / "fonts" / "ipaexg.ttf")

# Get font path
FONT_PATH = find_font_file()
print(f"Using font from: {FONT_PATH}")

# --- Settings Migration ---
def migrate_image_model_settings():
    """Migrates old Gemini image model name in settings to the new one."""
    try:
        settings = crud.read_settings()
        if settings and settings.models and \
           settings.models.image_api_model == "gemini-2.0-flash-preview-image-generation":  # 古い名前をチェック
            print("Migrating image_api_model name in settings...")
            settings.models.image_api_model = "gemini-2.0-flash-exp-image-generation"  # 新しい名前に設定
            crud.write_settings(settings)
            print("Settings migrated: Updated Gemini image model name to gemini-2.0-flash-exp-image-generation")
    except Exception as e:
        print(f"Settings migration for image model failed: {e}")

@app.on_event("startup")
async def startup_event():
    print("Running startup migration...")
    migrate_image_model_settings()
    
    # RAGデータベースの初期化（テーブル作成のみ、CSVは読み込まない）
    # データベースロックエラーを回避するため、エラーハンドリングを追加
    try:
        print("[RAG] Initializing database with lazy loading mode...")
        rag_processor.init_rag_database()
        
        # データベースファイルの存在確認
        from pathlib import Path
        import os
        
        # 実際のRAGプロセッサーと同じパスを使用
        rag_db_path = Path("./app_settings/rag_data.db")
        
        if rag_db_path.exists():
            print(f"[RAG] Database initialized at: {rag_db_path}")
            print("[RAG] CSV data will be loaded on-demand for each department")
        else:
            print(f"[RAG] Warning: Database file not found at: {rag_db_path}")
    except Exception as e:
        print(f"Error initializing RAG database: {e}")
        # データベースロックエラーの場合、アプリケーションを続行
        if "database is locked" in str(e):
            print("[RAG] Database is locked - likely another worker is initializing. Continuing...")
        else:
            # 他のエラーの場合はログに記録して続行
            print(f"[RAG] Continuing without RAG database: {e}")

# --- AI Client Initialization Helper --- 
def get_ai_client(model_name, api_key):
    """Initializes and returns the correct AI client based on model name."""
    if model_name.startswith("gpt"):
        if not OpenAI: raise ValueError("OpenAI library not loaded.")
        if not api_key: raise ValueError("OpenAI APIキーが設定されていません。")
        try: return OpenAI(api_key=api_key, timeout=120.0)  # タイムアウトを120秒に設定
        except Exception as e: raise ValueError(f"OpenAI Client Error: {e}")
    elif model_name.startswith("claude"):
        if not Anthropic: 
            raise ValueError("Anthropic library not loaded.")
        if not api_key: 
            raise ValueError("Anthropic APIキーが設定されていません。")
        try: 
            import httpx
            
            client = Anthropic(
                api_key=api_key,
                timeout=120.0,  # タイムアウトを120秒に設定
                max_retries=2,  # SDK レベルでのリトライ
            )
            return client
        except Exception as e: 
            raise ValueError(f"Anthropic Client Error: {e}")
    elif model_name.startswith("gemini"):
        if not api_key: raise ValueError("Google APIキーが設定されていません。")
        # Prefer new SDK
        if google_genai_sdk:
            try:
                return google_genai_sdk.Client(api_key=api_key)
            except Exception as e:
                pass  # エラーが発生した場合は古いSDKにフォールバック
        # Fallback to old SDK
        if old_gemini_sdk:
            try:
                old_gemini_sdk.configure(api_key=api_key)
                # For the old SDK, GenerativeModel is returned, not a 'Client' instance
                return old_gemini_sdk.GenerativeModel(model_name) 
            except Exception as e_old:
                raise ValueError(f"Google Client Error (old SDK): {e_old}")
        raise ValueError("Google AI SDK (new or old) not available or failed to initialize.")
    else:
        raise ValueError(f"未対応のモデルです: {model_name}")

# --- Function to generate text using AI ---
async def generate_text_response(prompt_text, model_name, api_key):
    """
    Generate text using the specified AI model.
    
    Args:
        prompt_text (str): The prompt to send to the AI model
        model_name (str): The model name (e.g., "gpt-4", "claude-3-opus", "gemini-pro")
        api_key (str): The API key for the AI service
        
    Returns:
        str: The generated text response, or None if generation fails
    """
    import time
    import socket
    import ssl
    
    try:
        # Initialize AI client
        client = get_ai_client(model_name, api_key)
        
        if model_name.startswith("gpt"):
            # OpenAI API call
            # GPT-5の判定とAPI選択
            if "gpt-5" in model_name:
                try:
                    # GPT-5の新しいresponses APIを使用（最小実装）
                    response = client.responses.create(
                        model="gpt-5",
                        input=prompt_text,
                        reasoning={"effort": "high"}  # ペルソナ生成は詳細な推論が必要
                    )
                    return response.output_text
                except Exception as e:
                    print(f"[ERROR] GPT-5 responses API error: {str(e)}")
                    print("[INFO] Trying GPT-5 with chat.completions API")
                    
                    # chat.completions APIでリトライ
                    try:
                        completion = client.chat.completions.create(
                            model=model_name,
                            messages=[{"role": "user", "content": prompt_text}],
                            temperature=1.0,  # GPT-5はデフォルト値のみサポート
                            max_completion_tokens=128000  # GPT-5最大値
                        )
                        return completion.choices[0].message.content
                    except Exception as e2:
                        print(f"[ERROR] GPT-5 chat API error: {str(e2)}")
                        return None
            else:
                # GPT-4以前のモデル
                completion = client.chat.completions.create(
                    model=model_name,
                    messages=[{"role": "user", "content": prompt_text}],
                    temperature=0.7,
                    max_tokens=128000  # 最大値に設定
                )
                return completion.choices[0].message.content
            
        elif model_name.startswith("claude"):
            # Anthropic API call with retry logic
            print(f"[DEBUG] Calling Claude API with model: {model_name}")
            
            # Network diagnostics (optional, can be removed in production)
            if os.environ.get("DEBUG_NETWORK", "false").lower() == "true":
                try:
                    api_host = "api.anthropic.com"
                    print(f"[DEBUG] DNS resolution test: {api_host}")
                    ip_addresses = socket.gethostbyname_ex(api_host)
                    print(f"[DEBUG] DNS resolution successful: IP addresses = {ip_addresses[2]}")
                except Exception as net_error:
                    print(f"[ERROR] Network diagnostic error: {type(net_error).__name__}: {net_error}")
            
            # Retry configuration
            max_retries = 3
            base_delay = 2  # Base delay in seconds
            
            messages_to_send = [{"role": "user", "content": prompt_text}]
            
            # Use direct HTTP if specified in environment
            use_direct_http = os.environ.get("USE_DIRECT_HTTP_FOR_CLAUDE", "false").lower() == "true"
            
            if use_direct_http:
                # Direct HTTP request implementation
                import httpx
                for attempt in range(max_retries):
                    try:
                        if attempt > 0:
                            wait_time = base_delay * (2 ** (attempt - 1))
                            print(f"[INFO] Retry attempt {attempt + 1}/{max_retries} for Claude API after {wait_time}s delay")
                            time.sleep(wait_time)
                        
                        async with httpx.AsyncClient(timeout=httpx.Timeout(120.0)) as http_client:  # タイムアウトを120秒に延長
                            http_response = await http_client.post(
                                "https://api.anthropic.com/v1/messages",
                                headers={
                                    "x-api-key": api_key,
                                    "anthropic-version": "2023-06-01",
                                    "content-type": "application/json"
                                },
                                json={
                                    "model": model_name,
                                    "max_tokens": 2500,  # 502エラー対策で削減
                                    "messages": messages_to_send,
                                    "temperature": 0.7
                                }
                            )
                            
                            if http_response.status_code == 200:
                                response_data = http_response.json()
                                if response_data.get("content") and len(response_data["content"]) > 0:
                                    return response_data["content"][0].get("text", "")
                            else:
                                error_detail = http_response.text
                                print(f"[ERROR] Claude API HTTP error: {http_response.status_code} - {error_detail}")
                                if http_response.status_code in [401, 403, 404]:
                                    raise ValueError(f"Claude API error (status {http_response.status_code}): {error_detail}")
                                
                    except Exception as e:
                        if attempt == max_retries - 1:
                            raise
            else:
                # SDK-based approach with retries
                for attempt in range(max_retries):
                    try:
                        if attempt > 0:
                            wait_time = base_delay * (2 ** (attempt - 1))
                            print(f"[INFO] Retry attempt {attempt + 1}/{max_retries} for Claude API after {wait_time}s delay")
                            time.sleep(wait_time)
                        
                        response = client.messages.create(
                            model=model_name,
                            max_tokens=64000,  # Claude Sonnet 4最大値
                            messages=messages_to_send,
                            temperature=0.7
                        )
                        
                        # Extract text from response
                        if response.content and len(response.content) > 0:
                            if hasattr(response.content[0], 'text'):
                                return response.content[0].text
                                
                    except Exception as e:
                        print(f"[ERROR] Claude API call failed (attempt {attempt + 1}): {type(e).__name__}: {e}")
                        
                        # Check for non-retryable errors
                        if hasattr(e, 'status_code') and e.status_code in [401, 403, 404]:
                            raise
                        
                        # Re-raise on last attempt
                        if attempt == max_retries - 1:
                            raise
                            
        elif model_name.startswith("gemini"):
            # Google AI API call
            if hasattr(client, 'models'):
                # New SDK
                print(f"[DEBUG] Using new Gemini SDK with model: {model_name}")
                response = client.models.generate_content(
                    model=model_name,
                    contents=prompt_text,
                    config=google_genai_types.GenerateContentConfig(
                        temperature=0.7,
                        max_output_tokens=64000  # Gemini 2.5 Pro最大値
                    )
                )
                print(f"[DEBUG] Gemini response received: {response}")
                if response.candidates and response.candidates[0].content.parts:
                    generated_text = response.candidates[0].content.parts[0].text
                else:
                    print(f"[ERROR] Gemini response has no candidates or parts")
                    return f"Gemini APIからの応答が空でした。別のモデルを試してください。"
            else:
                # Old SDK
                print(f"[DEBUG] Using old Gemini SDK")
                generation_config = {
                    "temperature": 0.7,
                    "max_output_tokens": 64000,  # Gemini 2.5 Pro最大値
                }
                response = client.generate_content(prompt_text, generation_config=generation_config)
                generated_text = response.text
            
            # Remove Gemini-specific character count annotations
            if generated_text:
                generated_text = re.sub(r'[（(]\d+文字程度[）)]\s*', '', generated_text)
                generated_text = re.sub(r'\d+文字程度[\s、。]', '', generated_text)
            
            return generated_text
            
        else:
            raise ValueError(f"Unsupported model: {model_name}")
            
    except Exception as e:
        print(f"[ERROR] Text generation failed: {type(e).__name__}: {e}")
        
        # Try fallback to Gemini if Claude fails with connection error
        if model_name.startswith("claude") and ("APIConnectionError" in str(e) or "Connection error" in str(e)):
            try:
                google_api_key = os.getenv("GOOGLE_API_KEY", "").strip()
                if google_api_key:
                    print("[INFO] Falling back to Gemini due to Claude connection error")
                    fallback_model = "gemini-2.0-flash-exp"
                    return await generate_text_response(prompt_text, fallback_model, google_api_key)
            except Exception as fallback_error:
                print(f"[ERROR] Fallback to Gemini also failed: {fallback_error}")
                return f"AI生成に失敗しました（フォールバックも失敗）: {str(fallback_error)}"
        
        return f"AI生成に失敗しました: {str(e)}"

# --- Function to build prompts ---
def build_prompt(data, limit_personality="100", limit_reason="100", limit_behavior="100", 
                 limit_reviews="100", limit_values="100", limit_demands="100", rag_context=""):

    # --- Patient Type Descriptions ---
    patient_type_details_map = {
        '利便性重視型': { "description": 'アクセスの良さ、待ち時間の短さ、診療時間の柔軟性など、便利さを最優先', "example": '忙しいビジネスパーソン、オンライン診療を好む患者' },
        '専門医療追求型': { "description": '専門医や高度専門医療機関での治療を希望し、医師の経歴や実績を重視', "example": '難病患者、複雑な症状を持つ患者' },
        '予防健康管理型': { "description": '病気になる前の予防や早期発見、健康維持に関心が高い', "example": '定期健診を欠かさない人、予防接種に積極的な人' },
        '代替医療志向型': { "description": '漢方、鍼灸、ホメオパシーなど、西洋医学以外の選択肢を積極的に取り入れる', "example": '自然療法愛好者、慢性疾患の患者' },
        '経済合理型': { "description": '自己負担額、保険適用の有無、費用対効果を重視', "example": '経済的制約のある患者、医療費控除を意識する人' },
        '情報探求型': { "description": '徹底的な情報収集、セカンドオピニオン取得、比較検討を行う', "example": '高学歴層、慎重な意思決定を好む患者' },
        '革新技術指向型': { "description": '最先端の医療技術、新薬、臨床試験などに積極的に関心を持つ', "example": '既存治療で効果が出なかった患者、医療イノベーションに関心がある人' },
        '対話重視型': { "description": '医師からの丁寧な説明や対話を求め、質問が多い', "example": '不安を感じやすい患者、医療従事者' },
        '信頼基盤型': { "description": 'かかりつけ医との長期的な関係や医療機関の評判を重視', "example": '地域密着型の患者、同じ医師に長期通院する患者' },
        '緊急解決型': { "description": '症状の即時改善を求め、緊急性を重視', "example": '急性疾患患者、痛みに耐性が低い患者' },
        '受動依存型': { "description": '医師の判断に全面的に依存し、自らの決定より医師の指示を優先', "example": '高齢者、医療知識が少ない患者' },
        '自律決定型': { "description": '自分の治療に主体的に関わり、最終決定権を持ちたいと考える', "example": '医療リテラシーが高い患者、自己管理を好む慢性疾患患者' }
    }

    prompt_parts = [
        "以下の情報に基づいて、医療系のペルソナを作成してください。",
        "各項目は自然な文章で記述し、**日本語で、指定されたおおよその文字数制限に従ってください**。",
        "",
        "# 利用者からの入力情報"
    ]

    # 基本情報
    basic_info = {
        "診療科": data.get('department'),
        "ペルソナ作成目的": data.get('purpose'),
        "名前": data.get('name'),
        "性別": data.get('gender'),
        "年齢": data.get('age'),
        "都道府県": data.get('prefecture'),
        "市区町村": data.get('municipality'),
        "家族構成": data.get('family'),
        "職業": data.get('occupation'),
        "年収": data.get('income'),
        "趣味": data.get('hobby'),
        "ライフイベント": data.get('life_events'),
        "患者タイプ": data.get('patient_type')
    }
    # setting_type と patient_type も追加
    if data.get('setting_type'):
        basic_info["基本情報設定タイプ"] = data.get('setting_type')

    for key, value in basic_info.items():
        if value: # 値が存在する場合のみプロンプトに追加
            prompt_parts.append(f"- {key}: {value}")
            if key == "患者タイプ" and value in patient_type_details_map:
                details = patient_type_details_map[value]
                prompt_parts.append(f"  - 患者タイプの特徴: {details['description']}")
                prompt_parts.append(f"  - 患者タイプの例: {details['example']}")
        # 「選択された患者タイプ」の(自動生成)条件から patient_type を除外 (「患者タイプ」で処理されるため)
        elif key in ["名前", "性別", "年齢", "都道府県", "市区町村", "家族構成", "職業", "年収", "趣味", "ライフイベント"] and key != "患者タイプ": 
             prompt_parts.append(f"- {key}: (自動生成)")
        elif key == "患者タイプ" and not value and data.get('setting_type') == 'patient_type': # setting_typeがpatient_typeで患者タイプが未選択の場合
             prompt_parts.append(f"- {key}: (指定なし/自動生成)")

    # Step 4の固定追加項目
    fixed_additional_info = {
        "座右の銘": data.get('motto'),
        "最近の悩み/関心": data.get('concerns'),
        "好きな有名人/尊敬する人物": data.get('favorite_person'),
        "よく見るメディア/SNS": data.get('media_sns'),
        "性格キーワード": data.get('personality_keywords'),
        "最近した健康に関する行動": data.get('health_actions'),
        "休日の過ごし方": data.get('holiday_activities'),
        "キャッチコピー": data.get('catchphrase'),
    }
    has_fixed_additional_info = any(fixed_additional_info.values())
    if has_fixed_additional_info:
        prompt_parts.append("\n## 追加情報（固定項目）")
        for key, value in fixed_additional_info.items():
            if value:
                prompt_parts.append(f"- {key}: {value}")

    # Step 4の動的追加項目
    additional_field_names = data.get('additional_field_name', [])
    additional_field_values = data.get('additional_field_value', [])
    
    dynamic_additional_info = []
    if isinstance(additional_field_names, list) and isinstance(additional_field_values, list):
        for i in range(len(additional_field_names)):
            field_name = additional_field_names[i]
            field_value = additional_field_values[i] if i < len(additional_field_values) else None
            if field_name and field_value: # 項目名と内容の両方がある場合のみ
                dynamic_additional_info.append(f"- {field_name}: {field_value}")
    
    if dynamic_additional_info:
        prompt_parts.append("\n## 追加情報（自由入力項目）")
        prompt_parts.extend(dynamic_additional_info)

    prompt_parts.append("\n# 生成項目")
    prompt_parts.append("以下の項目について、上記情報に基づいた自然な文章を生成してください。")
    prompt_parts.append("\n## 出力形式の重要な指示:")
    prompt_parts.append("- 必ず以下の形式で出力してください")
    prompt_parts.append("- 各項目は「番号. **項目名**: 内容」の形式で記述")
    prompt_parts.append("- 内容は項目名の後のコロン（:）の直後に続けて記述")
    prompt_parts.append("- 文字数の指定（例：「100文字程度」）は出力に含めない")
    prompt_parts.append("- ペルソナ名などの余分なヘッダーは含めない")
    prompt_parts.append("")
    # 主訴は必須項目
    chief_complaint = data.get('chief_complaint')
    
    prompt_parts.append("## 生成する項目（各項目を指定文字数で）:")
    prompt_parts.append(f"1. **個性（価値観・人生観）**: {limit_personality}文字程度の内容をここに記述")
    prompt_parts.append(f"2. **病院に行く理由（主訴「{chief_complaint}」を踏まえて）**: {limit_reason}文字程度の内容をここに記述")
    prompt_parts.append(f"3. **症状のパターンや受診頻度**: {limit_behavior}文字程度の内容をここに記述")
    prompt_parts.append(f"4. **口コミを見る際に重要視すること**: {limit_reviews}文字程度の内容をここに記述")
    prompt_parts.append(f"5. **医療機関に対する価値観や行動傾向**: {limit_values}文字程度の内容をここに記述")
    prompt_parts.append(f"6. **医療機関に求めるもの**: {limit_demands}文字程度の内容をここに記述")
    
    prompt_parts.append("\n注意事項:")
    prompt_parts.append(f"- 主訴「{chief_complaint}」を必ず反映させてください")
    prompt_parts.append("- 各項目は指定された文字数で簡潔に記述してください")
    prompt_parts.append("- 具体的でリアルな内容にしてください")
    
    return "\n".join(prompt_parts)

# --- Function to parse AI response ---
def parse_ai_response(text):
    """Extract sections from AI-generated text response."""
    sections = {
        "personality": "",
        "reason": "",
        "behavior": "",
        "reviews": "",
        "values": "",
        "demands": ""
    }
    
    if not text:
        return sections
    
    # Try to parse structured output
    try:
        lines = text.strip().split('\n')
        current_section = None
        
        
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
                
            # Skip persona name headers (e.g., **ペルソナ：本田 瑞季**)
            if line.startswith("**") and "ペルソナ" in line and line.endswith("**"):
                continue
                
            # Check if this is a section header for personality section
            # Support both numbered format and **header** format
            # Also support formats without colons and with content on next line
            if (line.startswith("1.") and ("個性" in line or "価値観" in line or "人生観" in line)) or \
               (line.startswith("**1.") and ("個性" in line or "価値観" in line or "人生観" in line)) or \
               (line.startswith("**") and ("個性" in line or "価値観" in line or "人生観" in line) and line.endswith("**")):
                current_section = "personality"
                # Check if content is on the same line after colon
                content = line.split(':', 1)[1].strip() if ':' in line else ""
                if content:
                    # 「(100文字程度)」のような文字数指定を削除
                    content = re.sub(r'[（(]\d+文字程度[）)]\s*', '', content)
                    content = re.sub(r'\d+文字程度\s*', '', content)
                    sections[current_section] = content
                continue
                
            # Check for reason section
            elif (line.startswith("2.") and "病院に行く理由" in line) or \
                 (line.startswith("**2.") and "病院に行く理由" in line) or \
                 (line.startswith("**") and "病院に行く理由" in line and line.endswith("**")):
                current_section = "reason"
                content = line.split(':', 1)[1].strip() if ':' in line else ""
                if content:
                    content = re.sub(r'\(\d+文字程度\)\s*', '', content)
                    sections[current_section] = content
                continue
                
            # Check for behavior section
            elif (line.startswith("3.") and ("症状" in line or "通院頻度" in line or "行動パターン" in line)) or \
                 (line.startswith("**3.") and ("症状" in line or "通院頻度" in line or "行動パターン" in line)) or \
                 (line.startswith("**") and ("症状" in line or "通院頻度" in line or "行動パターン" in line) and line.endswith("**")):
                current_section = "behavior"
                content = line.split(':', 1)[1].strip() if ':' in line else ""
                if content:
                    content = re.sub(r'\(\d+文字程度\)\s*', '', content)
                    sections[current_section] = content
                continue
                
            # Check for reviews section
            elif (line.startswith("4.") and ("口コミ" in line or "重視ポイント" in line)) or \
                 (line.startswith("**4.") and ("口コミ" in line or "重視ポイント" in line)) or \
                 (line.startswith("**") and ("口コミ" in line or "重視ポイント" in line) and line.endswith("**")):
                current_section = "reviews"
                content = line.split(':', 1)[1].strip() if ':' in line else ""
                if content:
                    content = re.sub(r'\(\d+文字程度\)\s*', '', content)
                    sections[current_section] = content
                continue
                
            # Check for values section
            elif (line.startswith("5.") and ("医療機関" in line or "価値観" in line or "行動傾向" in line)) or \
                 (line.startswith("**5.") and ("医療機関" in line or "価値観" in line or "行動傾向" in line)) or \
                 (line.startswith("**") and ("医療機関" in line and ("価値観" in line or "行動傾向" in line)) and line.endswith("**")):
                current_section = "values"
                content = line.split(':', 1)[1].strip() if ':' in line else ""
                if content:
                    content = re.sub(r'\(\d+文字程度\)\s*', '', content)
                    sections[current_section] = content
                continue
                
            # Check for demands section
            elif (line.startswith("6.") and ("医療機関" in line or "求めるもの" in line)) or \
                 (line.startswith("**6.") and ("医療機関" in line or "求めるもの" in line)) or \
                 (line.startswith("**") and ("医療機関" in line and "求めるもの" in line) and line.endswith("**")):
                current_section = "demands"
                content = line.split(':', 1)[1].strip() if ':' in line else ""
                if content:
                    content = re.sub(r'\(\d+文字程度\)\s*', '', content)
                    sections[current_section] = content
                continue
                
            # If we're in a section, append text (skip headers and numbered lines)
            if current_section and not line.startswith(("#", "##", "###", "**")) and not re.match(r'^\d+\.', line):
                # 「(100文字程度)」のような文字数指定を削除
                cleaned_line = re.sub(r'[（(]\d+文字程度[）)]\s*', '', line)
                cleaned_line = re.sub(r'\d+文字程度\s*', '', cleaned_line)
                # Remove extra asterisks that might be in the content
                cleaned_line = cleaned_line.strip('*').strip()
                if cleaned_line:
                    if sections[current_section]:
                        sections[current_section] += " " + cleaned_line
                    else:
                        sections[current_section] = cleaned_line
    
    except Exception as e:
        print(f"Error parsing AI response: {e}")
        traceback.print_exc()
    
    # If all sections are empty, try a more flexible parsing approach
    if all(not v for v in sections.values()):
        print("[DEBUG] All sections empty, trying flexible parsing")
        # Try to find content by keywords anywhere in the text
        keywords_map = {
            "personality": ["個性", "価値観", "人生観", "性格"],
            "reason": ["病院に行く理由", "来院理由", "受診理由"],
            "behavior": ["症状", "通院頻度", "行動パターン"],
            "reviews": ["口コミ", "重視ポイント", "評価"],
            "values": ["医療機関への価値観", "医療に対する価値観"],
            "demands": ["医療機関に求めるもの", "希望", "要望"]
        }
        
        for section, keywords in keywords_map.items():
            for keyword in keywords:
                if keyword in text:
                    # Find the content after the keyword
                    pattern = f"{keyword}[：:]*\\s*([^\\n]+)"
                    match = re.search(pattern, text)
                    if match:
                        content = match.group(1).strip()
                        content = re.sub(r'[（(]\d+文字程度[）)]\s*', '', content)
                        content = content.strip('*').strip()
                        if content:
                            sections[section] = content
                            break
    
    # Debug output
    print(f"[DEBUG] Parsed sections:")
    for key, value in sections.items():
        print(f"  {key}: {value[:50]}..." if value else f"  {key}: (empty)")
    
    return sections

@app.post("/api/generate")
async def generate_persona(request: Request, username: str = Depends(verify_admin_credentials)):
    try:
        data = await request.json()
        
        # --- 設定をcrudから読み込む ---
        app_settings = crud.read_settings() # AdminSettings インスタンスが返る
        
        selected_text_model = app_settings.models.text_api_model if app_settings.models else "gpt-5" # デフォルト値をGPT-5に
        selected_image_model = app_settings.models.image_api_model if app_settings.models else "dall-e-3" # デフォルト値
        
        # ===== モデル使用ログ =====
        print("="*60)
        print("[MODEL] ===== AI MODEL SELECTION =====")
        print(f"[MODEL] Text Generation Model: {selected_text_model}")
        print(f"[MODEL] Image Generation Model: {selected_image_model}")
        print("="*60)

        # --- APIキーの取得 ---
        openai_api_key = os.environ.get("OPENAI_API_KEY", "").strip()
        anthropic_api_key = os.environ.get("ANTHROPIC_API_KEY", "").strip()
        google_api_key = os.environ.get("GOOGLE_API_KEY", "").strip()

        # テキスト生成用APIキーの選択
        text_api_key_to_use = None
        if selected_text_model.startswith("gpt"):
            text_api_key_to_use = openai_api_key
        elif selected_text_model.startswith("claude"):
            text_api_key_to_use = anthropic_api_key
        elif selected_text_model.startswith("gemini"):
            text_api_key_to_use = google_api_key
        
        # テスト用ダミーデータ (テキスト生成) - APIキーがない場合
        if not text_api_key_to_use and not (selected_text_model.startswith("gemini") and google_api_key): # Geminiはキーがなくてもgenai.configureされるため別途判定
             if selected_text_model != "dummy": # "dummy" モデル選択時以外でキーがない場合
                print(f"Warning: API key for {selected_text_model} not found. Text generation might fail or use defaults.")
                # ここでエラーを返すか、ダミーデータを返すかは設計次第
                # 今回はget_ai_clientやAPI呼び出しでエラーになることを許容


        # 文字数制限を取得（管理画面の設定を優先）
        if app_settings.limits:
            limit_personality = app_settings.limits.get("personality", "100")
            limit_reason = app_settings.limits.get("reason", "100")
            limit_behavior = app_settings.limits.get("behavior", "100")
            limit_reviews = app_settings.limits.get("reviews", "100")
            limit_values = app_settings.limits.get("values", "100")
            limit_demands = app_settings.limits.get("demands", "100")
        else:
            # フォールバック：環境変数から取得
            limit_personality = os.environ.get("LIMIT_PERSONALITY", "100")
            limit_reason = os.environ.get("LIMIT_REASON", "100")
            limit_behavior = os.environ.get("LIMIT_BEHAVIOR", "100")
            limit_reviews = os.environ.get("LIMIT_REVIEWS", "100")
            limit_values = os.environ.get("LIMIT_VALUES", "100")
            limit_demands = os.environ.get("LIMIT_DEMANDS", "100")
        
        # RAGデータベースは起動時に一度だけ初期化されているので、ここでは不要
        # rag_processor.init_rag_database()  # コメントアウト - 遅延読み込みのため不要
        
        # RAGデータの検索
        rag_context = ""
        rag_results = []
        department = data.get('department')
        age = data.get('age')
        gender = data.get('gender')
        chief_complaint = data.get('chief_complaint')  # 主訴を取得
        
        if department:
            # 年齢グループの判定
            age_group = None
            if age:
                try:
                    age_str = str(age)
                    # 年齢を数値に変換
                    if 'y' in age_str:
                        # "45y3m" -> 45
                        age_num = int(age_str.split('y')[0])
                    elif age_str.isdigit():
                        # "45" -> 45
                        age_num = int(age_str)
                    else:
                        age_num = None
                    
                    if age_num is not None:
                        if age_num < 20:
                            age_group = "10s"
                        elif age_num < 30:
                            age_group = "20s"
                        elif age_num < 40:
                            age_group = "30s"
                        elif age_num < 50:
                            age_group = "40s"
                        elif age_num < 60:
                            age_group = "50s"
                        elif age_num < 70:
                            age_group = "60s"
                        else:
                            age_group = "70s"
                except (ValueError, AttributeError):
                    print(f"Warning: Could not parse age '{age}' for RAG search")
            
            # RAGデータの検索（主訴を含めて検索）
            print("="*60)
            print("[RAG] ===== STARTING RAG DATA SEARCH =====")
            print(f"[RAG] Department: {department}")
            print(f"[RAG] Chief Complaint: {chief_complaint}")
            print(f"[RAG] Age Group: {age_group}")
            print(f"[RAG] Gender: {gender}")
            print("="*60)
            
            rag_results = rag_processor.search_rag_data(
                specialty=department,
                age_group=age_group,
                gender=gender,
                chief_complaint=chief_complaint,  # 主訴を追加
                limit=5
            )
            print(f"[RAG] Search completed: {len(rag_results)} results found")
            
            if rag_results:
                rag_context = "\n\n# 参考情報（この診療科・主訴の患者が検索するキーワード）\n"
                if chief_complaint:
                    rag_context += f"以下は、{department}で「{chief_complaint}」に関連し、同じ年代・性別の患者がよく検索するキーワードです。ペルソナ作成の参考にしてください：\n"
                else:
                    rag_context += "以下は、同じ診療科・年代・性別の患者がよく検索するキーワードです。ペルソナ作成の参考にしてください：\n"
                for i, result in enumerate(rag_results, 1):
                    rag_context += f"{i}. {result['keyword']} (検索数: {result['search_volume']}人)\n"
                    # カテゴリーフィールドは削除（CSVに存在しないため）
                
                # RAGデータ使用の詳細ログ
                print("="*60)
                print("[RAG] ===== RAG DATA SUCCESSFULLY LOADED =====")
                print(f"[RAG] Using data for:")
                print(f"[RAG]   - Department: {department}")
                print(f"[RAG]   - Chief Complaint: {chief_complaint if chief_complaint else 'Not specified'}")
                print(f"[RAG]   - Age Group: {age_group if age_group else 'All ages'}")
                print(f"[RAG]   - Gender: {gender if gender else 'All genders'}")
                print(f"[RAG] Top {len(rag_results)} Keywords Found:")
                for i, result in enumerate(rag_results, 1):
                    print(f"[RAG]   {i}. {result['keyword']} (検索数: {result['search_volume']}人)")
                print("[RAG] This data will be included in the AI prompt")
                print("="*60)
            else:
                print("="*60)
                print(f"[RAG] ⚠️ WARNING: No RAG data found")
                print(f"[RAG] Searched for: {department} / {chief_complaint}")
                print("[RAG] Proceeding without RAG context")
                print("="*60)
        
        # RAGデータベース情報を追加（フロントエンドで表示するため）
        rag_info = {
            "database_path": str(rag_processor.RAG_DB_PATH),
            "is_local": "app_settings" in str(rag_processor.RAG_DB_PATH),
            "results_count": len(rag_results) if rag_results else 0,
            "department": department
        }
        
        # 入力データの内容をログ出力
        print(f"[DEBUG] Input data keys: {list(data.keys())}")
        print(f"[DEBUG] Input data sample:")
        for key in ['department', 'purpose', 'name', 'gender', 'age', 'occupation']:
            if key in data:
                print(f"  - {key}: {data[key]}")
        
        # プロンプト構築（RAGコンテキストを含む）
        char_limits = {
            "personality": limit_personality,
            "reason": limit_reason,
            "behavior": limit_behavior,
            "reviews": limit_reviews,
            "values": limit_values,
            "demands": limit_demands
        }
        # build_prompt関数を使用（主訴対応済み）
        prompt_text = build_prompt(
            data,
            limit_personality=str(char_limits["personality"]),
            limit_reason=str(char_limits["reason"]),
            limit_behavior=str(char_limits["behavior"]),
            limit_reviews=str(char_limits["reviews"]),
            limit_values=str(char_limits["values"]),
            limit_demands=str(char_limits["demands"]),
            rag_context=rag_context
        )
        
        
        # ===== 非同期処理開始 =====
        # 画像生成タスクを先に開始（バックグラウンドで実行）
        print("[Async] Starting image generation task in background")
        image_generation_task = asyncio.create_task(
            generate_image_async(
                data=data,
                selected_image_model=selected_image_model,
                openai_api_key=openai_api_key,
                google_api_key=google_api_key
            )
        )
        
        # AIクライアント初期化 (テキスト生成用)
        text_generation_client = None
        client_init_error = None
        if text_api_key_to_use or (selected_text_model.startswith("gemini") and google_api_key):
            try:
                text_generation_client = get_ai_client(selected_text_model, text_api_key_to_use)
            except Exception as e:
                client_init_error = str(e)
        
        # テキスト生成実行（画像生成と並列）
        print("="*60)
        print(f"[AI] ===== STARTING TEXT GENERATION =====")
        print(f"[AI] Model: {selected_text_model}")
        print(f"[AI] RAG Context: {'Yes' if rag_context else 'No'}")
        print(f"[AI] Prompt Length: {len(prompt_text)} characters")
        print("="*60)
        
        generated_text_str = None
        if text_generation_client:
            try:
                generated_text_str = await generate_text_response(prompt_text, selected_text_model, text_api_key_to_use)
                print("="*60)
                print(f"[AI] ✓ Text generation completed successfully")
                print(f"[AI] Response length: {len(generated_text_str) if generated_text_str else 0} characters")
                print("="*60)
            except Exception as e:
                print("="*60)
                print(f"[AI] ✗ ERROR: Text generation failed")
                print(f"[AI] Error type: {type(e).__name__}")
                print(f"[AI] Error message: {e}")
                print("="*60)
                generated_text_str = None
        
        if generated_text_str is None: # AI生成失敗またはスキップ時のフォールバック
            error_msg = "Text generation failed or was skipped."
            if client_init_error:
                error_msg += f" Client initialization error: {client_init_error}"
            generated_details = {
                "personality": "真面目で責任感が強く、家族を大切にする。健康意識が高く、予防医療に関心がある。",
                "reason": "定期的な健康診断と、軽度の高血圧の管理のため。",
                "behavior": "3ヶ月に一度の定期検診に欠かさず通院。処方された降圧剤を規則正しく服用している。",
                "reviews": "医師の説明がわかりやすいこと、待ち時間が短いこと、スタッフの対応が丁寧であることを重視する。",
                "values": "信頼できる医師との長期的な関係を望む。予防医療に前向きで、医師のアドバイスを真摯に受け止める。",
                "demands": "わかりやすい説明と、必要に応じて専門医への適切な紹介。予防医療のアドバイスも欲しい。"
            }
        else:
            # Debug: Log the raw AI response
            print(f"[DEBUG] Raw AI response from {selected_text_model}:")
            print(f"{generated_text_str[:500]}..." if len(generated_text_str) > 500 else generated_text_str)
            generated_details = parse_ai_response(generated_text_str)
        
        # 画像生成の完了を待つ
        print("[Async] Waiting for image generation to complete")
        try:
            image_url = await image_generation_task
            print(f"[Async] Image generation completed successfully")
        except Exception as e:
            print(f"[Async] Image generation task failed: {e}")
            traceback.print_exc()
            image_url = "https://placehold.jp/300x200/FF0000/FFFFFF?text=Async+Error"
        
        # ===== 非同期処理完了 =====
        print(f"[Async] Both text and image generation completed")
        
        # レスポンスデータ作成
        response_data = {
            "profile": data, # フロントから送られてきた入力データをそのまま返す
            "details": generated_details,
            "image_url": image_url, # DALL-EならURL、GeminiならBase64 Data URI
            "rag_info": rag_info if 'rag_info' in locals() else None # RAGデータベース情報
        }
        return response_data

    except Exception as e:
        
        # エラーメッセージの詳細化
        error_message = f"サーバーエラー: {str(e)}"
        error_details = {
            "error": error_message,
            "error_type": type(e).__name__
        }
        
        # HTTPステータスコードがある場合は502として返す
        status_code = 500
        if hasattr(e, 'status_code'):
            if e.status_code == 502:
                status_code = 502
                error_details["error"] = "AI API サービスへの接続に失敗しました。しばらく待ってから再度お試しください。"
                error_details["original_error"] = str(e)
        
        # Claude固有のエラーメッセージ
        if 'selected_text_model' in locals() and selected_text_model and "claude" in selected_text_model.lower():
            if "APIConnectionError" in str(e) or "Connection error" in str(e):
                error_details["error"] = "Claude APIへの接続に失敗しました。Renderの無料プランを使用している場合、サービスがスピンダウンしている可能性があります。"
                error_details["model"] = selected_text_model
                error_details["suggestion"] = "数秒後に再試行するか、他のモデル（GPT、Gemini）を使用してください。"
                error_details["render_note"] = "Renderの無料プランでは15分間アイドル後にサービスが停止し、次回リクエスト時に30-60秒の起動時間が必要です。"
            elif "502" in str(e) or "Bad Gateway" in str(e):
                error_details["error"] = "Claude APIのゲートウェイエラーが発生しました。"
                error_details["model"] = selected_text_model
                error_details["suggestion"] = "APIサービスが一時的に利用できない可能性があります。"
        
        return JSONResponse(
            status_code=status_code,
            content=error_details
        )

# [削除済み] /api/generate-by-complaintエンドポイントは/api/generateに統合されました

@app.post("/api/download/pdf")
async def download_pdf(request: Request, username: str = Depends(verify_admin_credentials)):
    """ペルソナデータをPDFとしてダウンロードするエンドポイント"""
    try:
        data = await request.json()
        if not data:
            return JSONResponse(
                status_code=400,
                content={"error": "No data provided"}
            )
        

        # PDF生成
        try:
            pdf_buffer = generate_pdf(data)
        except Exception as pdf_error:
            print(f"[ERROR] PDF generation failed: {pdf_error}")
            import traceback
            traceback.print_exc()
            return JSONResponse(
                status_code=500,
                content={"error": f"PDF generation failed: {str(pdf_error)}"}
            )
        
        # ファイル名の設定
        profile_name = data.get('profile', {}).get('name', 'persona')
        safe_profile_name = "".join(c if c.isalnum() or c in [' ', '(', ')'] else '_' for c in profile_name)
        filename_utf8 = f"{safe_profile_name}_persona.pdf"
        filename_encoded = quote(filename_utf8)

        # レスポンス作成
        response = Response(pdf_buffer.getvalue(), media_type='application/pdf')
        response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{filename_encoded}"
        
        return response
        
    except Exception as e:
        print(f"Error generating PDF: {e}")
        traceback.print_exc()
        return JSONResponse(
            status_code=500,
            content={"error": f"Failed to generate PDF: {str(e)}"}
        )

@app.post("/api/download/ppt")
async def download_ppt(request: Request, username: str = Depends(verify_admin_credentials)):
    """ペルソナデータをPPTXとしてダウンロードするエンドポイント"""
    try:
        data = await request.json()
        if not data:
            return JSONResponse(
                status_code=400,
                content={"error": "No data provided"}
            )
        
        # PPTXデータ作成のためのデータ変換
        persona_data = {}
        for key, value in data.get('profile', {}).items():
            persona_data[key] = value
        
        # details からの転送
        for key, value in data.get('details', {}).items():
            persona_data[key] = value
        
        # timeline_analysisデータを転送
        if 'timeline_analysis' in data:
            persona_data['timeline_analysis'] = data['timeline_analysis']
            print(f"[DEBUG] PPT: timeline_analysis included with keys: {list(data['timeline_analysis'].keys())}")
        
        # 画像URL
        image_url = data.get('image_url')
        image_path = None
        
        # 診療科と目的の取得
        department_val = persona_data.get('department', '-')
        print(f"[DEBUG] Department value: {department_val}, type: {type(department_val)}")
        # department_valが文字列でない場合も考慮
        if department_val and department_val != '-':
            department_text = DEPARTMENT_MAP.get(str(department_val).lower(), str(department_val))
            print(f"[DEBUG] Department text: {department_text}")
        else:
            department_text = '-'
        
        purpose_val = persona_data.get('purpose', '-')
        # purpose_valが文字列でない場合も考慮
        if purpose_val and purpose_val != '-':
            purpose_text = PURPOSE_MAP.get(str(purpose_val).lower(), str(purpose_val))
        else:
            purpose_text = '-'
        
        # 画像があれば一時ファイルに保存
        if image_url:
            try:
                import tempfile
                import requests
                from PIL import Image
                
                # Data URLの場合
                if image_url.startswith('data:'):
                    # data:image/png;base64,xxxxx の形式から画像データを抽出
                    try:
                        header, base64_data = image_url.split(',', 1)
                        image_data = base64.b64decode(base64_data)
                        
                        # MIME typeから拡張子を判定
                        if 'jpeg' in header or 'jpg' in header:
                            suffix = '.jpg'
                        elif 'webp' in header:
                            suffix = '.webp'
                        else:
                            suffix = '.png'
                        
                        # 一時ファイルに保存
                        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as temp_file:
                            temp_file.write(image_data)
                            image_path = temp_file.name
                    except (ValueError, base64.binascii.Error) as e:
                        print(f"Error decoding base64 image: {e}")
                        image_path = None
                        
                # HTTPSやHTTPのURLの場合
                elif image_url.startswith(('http://', 'https://')):
                    response = requests.get(image_url, timeout=10, headers={'User-Agent': 'Mozilla/5.0'})
                    response.raise_for_status()
                    
                    # Content-Typeから拡張子を判定
                    content_type = response.headers.get('content-type', '').lower()
                    if 'jpeg' in content_type or 'jpg' in content_type:
                        suffix = '.jpg'
                    elif 'webp' in content_type:
                        suffix = '.webp'
                    else:
                        suffix = '.png'
                    
                    # 一時ファイルに保存
                    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as temp_file:
                        temp_file.write(response.content)
                        image_path = temp_file.name
                
            except Exception as e:
                print(f"Error downloading image: {e}")
                image_path = None
        
        # PPTX生成
        try:
            pptx_buffer = generate_ppt(persona_data, image_path, department_text, purpose_text)
        except Exception as ppt_error:
            print(f"[ERROR] PPT generation failed: {ppt_error}")
            import traceback
            traceback.print_exc()
            return JSONResponse(
                status_code=500,
                content={"error": f"PPT generation failed: {str(ppt_error)}"}
            )
        
        # ファイル名の設定
        profile_name = data.get('profile', {}).get('name', 'persona')
        safe_profile_name = "".join(c if c.isalnum() or c in [' ', '(', ')'] else '_' for c in profile_name)
        filename_utf8 = f"{safe_profile_name}_persona.pptx"
        filename_encoded = quote(filename_utf8)
        
        # レスポンス作成
        response = Response(
            pptx_buffer.getvalue(),
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{filename_encoded}"
        
        return response
        
    except Exception as e:
        print(f"Error generating PPTX: {e}")
        traceback.print_exc()
        return JSONResponse(
            status_code=500,
            content={"error": f"Failed to generate PPTX: {str(e)}"}
        )
    finally:
        # 一時ファイルのクリーンアップ（エラーが発生しても必ず実行）
        if 'image_path' in locals() and image_path and os.path.exists(image_path):
            try:
                os.unlink(image_path)
            except Exception as e:
                pass  # エラーが発生しても無視して続行

@app.get("/health", summary="Health check endpoint", tags=["Health"])
async def health_check():
    """
    Simple health check endpoint to confirm the API is running.
    """
    return {"status": "ok"}

@app.post("/api/search-timeline")
async def get_search_timeline(request: Request, username: str = Depends(verify_admin_credentials)):
    """検索キーワードの時系列データを取得"""
    try:
        data = await request.json()
        
        # 必須パラメータの確認
        department = data.get('department')
        chief_complaint = data.get('chief_complaint')
        
        if not department or not chief_complaint:
            return JSONResponse(
                status_code=400,
                content={"error": "department and chief_complaint are required"}
            )
        
        # オプションパラメータ
        gender = data.get('gender')
        age = data.get('age')
        
        # 時系列データを分析
        result = timeline_analyzer.analyze_search_timeline(
            department=department,
            chief_complaint=chief_complaint,
            gender=gender,
            age=age
        )
        
        return result
        
    except Exception as e:
        print(f"Error in search-timeline: {e}")
        return JSONResponse(
            status_code=500,
            content={"error": f"Failed to analyze timeline: {str(e)}"}
        )

@app.post("/api/search-timeline-analysis")
async def analyze_search_behavior(request: Request, username: str = Depends(verify_admin_credentials)):
    """検索行動をAIで分析"""
    try:
        data = await request.json()
        
        # 必須パラメータの確認
        filtered_keywords = data.get('filtered_keywords', [])
        persona_profile = data.get('persona_profile', {})
        
        if not filtered_keywords or not persona_profile:
            return JSONResponse(
                status_code=400,
                content={"error": "filtered_keywords and persona_profile are required"}
            )
        
        # 診断前後のキーワードを分類
        pre_diagnosis = [k for k in filtered_keywords if k['time_diff_days'] < 0]
        post_diagnosis = [k for k in filtered_keywords if k['time_diff_days'] >= 0]
        
        # キーワード分析を実行
        from .services.keyword_analyzer import analyze_search_patterns, extract_emotional_keywords, calculate_demographic_match
        pattern_analysis = analyze_search_patterns(pre_diagnosis, post_diagnosis)
        emotional_keywords = extract_emotional_keywords(filtered_keywords)
        
        # デモグラフィックマッチ度を計算
        demographic_match = calculate_demographic_match(
            filtered_keywords, 
            persona_profile.get('gender', ''),
            persona_profile.get('age', '')
        )
        
        # プロンプトを構築
        prompt = f"""あなたは医療マーケティングの専門家です。
提供されたペルソナの全情報と検索行動データを総合的に分析し、深い洞察を提供してください。

【ペルソナ詳細情報】
■ 基本属性
- 性別：{persona_profile.get('gender')}、年齢：{persona_profile.get('age')}
- 職業：{persona_profile.get('occupation')}、年収：{persona_profile.get('income')}
- 居住地：{persona_profile.get('prefecture')}{persona_profile.get('municipality', '')}

■ 性格・価値観
- 性格：{persona_profile.get('personality_keywords')}
- 座右の銘：{persona_profile.get('motto')}
- キャッチフレーズ：{persona_profile.get('catchphrase')}
- 患者タイプ：{persona_profile.get('patient_type')}

■ ライフスタイル
- 趣味：{persona_profile.get('hobby')}
- 休日の過ごし方：{persona_profile.get('holiday_activities')}
- 利用メディア/SNS：{persona_profile.get('media_sns')}
- 好きな有名人：{persona_profile.get('favorite_person')}
- 家族関係：{persona_profile.get('family')}
- 健康への取り組み：{persona_profile.get('health_actions')}
- 悩み事：{persona_profile.get('concerns')}
- 最近のライフイベント：{persona_profile.get('life_events')}

■ 生成された詳細情報
- 個性・性格：{persona_profile.get('personality')}
- 来院理由：{persona_profile.get('reason')}
- 医療機関への期待：{persona_profile.get('demands')}
- 口コミの見方：{persona_profile.get('reviews')}
- 医療に対する価値観：{persona_profile.get('values')}
- 行動特性：{persona_profile.get('behavior')}

【{persona_profile.get('chief_complaint')}に関する検索行動データ】
診断前の検索（時系列順）：
"""
        
        for k in sorted(pre_diagnosis, key=lambda x: x['time_diff_days']):
            prompt += f"- {k['keyword']} ({abs(k['time_diff_days']):.1f}日前、推定{k['estimated_volume']}人)\n"
        
        prompt += "\n診断後の検索（時系列順）：\n"
        for k in sorted(post_diagnosis, key=lambda x: x['time_diff_days']):
            prompt += f"- {k['keyword']} ({k['time_diff_days']:.1f}日後、推定{k['estimated_volume']}人)\n"
        
        # 検索パターン分析を追加
        prompt += f"\n【検索行動の分析結果】\n"
        prompt += f"- 診断前の主な関心: {', '.join(pattern_analysis['pre_diagnosis_focus']) or 'なし'}\n"
        prompt += f"- 診断後の主な関心: {', '.join(pattern_analysis['post_diagnosis_focus']) or 'なし'}\n"
        prompt += f"- 緊急度: {pattern_analysis['urgency_level']}\n"
        prompt += f"- 解決志向度: {pattern_analysis['solution_seeking_rate']:.1%}\n"
        prompt += f"- ペルソナ属性との一致度: {demographic_match:.1%}\n"
        if emotional_keywords:
            prompt += f"- 感情的キーワード: {', '.join(emotional_keywords[:5])}\n"
        
        prompt += f"""
【分析項目】
以下の形式で簡潔に回答してください：

心理変化の分析：
{persona_profile.get('chief_complaint', '')}診断前後の心理状態の変化を、性格、価値観、家族関係、ライフイベントなどから分析（200文字程度）

隠れたニーズ：
職業、年収、趣味、悩み事、健康への取り組みと検索行動から、潜在的なニーズを洞察（200文字程度）

マーケティング提案：
性格、価値観、ライフスタイル、経済状況を考慮し、{persona_profile.get('department', '')}がとるべき具体的な戦略（200文字程度）

重要：
- 各項目は「項目名：」の後に内容を書く形式にすること
- 番号、記号（###、-など）、文字数表記は含めない
- 前置きや挨拶は不要"""
        
        # AI分析を実行
        try:
            # モデル設定を取得
            from .api.admin_settings import get_admin_settings
            app_settings = await get_admin_settings()
            selected_text_model = app_settings.models.text_api_model or os.getenv("DEFAULT_TEXT_MODEL", "gpt-5-2025-08-07")
            
            # モデル使用ログ
            print("="*60)
            print("[TimelineAnalysis] ===== GENERATING AI ANALYSIS =====")
            print(f"[TimelineAnalysis] Model: {selected_text_model}")
            print(f"[TimelineAnalysis] Keywords to analyze: {len(filtered_keywords)}")
            print("="*60)
            
            # モデルに応じて適切なAPIキーを選択
            if "gemini" in selected_text_model.lower():
                api_key = os.getenv("GOOGLE_API_KEY")
            elif "claude" in selected_text_model.lower():
                api_key = os.getenv("ANTHROPIC_API_KEY")
            else:
                api_key = os.getenv("OPENAI_API_KEY")
            
            # テキスト生成
            analysis_result = await generate_text_response(
                prompt_text=prompt,
                model_name=selected_text_model,
                api_key=api_key
            )
            
            return {
                "ai_analysis": analysis_result,
                "keywords_analyzed": len(filtered_keywords),
                "model_used": selected_text_model
            }
            
        except Exception as ai_error:
            print(f"AI analysis error: {ai_error}")
            return JSONResponse(
                status_code=500,
                content={"error": f"AI analysis failed: {str(ai_error)}"}
            )
            
    except Exception as e:
        print(f"Error in search-timeline-analysis: {e}")
        return JSONResponse(
            status_code=500,
            content={"error": f"Failed to analyze behavior: {str(e)}"}
        )


# To run this app locally (from the 'backend' directory, assuming main.py is in 'app'):
# Ensure 'app' is a package (has __init__.py if needed, though often not for FastAPI structure like this)
# cd backend
# python -m uvicorn app.main:app --reload --port 8000

# If your structure is flatter, e.g. main.py in 'backend' root:
# cd backend
# uvicorn main:app --reload --port 8000 

# --- Global Maps and Helper Functions for PDF/PPT Generation ---
GENDER_MAP = {
    "male": "男性", "female": "女性", "other": "その他",
}
INCOME_MAP = {
    "<100": "100万円未満", "100-200": "100-200万円", "200-300": "200-300万円",
    "300-400": "300-400万円", "400-500": "400-500万円", "500-600": "500-600万円",
    "600-700": "600-700万円", "700-800": "700-800万円", "800-900": "800-900万円",
    "900-1000": "900-1000万円", "1000-1100": "1000-1100万円", "1100-1200": "1100-1200万円",
    "1200-1300": "1200-1300万円", "1300-1400": "1300-1400万円", "1400-1500": "1400-1500万円",
    "1500-1600": "1500-1600万円", "1600-1700": "1600-1700万円", "1700-1800": "1700-1800万円",
    "1800-1900": "1800-1900万円", "1900-2000": "1900-2000万円", "2000-2100": "2000-2100万円",
    "2100-2200": "2100-2200万円", "2200-2300": "2200-2300万円", "2300-2400": "2300-2400万円",
    "2400-2500": "2400-2500万円", "2500-2600": "2500-2600万円", "2600-2700": "2600-2700万円",
    "2700-2800": "2700-2800万円", "2800-2900": "2800-2900万円", "2900-3000": "2900-3000万円",
    "3000-3100": "3000-3100万円", "3100-3200": "3100-3200万円", "3200-3300": "3200-3300万円",
    "3300-3400": "3300-3400万円", "3400-3500": "3400-3500万円", "3500-3600": "3500-3600万円",
    "3600-3700": "3600-3700万円", "3700-3800": "3700-3800万円", "3800-3900": "3800-3900万円",
    "3900-4000": "3900-4000万円", "4000-4100": "4000-4100万円", "4100-4200": "4100-4200万円",
    "4200-4300": "4200-4300万円", "4300-4400": "4300-4400万円", "4400-4500": "4400-4500万円",
    "4500-4600": "4500-4600万円", "4600-4700": "4600-4700万円", "4700-4800": "4700-4800万円",
}
DEPARTMENT_MAP = {
    "internal_medicine": "内科", "surgery": "外科", "pediatrics": "小児科",
    "orthopedics": "整形外科", "dermatology": "皮膚科", "ophthalmology": "眼科",
    "cardiology": "循環器内科", "psychiatry": "精神科", "dentistry": "歯科",
    "pediatric_dentistry": "小児歯科", "otorhinolaryngology": "耳鼻咽喉科",
    "ent": "耳鼻咽喉科",
    "gynecology": "婦人科",
    "urology": "泌尿器科",
    "neurosurgery": "脳神経外科",
    "general_dentistry": "一般歯科",
    "orthodontics": "矯正歯科",
    "cosmetic_dentistry": "審美歯科",
    "oral_surgery": "口腔外科",
    "anesthesiology": "麻酔科",
    "radiology": "放射線科",
    "rehabilitation": "リハビリテーション科",
    "allergy": "アレルギー科",
    "gastroenterology": "消化器内科",
    "respiratory_medicine": "呼吸器内科",
    "diabetes_medicine": "糖尿病内科",
    "nephrology": "腎臓内科",
    "neurology": "神経内科",
    "hematology": "血液内科",
    "plastic_surgery": "形成外科",
    "beauty_surgery": "美容外科",
}
PURPOSE_MAP = {
    "increase_patients": "患者数を増やす",
    "increase_frequency": "来院頻度を増やす",
    "increase_spend": "客単価を増やす",
}
HEADER_MAP = {
    "personality": "性格（価値観・人生観）",
    "reason": "通院理由",
    "behavior": "症状通院頻度・行動パターン",
    "reviews": "口コミの重視ポイント",
    "values": "医療機関への価値観・行動傾向",
    "demands": "医療機関に求めるもの"
}

def format_age_for_pdf_ppt(age_value):
    if not age_value: return '-'
    try:
        age_value_str = str(age_value) # Ensure it's a string
        if 'm' in age_value_str and 'y' in age_value_str:
            parts = age_value_str.split('y')
            years = parts[0]
            months = parts[1].replace('m', '')
            if years == '0' and months == '0':
                return "0歳"
            return f"{years}歳{months}ヶ月"
        elif 'y' in age_value_str:
            return f"{age_value_str.replace('y', '')}歳"
        elif 'm' in age_value_str:
            months = age_value_str.replace('m', '')
            if months == '0':
                return "0歳"
            return f"0歳{months}ヶ月"
        elif age_value_str.isdigit():
            return f"{age_value_str}歳"
        return age_value_str
    except Exception:
        return str(age_value) 

def sanitize_for_ppt(text):
    """PPT用にテキストをサニタイズ"""
    if not text:
        return ''
    # 文字列に変換
    text = str(text)
    # 制御文字を除去（タブと改行は保持）
    text = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', text)
    # XMLで問題になる文字をエスケープ
    text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    return text

# Helper function to generate PPT
def add_text_to_shape(shape, text, font_size=Pt(9), is_bold=False, alignment=PP_ALIGN.LEFT, font_name='Meiryo UI', fill_color=None, add_border=False):
    # テキストフレームの処理を先に行う
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    
    # 背景色を設定（fillColorが指定されている場合）
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    
    # 枠線を追加（必要な場合）
    if add_border:
        shape.line.color.rgb = RGBColor(200, 200, 200)  # 薄いグレーの枠線
        shape.line.width = Pt(0.5)
    # テキストが長い場合は、テキストサイズを自動調整するのではなく、形状を固定
    if len(text) > 100:
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    else:
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    # 余白を設定
    text_frame.margin_left = Cm(0.1)
    text_frame.margin_right = Cm(0.1)
    text_frame.margin_top = Cm(0.05)
    text_frame.margin_bottom = Cm(0.05)

    # 既存の段落をクリアして新しいテキストを設定
    if text_frame.paragraphs:
        p = text_frame.paragraphs[0]
        p.clear()
    else:
        p = text_frame.add_paragraph()

    run = p.add_run()
    run.text = sanitize_for_ppt(text)
    font = run.font
    font.name = font_name
    font.size = font_size
    font.bold = is_bold
    font.color.rgb = RGBColor(0, 0, 0) # Black text
    p.alignment = alignment

    # Estimate height based on text length and font size
    num_lines = len(text.split('\\n'))
    estimated_height = num_lines * font_size.pt * 1.5 # Approximate line height factor
    return Cm(estimated_height / 28.3465 / 2.54) # Convert points to cm (rough estimate)

def generate_timeline_graph(timeline_data, output_path):
    """タイムライン分析用の散布図を生成"""
    if not GRAPH_ENABLED:
        print("[WARNING] Graph generation is disabled (matplotlib not installed)")
        return False
        
    try:
        # シンプルに英語で表示（日本語フォントの問題を回避）
        import os
        use_japanese = False
        
        # ローカル環境でのみ日本語を試す
        if 'RENDER' not in os.environ:
            try:
                import japanize_matplotlib
                plt.rcParams['font.family'] = 'IPAexGothic'
                use_japanese = True
                print("[INFO] Using Japanese fonts in local environment")
            except:
                print("[INFO] Using English labels (japanize-matplotlib not available)")
        else:
            print("[INFO] Running on Render, using English labels for compatibility")
        
        # グラフのサイズとスタイル設定
        plt.figure(figsize=(12, 6))
        plt.style.use('default')
        
        # データ取得（フロントエンドからの構造に合わせる）
        keywords = timeline_data.get('keywords', [])
        
        # 診断前後でキーワードを分割
        pre_keywords = [kw for kw in keywords if kw.get('time_diff_days', 0) < 0]
        post_keywords = [kw for kw in keywords if kw.get('time_diff_days', 0) >= 0]
        
        # 散布図データの準備
        pre_x = [kw.get('time_diff_days', 0) for kw in pre_keywords]
        pre_y = [kw.get('estimated_volume', kw.get('search_volume', 0)) for kw in pre_keywords]
        
        post_x = [kw.get('time_diff_days', 0) for kw in post_keywords]
        post_y = [kw.get('estimated_volume', kw.get('search_volume', 0)) for kw in post_keywords]
        
        # 散布図を描画（英語ラベルで統一）
        plt.scatter(pre_x, pre_y, c='#3b82f6', alpha=0.6, s=40, label='Pre-diagnosis')
        plt.scatter(post_x, post_y, c='#ef4444', alpha=0.6, s=40, label='Post-diagnosis')
        
        # 番号表示は削除（主要検索キーワードリストで確認）
        
        # 診断日に縦線を追加
        plt.axvline(x=0, color='gray', linestyle='--', alpha=0.5, label='Diagnosis Date')
        
        # グラフの装飾（英語で統一）
        plt.xlabel('Days from Diagnosis', fontsize=12)
        plt.ylabel('Search Volume', fontsize=12)
        plt.title('Timeline Analysis of Search Keywords', fontsize=14, fontweight='bold')
        plt.legend(loc='upper right')
        plt.grid(True, alpha=0.3)
        
        # Y軸を対数スケールに設定（検索ボリュームの範囲が広い場合）
        if pre_y + post_y:  # データが存在する場合
            max_volume = max(pre_y + post_y)
            min_volume = min([v for v in pre_y + post_y if v > 0] or [1])
            if max_volume / min_volume > 100:  # 100倍以上の差がある場合
                plt.yscale('log')
                plt.ylabel('Search Volume (Log Scale)', fontsize=12)
        
        plt.tight_layout()
        plt.savefig(output_path, dpi=150, bbox_inches='tight')
        plt.close()
        
        return True
    except Exception as e:
        print(f"Error generating graph: {e}")
        import traceback
        traceback.print_exc()
        return False

def generate_pdf(data):
    # A4サイズ横長に設定、レイアウト最適化
    pdf = FPDF(orientation='L', unit='mm', format='A4')
    
    # 余白を小さく設定 (左右上下の余白を8mmに)
    pdf.set_margins(8, 8, 8)
    pdf.set_auto_page_break(True, margin=8)
    
    pdf.add_page()
    
    # フォント設定
    try:
        # フォントファイルのパスを取得
        font_path = find_font_file()
        print(f"PDF generation using font from: {font_path}")
        
        # Register Regular font
        pdf.add_font("ipa", "", font_path, uni=True)
        # Register Bold style using the same file
        pdf.add_font("ipa", "B", font_path, uni=True)
        pdf.set_font("ipa", size=10) # 全体的に小さいフォントサイズをデフォルトに
    except RuntimeError as e:
        print(f"WARNING: Could not load/register font from '{font_path}'. Error: {e}. Using default font.")
        pdf.set_font("Arial", size=10) # フォールバックフォントも小さく

    profile = data.get('profile', {})
    details = data.get('details', {})
    image_url = data.get('image_url')

    # --- 定数定義 ---
    name_line_height = 6  # ペルソナ名の行の高さ (mm)
    section_title_height = 7 # セクションタイトルのセルの高さ (mm)

    # --- ページとカラムの基本設定 ---
    left_column_content_x = pdf.l_margin
    page_width = pdf.w - pdf.l_margin - pdf.r_margin
    left_column_width = page_width * 0.35
    right_column_width = page_width * 0.65
    column_gap = 5
    right_column_x = left_column_content_x + left_column_width + column_gap
    
    content_start_y = pdf.get_y() # 全体のコンテンツ開始Y座標（一番上の要素に合わせる）
    header_end_y = content_start_y # header_end_y を定義 (アイコンなどの開始Y座標の基準)

    # --- 左カラムの描画開始 ---
    pdf.set_xy(left_column_content_x, header_end_y + 5) # ヘッダーの下に少しスペース

    # --- ペルソナアイコンと名前 ---
    icon_size = 30 # アイコンのサイズ (mm)
    icon_y_position = pdf.get_y()

    if image_url:
        try:
            print(f"Fetching image from URL: {image_url}")
            # Data URLの場合の処理
            if image_url.startswith('data:'):
                # data:image/png;base64,... の形式から画像データを抽出
                try:
                    header, encoded = image_url.split(',', 1)
                    image_data = base64.b64decode(encoded)
                except (ValueError, base64.binascii.Error) as e:
                    print(f"Error decoding data URL: {e}")
                    raise
            else:
                # 通常のURLの場合
                image_data = urlopen(image_url).read()
            
            # PIL Imageで画像を開く
            img_file_obj = io.BytesIO(image_data)
            pil_image = Image.open(img_file_obj)
            
            # アスペクト比を保ちながら正方形にリサイズ（中央クロップ）
            original_width, original_height = pil_image.size
            
            # 正方形にクロップ
            min_dimension = min(original_width, original_height)
            left = (original_width - min_dimension) // 2
            top = (original_height - min_dimension) // 2
            right = left + min_dimension
            bottom = top + min_dimension
            
            # クロップして正方形にする
            cropped_image = pil_image.crop((left, top, right, bottom))
            
            # リサイズして最終サイズに調整
            final_size = int(icon_size * 10)  # mmをpixelに変換（おおよそ）
            resized_image = cropped_image.resize((final_size, final_size), Image.Resampling.LANCZOS)
            
            # RGB形式に変換（PDFで確実に表示されるように）
            if resized_image.mode != 'RGB':
                resized_image = resized_image.convert('RGB')
            
            # バイトストリームに保存
            processed_img_obj = io.BytesIO()
            resized_image.save(processed_img_obj, format='JPEG', quality=85)
            processed_img_obj.seek(0)
            
            # PDFに画像を追加
            pdf.image(processed_img_obj, x=left_column_content_x, y=icon_y_position, w=icon_size, h=icon_size)
            
        except Exception as e:
            print(f"Error loading image: {e}")
            # アイコン失敗時は代替テキストや枠線などを表示
            pdf.rect(left_column_content_x, icon_y_position, icon_size, icon_size, style='D')
            pdf.set_xy(left_column_content_x + 1, icon_y_position + icon_size / 2 - 2)
            pdf.multi_cell(icon_size - 2, 4, "No Img", 0, 'C')

    # 名前の描画開始位置をアイコンの右隣に設定
    name_x_position = left_column_content_x + icon_size + 3 # アイコンの右に少しスペース
    pdf.set_xy(name_x_position, icon_y_position + (icon_size / 2) - (name_line_height / 2) - 2) # 上下中央揃えっぽく調整
    
    pdf.set_font("ipa", 'B', 14) # 名前は少し大きく太字に
    # 名前の最大幅を左カラムの残り幅に制限
    name_max_width = left_column_width - (icon_size + 3) # アイコンとスペース分を引く
    pdf.multi_cell(name_max_width, name_line_height, profile.get('name', '-'), 0, 'L')
    
    # アイコンと名前の下に基本情報を配置するためのY座標を設定
    # アイコンの高さ、または名前の高さのうち、大きい方を基準にする
    current_y_after_icon_name = icon_y_position + icon_size + 5 # アイコンの下端 + 余白

    # --- 基本情報セクション ---
    pdf.set_xy(left_column_content_x, current_y_after_icon_name)
    pdf.set_font("ipa", 'B', 11) # セクションタイトル (太字、下線削除)
    pdf.set_fill_color(200, 220, 240)  # 薄い青色の背景
    pdf.set_text_color(0, 0, 0)  # テキストは黒
    pdf.cell(left_column_width, 7, "基本情報", 0, 1, 'L', fill=True)
    current_y_after_icon_name = pdf.get_y() # 「基本情報」タイトルの後にY座標を更新
    pdf.set_font("ipa", '', 10) # 内容のフォントに戻す

    # 1. 診療科 (左カラム)
    pdf.set_xy(left_column_content_x, current_y_after_icon_name) # 更新されたY座標を使用
    pdf.set_font("ipa", '', 10)
    department_val = profile.get('department', '-')
    department_display = DEPARTMENT_MAP.get(department_val.lower(), department_val) # Lowercase for map key
    pdf.multi_cell(left_column_width, 7, f"診療科: {department_display}", 0, 'L')
    current_y_after_icon_name = pdf.get_y()

    # 2. 主訴 (左カラム)
    pdf.set_xy(left_column_content_x, current_y_after_icon_name)
    chief_complaint_val = profile.get('chief_complaint', '-')
    pdf.multi_cell(left_column_width, 7, f"主訴: {chief_complaint_val}", 0, 'L')
    current_y_after_icon_name = pdf.get_y()

    # 3. 作成目的 (左カラム)
    pdf.set_xy(left_column_content_x, current_y_after_icon_name)
    purpose_val = profile.get('purpose', '-')
    purpose_display = PURPOSE_MAP.get(purpose_val.lower(), purpose_val) # Lowercase for map key
    pdf.multi_cell(left_column_width, 7, f"作成目的: {purpose_display}", 0, 'L')
    current_y_after_icon_name = pdf.get_y()
    pdf.ln(3) # 少しスペース
    current_y_after_icon_name = pdf.get_y()

    # 3. 基本情報セクション (左カラム)
    pdf.set_xy(left_column_content_x, current_y_after_icon_name) 
    pdf.set_font("ipa", '', 9)
    info_items = [
        ("性別", GENDER_MAP.get(profile.get('gender', '-'), profile.get('gender', '-'))),
        ("年齢", format_age_for_pdf_ppt(profile.get('age', '-'))),
        ("都道府県", profile.get('prefecture', '-')),
        ("市区町村", profile.get('municipality', '-')),
        ("職業", profile.get('occupation', '-')),
        ("年収", INCOME_MAP.get(profile.get('income', '-'), profile.get('income', '-'))),
        ("家族構成", profile.get('family', '-')),
        ("趣味", profile.get('hobby', '-')),
        ("ライフイベント", profile.get('life_events', '-')),
        ("患者タイプ", profile.get('patient_type', '-'))
    ]
    
    item_height = 4 
    key_width = 25 
    value_width = left_column_width - key_width

    for i, (key, value_display) in enumerate(info_items):
        pdf.set_xy(left_column_content_x, current_y_after_icon_name)
        # 偶数行に薄い背景色を追加（ストライプ効果）
        if i % 2 == 0:
            pdf.set_fill_color(245, 245, 245)  # 非常に薄いグレー
            pdf.rect(left_column_content_x, current_y_after_icon_name, left_column_width, item_height + 1, 'F')
        pdf.set_font("ipa", '', 9)
        pdf.set_text_color(0, 0, 0)  # テキストは黒
        pdf.set_xy(left_column_content_x, current_y_after_icon_name)
        pdf.cell(key_width, item_height, f"{key}:", 0, 0, 'L')
        pdf.set_font("ipa", '', 9)
        pdf.set_xy(left_column_content_x + key_width, current_y_after_icon_name)
        pdf.multi_cell(value_width, item_height, str(value_display), 0, 'L')
        current_y_after_icon_name = pdf.get_y() 

    current_y_after_icon_name += 3 
    
    # 4. その他の特徴セクション (左カラム)
    pdf.set_xy(left_column_content_x, current_y_after_icon_name)
    pdf.set_font("ipa", '', 11)
    pdf.set_fill_color(200, 220, 240)  # 薄い青色の背景
    pdf.set_text_color(0, 0, 0)  # テキストは黒
    pdf.cell(left_column_width, 7, "その他の特徴", 0, 1, 'L', fill=True)  # 高さを6から7に変更
    current_y_after_icon_name = pdf.get_y() 
    pdf.ln(1)
    current_y_after_icon_name = pdf.get_y() 
    
    additional_items = [
        ("座右の銘", profile.get('motto', '-')),
        ("最近の悩み/関心", profile.get('concerns', '-')),
        ("好きな有名人", profile.get('favorite_person', '-')),
        ("よく見るメディア", profile.get('media_sns', '-')),
        ("性格キーワード", profile.get('personality_keywords', '-')),
        ("健康に関する行動", profile.get('health_actions', '-')),
        ("休日の過ごし方", profile.get('holiday_activities', '-')),
        ("キャッチコピー", profile.get('catchphrase', '-'))
    ]
    
    additional_key_width = 30
    additional_value_width = left_column_width - additional_key_width

    for i, (key, value) in enumerate(additional_items):
        pdf.set_xy(left_column_content_x, current_y_after_icon_name)
        # 偶数行に薄い背景色を追加（ストライプ効果）
        if i % 2 == 0:
            pdf.set_fill_color(245, 245, 245)  # 非常に薄いグレー
            pdf.rect(left_column_content_x, current_y_after_icon_name, left_column_width, item_height + 1, 'F')
        pdf.set_font("ipa", '', 9)
        pdf.set_text_color(0, 0, 0)  # テキストは黒
        pdf.set_xy(left_column_content_x, current_y_after_icon_name)
        pdf.cell(additional_key_width, item_height, f"{key}:", 0, 0, 'L')
        pdf.set_font("ipa", '', 9)
        pdf.set_xy(left_column_content_x + additional_key_width, current_y_after_icon_name)
        pdf.multi_cell(additional_value_width, item_height, str(value), 0, 'L')
        pdf.ln(1)  # 項目間に1mmの余白を追加
        current_y_after_icon_name = pdf.get_y()

    if profile.get('additional_field_name') and profile.get('additional_field_value'):
        additional_fields = zip(profile.get('additional_field_name'), profile.get('additional_field_value'))
        current_y_after_icon_name +=1 
        for j, (field_name, field_value) in enumerate(additional_fields):
            if field_name or field_value:
                pdf.set_xy(left_column_content_x, current_y_after_icon_name)
                # 追加フィールドも含めた総数で偶数行判定
                total_index = len(additional_items) + j
                if total_index % 2 == 0:
                    pdf.set_fill_color(245, 245, 245)  # 非常に薄いグレー
                    pdf.rect(left_column_content_x, current_y_after_icon_name, left_column_width, item_height + 1, 'F')
                pdf.set_font("ipa", '', 9)
                pdf.set_text_color(0, 0, 0)  # テキストは黒
                pdf.set_xy(left_column_content_x, current_y_after_icon_name)
                pdf.cell(additional_key_width, item_height, f"{field_name if field_name else ''}:", 0, 0, 'L')
                pdf.set_font("ipa", '', 9)
                pdf.set_xy(left_column_content_x + additional_key_width, current_y_after_icon_name)
                pdf.multi_cell(additional_value_width, item_height, str(field_value) if field_value else '-', 0, 'L')
                pdf.ln(1)  # 項目間に1mmの余白を追加
                current_y_after_icon_name = pdf.get_y()
    
    max_left_y = current_y_after_icon_name # 左カラムの最終Y座標

    # --- 右カラムの描画開始 ---
    # 右カラムの開始Y座標を、左カラムのアイコンの開始高さに合わせる
    right_column_current_y = icon_y_position 
    
    for detail_key, japanese_header_text in HEADER_MAP.items():
        value = details.get(detail_key)
        if value and str(value).strip(): # 値が存在し、空でない場合のみ描画
            pdf.set_xy(right_column_x, right_column_current_y)
            
            # セクションヘッダー
            pdf.set_font("ipa", 'B', 12) # 太字、サイズ12に変更（下線削除）
            pdf.set_fill_color(200, 230, 200) # 薄い緑色の背景に変更
            pdf.set_text_color(0, 0, 0)  # テキストは黒
            pdf.cell(right_column_width, 8, str(japanese_header_text), 0, 1, 'L', fill=True) # 高さを6から8に変更
            # pdf.get_y() はこのセルの後に自動更新される

            # コンテンツ
            pdf.set_x(right_column_x) # multi_cell のためにX座標を右カラムの開始位置にリセット
            pdf.set_font("ipa", '', 10) # 通常フォント、サイズ10に変更
            pdf.set_text_color(0, 0, 0)  # テキストは黒にリセット
            pdf.ln(2) # コンテンツ前に少し余白を追加
            pdf.set_x(right_column_x)
            pdf.multi_cell(right_column_width, 6, str(value), 0, 'L') # 行の高さ5mmから6mmに変更
            right_column_current_y = pdf.get_y() # multi_cell 後のY座標を更新
            
            pdf.ln(5) # セクション間のスペースを3mmから5mmに増加
            right_column_current_y = pdf.get_y() # スペース後のY座標を更新

    # --- タイムライン分析セクション（新しいページに追加） ---
    timeline_analysis = data.get('timeline_analysis')
    if timeline_analysis and timeline_analysis.get('ai_analysis'):
        pdf.add_page()
        
        # セクションタイトル
        pdf.set_font("ipa", "B", 14)
        pdf.cell(pdf.w - pdf.l_margin - pdf.r_margin, 10, 'タイムライン分析', 0, 1, 'C')
        pdf.ln(5)
        
        # グラフを生成して追加
        try:
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                graph_path = tmp_file.name
                
            if generate_timeline_graph(timeline_analysis, graph_path):
                # グラフをPDFに追加（ページ幅の80%を使用）
                page_width = pdf.w - pdf.l_margin - pdf.r_margin
                graph_width = page_width * 0.8
                graph_x = pdf.l_margin + (page_width - graph_width) / 2
                
                # グラフの高さを計算（幅の約半分）
                graph_height = graph_width * 0.5
                current_y = pdf.get_y()
                
                pdf.image(graph_path, x=graph_x, y=current_y, w=graph_width, h=graph_height)
                
                # グラフの後に適切なスペースを追加
                pdf.set_y(current_y + graph_height + 10)  # グラフの高さ + 10mmの余白
                
                # 一時ファイルを削除
                import os
                os.unlink(graph_path)
        except Exception as e:
            print(f"Error adding graph to PDF: {e}")
        
        # スペースが不足している場合は新しいページを追加
        current_y_after_graph = pdf.get_y()
        remaining_space = pdf.h - pdf.b_margin - current_y_after_graph
        if remaining_space < 50:  # 50mm以下の場合
            pdf.add_page()
        
        
        # 主要検索キーワード（上位10件）を表示（グラフの番号と対応）
        if timeline_analysis.get('keywords'):
            pdf.set_font("ipa", "B", 11)
            pdf.cell(pdf.w - pdf.l_margin - pdf.r_margin, 7, '主要検索キーワード（上位10件）', 0, 1)
            pdf.set_font("ipa", "", 9)
            
            # 検索ボリューム順にソートして上位10件を取得
            all_keywords = timeline_analysis['keywords']
            sorted_keywords = sorted(all_keywords, 
                                   key=lambda k: k.get('estimated_volume', k.get('search_volume', 0)), 
                                   reverse=True)
            keywords = sorted_keywords[:10]
            for i, kw in enumerate(keywords, 1):
                keyword_text = f"{i}. {kw['keyword']} ({kw['time_diff_days']:.1f}日"
                if kw['time_diff_days'] < 0:
                    keyword_text += "前)"
                else:
                    keyword_text += "後)"
                keyword_text += f" - 検索ボリューム: {kw.get('estimated_volume', kw.get('search_volume', 0)):,}"
                pdf.cell(pdf.w - pdf.l_margin - pdf.r_margin, 5, keyword_text, 0, 1)
            pdf.ln(8)  # キーワードリストの後にスペース
        
        # AI分析レポート
        pdf.set_font("ipa", "B", 11)
        pdf.cell(pdf.w - pdf.l_margin - pdf.r_margin, 7, 'AI分析レポート', 0, 1)
        pdf.ln(2)  # タイトルの後に少しスペース
        pdf.set_font("ipa", "", 9)
        
        ai_analysis_text = timeline_analysis.get('ai_analysis', '')
        
        if ai_analysis_text:
            # 現在のX座標を保存
            left_x = pdf.l_margin
            
            # テキストを段落ごとに処理
            paragraphs = ai_analysis_text.split('\n\n')
            for i, paragraph in enumerate(paragraphs):
                if paragraph.strip():
                    # X座標を左マージンに設定（左寄せを確実にする）
                    pdf.set_x(left_x)
                    # 段落全体を一つのmulti_cellで表示
                    pdf.multi_cell(pdf.w - pdf.l_margin - pdf.r_margin, 5, paragraph.strip(), 0, 'L')
                    
                    # 段落間のスペース（最後の段落以外）
                    if i < len(paragraphs) - 1:
                        pdf.ln(3)
        else:
            pdf.multi_cell(pdf.w - pdf.l_margin - pdf.r_margin, 5, "AI分析データがありません", 0, 'L')
    
    # Generate PDF in memory
    pdf_output = pdf.output() # Get output as bytes directly
    buffer = io.BytesIO(pdf_output)
    buffer.seek(0)
    return buffer

def generate_ppt(persona_data, image_path=None, department_text=None, purpose_text=None):
    import os  # osモジュールをインポート
    
    prs = Presentation()
    prs.slide_width = Inches(11.69)  # A4 Landscape width
    prs.slide_height = Inches(8.27) # A4 Landscape height
    # Try to use the most blank layout available
    try:
        slide_layout = prs.slide_layouts[6]  # Sometimes layout 6 is even more blank
    except:
        slide_layout = prs.slide_layouts[5]  # Fall back to layout 5 (blank layout)
    slide = prs.slides.add_slide(slide_layout)
    
    # スライド上のすべてのプレースホルダーを削除
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            shapes_to_remove.append(shape)
    
    # プレースホルダーを逆順で削除（インデックスの問題を避けるため）
    for shape in reversed(shapes_to_remove):
        sp = shape.element
        sp.getparent().remove(sp)
    
    # 余分なテキストフレームをさらに削除
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            shape.text_frame.clear()

    # Margins (approximated from PDF's 8mm)
    left_margin_ppt = Cm(0.8)
    right_margin_ppt = Cm(0.8)
    top_margin_ppt = Cm(0.8)
    bottom_margin_ppt = Cm(0.8)
    content_width = prs.slide_width - left_margin_ppt - right_margin_ppt
    content_height = prs.slide_height - top_margin_ppt - bottom_margin_ppt
    
    item_spacing_ppt = Cm(0.15) # Spacing between items

    # Title
    title_shape = slide.shapes.add_textbox(left=Cm(0.5), top=Cm(0.2), width=prs.slide_width - Cm(1.0), height=Cm(1.0))
    text_frame = title_shape.text_frame
    text_frame.clear()  # 既存のテキストをクリア
    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    p.text = "生成されたペルソナ"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.name = 'Meiryo UI'

    # Icon
    icon_left = left_margin_ppt
    icon_top = top_margin_ppt + Cm(1.0) # Below title
    icon_size = Cm(3.0)
    
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, icon_left, icon_top, height=icon_size)
        except Exception as e:
            print(f"Error adding image to PPT: {e}")
            # Add a placeholder if image fails
            icon_placeholder = slide.shapes.add_textbox(icon_left, icon_top, icon_size, icon_size)
            add_text_to_shape(icon_placeholder, "画像エラー", font_size=Pt(8))
    else:
        icon_placeholder = slide.shapes.add_textbox(icon_left, icon_top, icon_size, icon_size)
        add_text_to_shape(icon_placeholder, "画像なし", font_size=Pt(8))

    # Persona Name (right of icon)
    name_left = icon_left + icon_size + Cm(0.3)
    name_top = icon_top 
    name_width = content_width * 0.35 - icon_size - Cm(0.3) # Available width in the 35% column part
    name_height = icon_size # Align height with icon
    
    name_text_box = slide.shapes.add_textbox(name_left, name_top, name_width, name_height)
    # Vertical centering
    tf_name = name_text_box.text_frame
    tf_name.word_wrap = True
    tf_name.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE 
    p_name = tf_name.paragraphs[0] if tf_name.paragraphs else tf_name.add_paragraph()
    p_name.clear()
    run_name = p_name.add_run()
    run_name.text = persona_data.get('name', '')
    font_name_style = run_name.font
    font_name_style.name = 'Meiryo UI'
    font_name_style.size = Pt(12)
    font_name_style.bold = True
    p_name.alignment = PP_ALIGN.LEFT

    # Department and Purpose below icon and name
    if not department_text:
        department_text = DEPARTMENT_MAP.get(persona_data.get('department', '').lower(), persona_data.get('department', '-'))
    if not purpose_text:
        purpose_text = PURPOSE_MAP.get(persona_data.get('purpose', '').lower(), persona_data.get('purpose', '-'))

    header_info_top = icon_top + icon_size + Cm(0.2)
    header_info_shape = slide.shapes.add_textbox(icon_left, header_info_top, content_width * 0.33, Cm(1.0))
    text_frame = header_info_shape.text_frame
    text_frame.clear()
    
    header_items = [
        ("診療科", department_text),
        ("主訴", persona_data.get('chief_complaint', '-')),
        ("作成目的", purpose_text)
    ]
    for i, (label, value) in enumerate(header_items):
        if i == 0:
            # 最初の段落を使用
            item_p = text_frame.paragraphs[0]
        else:
            # 2番目以降は新しい段落を追加
            item_p = text_frame.add_paragraph()
        
        run_label = item_p.add_run()
        run_label.text = f"{label}: "
        font_label = run_label.font
        font_label.name = 'Meiryo UI'
        font_label.size = Pt(9)
        font_label.bold = True

        run_value = item_p.add_run()
        run_value.text = sanitize_for_ppt(value)
        font_value = run_value.font
        font_value.name = 'Meiryo UI'
        font_value.size = Pt(9)
        item_p.alignment = PP_ALIGN.LEFT

    y_position_left = header_info_top + Cm(1.0) + Cm(0.3) # Start Y for left column content
    y_position_right = icon_top # Start Y for right column content (aligned with icon top)

    # Left Column (Basic Information)
    left_column_x = left_margin_ppt
    left_width = content_width * 0.33
    
    current_y_left = y_position_left

    # Section Title: 基本情報
    shape_title_basic = slide.shapes.add_textbox(left_column_x, current_y_left, left_width, Cm(0.6))
    add_text_to_shape(shape_title_basic, "基本情報", font_size=Pt(11), is_bold=True, font_name='Meiryo UI')
    current_y_left += Cm(0.6) + item_spacing_ppt # Update Y position

    basic_info_keys = ["gender", "age", "prefecture", "municipality", "family", "occupation", "income", "hobby", "life_events", "patient_type"]
    basic_info_labels = {
        "gender": "性別", "age": "年齢", "prefecture": "都道府県", "municipality": "市区町村",
        "family": "家族構成", "occupation": "職業", "income": "年収", "hobby": "趣味",
        "life_events": "ライフイベント", "patient_type": "患者タイプ"
    }

    for idx, key in enumerate(basic_info_keys):
        value = persona_data.get(key, "-")
        if key == "gender": value = GENDER_MAP.get(value, value)
        elif key == "age": value = format_age_for_pdf_ppt(value)
        elif key == "income": value = INCOME_MAP.get(value, value)
        
        item_text = f"{basic_info_labels.get(key, key)}: {value}"
        item_shape = slide.shapes.add_textbox(left_column_x, current_y_left, left_width, Cm(0.5))
        # 偶数行に薄いグレーの背景色を設定
        fill_color = RGBColor(245, 245, 245) if idx % 2 == 0 else None
        add_text_to_shape(item_shape, item_text, font_size=Pt(9), font_name='Meiryo UI', fill_color=fill_color)
        current_y_left += Cm(0.5) + item_spacing_ppt

    # Additional Fixed Fields (Left Column)
    additional_fixed_keys_labels = {
        "motto": "座右の銘", "concerns": "最近の悩み/関心", "favorite_person": "好きな有名人/尊敬する人物",
        "media_sns": "よく見るメディア/SNS", "personality_keywords": "性格キーワード",
        "health_actions": "最近した健康に関する行動", "holiday_activities": "休日の過ごし方",
        "catchphrase": "キャッチコピー"
    }
    # Section Title: 追加情報 (Fixed)
    current_y_left += item_spacing_ppt # Extra space before next section
    shape_title_add_fixed = slide.shapes.add_textbox(left_column_x, current_y_left, left_width, Cm(0.6))
    add_text_to_shape(shape_title_add_fixed, "追加情報", font_size=Pt(11), is_bold=True, font_name='Meiryo UI')
    current_y_left += Cm(0.6) + item_spacing_ppt

    for idx, (key, label) in enumerate(additional_fixed_keys_labels.items()):
        value = persona_data.get(key, persona_data.get(key.replace('_input', ''), "-"))
        item_text = f"{label}: {value}"
        item_shape = slide.shapes.add_textbox(left_column_x, current_y_left, left_width, Cm(0.5))
        # 偶数行に薄いグレーの背景色を設定
        fill_color = RGBColor(245, 245, 245) if idx % 2 == 0 else None
        add_text_to_shape(item_shape, item_text, font_size=Pt(9), font_name='Meiryo UI', fill_color=fill_color)
        current_y_left += Cm(0.5) + item_spacing_ppt
        
    # Dynamic Additional Fields (Left Column)
    if persona_data.get("additional_field_name") and persona_data.get("additional_field_value"):
        fields = list(zip(persona_data.get("additional_field_name"), persona_data.get("additional_field_value")))
        has_fields = False
        field_idx = 0
        for field_name, field_value in fields:
            if field_name and field_value:
                if not has_fields:
                    # Section title on first valid field
                    current_y_left += item_spacing_ppt
                    shape_title_add_dyn = slide.shapes.add_textbox(left_column_x, current_y_left, left_width, Cm(0.6))
                    add_text_to_shape(shape_title_add_dyn, "自由記述項目", font_size=Pt(11), is_bold=True, font_name='Meiryo UI')
                    current_y_left += Cm(0.6) + item_spacing_ppt
                    has_fields = True
                
                item_text = f"{field_name}: {field_value}"
                item_shape = slide.shapes.add_textbox(left_column_x, current_y_left, left_width, Cm(0.5))
                # 偶数行に薄いグレーの背景色を設定
                fill_color = RGBColor(245, 245, 245) if field_idx % 2 == 0 else None
                add_text_to_shape(item_shape, item_text, font_size=Pt(9), font_name='Meiryo UI', fill_color=fill_color)
                current_y_left += Cm(0.5) + item_spacing_ppt
                field_idx += 1

    # Right Column (Detailed Information)
    right_column_x = left_margin_ppt + content_width * 0.35 + Cm(0.3) # Start after left column + gap
    right_width = content_width * 0.65 - Cm(0.3) # Remaining width
    
    current_y_right = y_position_right

    detail_key_map = {
        "personality": "性格（価値観・人生観）",
        "reason": "通院理由",
        "behavior": "症状通院頻度・行動パターン",
        "reviews": "口コミの重視ポイント",
        "values": "医療機関への価値観・行動傾向",
        "demands": "医療機関に求めるもの"
    }
    
    # Draw "性格（価値観・人生観）" first as it's a primary section
    if persona_data.get('personality'):
        section_title = detail_key_map['personality']
        value = persona_data['personality']
        
        title_shape = slide.shapes.add_textbox(right_column_x, current_y_right, right_width, Cm(0.8))
        # セクションヘッダーに薄い緑色の背景を設定
        add_text_to_shape(title_shape, section_title, font_size=Pt(12), is_bold=True, font_name='Meiryo UI', fill_color=RGBColor(200, 230, 200))
        current_y_right += Cm(0.8)

        content_shape = slide.shapes.add_textbox(right_column_x, current_y_right, right_width, Cm(2.5))
        add_text_to_shape(content_shape, value, font_size=Pt(10.5), font_name='Meiryo UI')
        current_y_right += Cm(2.5) + item_spacing_ppt * 2
        
    # Other detailed sections
    for key in ["reason", "behavior", "reviews", "values", "demands"]:
        if key in persona_data and persona_data[key]:
            section_title = detail_key_map.get(key, key)
            value = persona_data[key]
            
            title_shape = slide.shapes.add_textbox(right_column_x, current_y_right, right_width, Cm(0.8))
            # セクションヘッダーに薄い緑色の背景を設定
            add_text_to_shape(title_shape, section_title, font_size=Pt(12), is_bold=True, font_name='Meiryo UI', fill_color=RGBColor(200, 230, 200))
            current_y_right += Cm(0.8)

            content_shape = slide.shapes.add_textbox(right_column_x, current_y_right, right_width, Cm(2.0))
            add_text_to_shape(content_shape, value, font_size=Pt(10.5), font_name='Meiryo UI')
            current_y_right += Cm(2.0) + item_spacing_ppt * 2
    
    # --- タイムライン分析スライド ---
    timeline_analysis = persona_data.get('timeline_analysis')
    if timeline_analysis and timeline_analysis.get('ai_analysis'):
        # 新しいスライドを追加
        slide = prs.slides.add_slide(slide_layout)
        
        # スライド上のすべてのプレースホルダーを削除
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                shapes_to_remove.append(shape)
        
        # プレースホルダーを逆順で削除（インデックスの問題を避けるため）
        for shape in reversed(shapes_to_remove):
            sp = shape.element
            sp.getparent().remove(sp)
        
        # 余分なテキストフレームをさらに削除
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                shape.text_frame.clear()
        
        # タイトル
        title_shape = slide.shapes.add_textbox(left_margin_ppt, Cm(1), prs.slide_width - left_margin_ppt * 2, Cm(1.5))
        add_text_to_shape(title_shape, 'タイムライン分析', font_size=Pt(20), is_bold=True, 
                         font_name='Meiryo UI')
        
        # グラフを生成して追加
        try:
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                graph_path = tmp_file.name
                
            if generate_timeline_graph(timeline_analysis, graph_path):
                # グラフをスライドに追加（上部に配置、サイズを調整）
                slide_width = prs.slide_width
                graph_width = Cm(18)  # 18cm幅に縮小
                graph_height = Cm(7)  # 7cm高さに縮小
                graph_x = (slide_width - graph_width) / 2
                graph_y = Cm(2.0)  # 2cmから開始
                
                slide.shapes.add_picture(graph_path, graph_x, graph_y, width=graph_width, height=graph_height)
                
                # 一時ファイルを削除
                import os
                os.unlink(graph_path)
        except Exception as e:
            print(f"Error adding graph to PPT: {e}")
        
        # 左カラム：主要キーワード（グラフの下に配置）
        # グラフが2.0cmから始まり、7cm高さなので、9.0cmで終わる
        left_column_y = Cm(9.5)  # グラフの下0.5cmの余白
        left_width = content_width * 0.35  # 左カラムの幅を定義
        
        # 主要キーワード（上位5件 - グラフの番号と対応）
        if timeline_analysis.get('keywords'):
            keywords_title = slide.shapes.add_textbox(left_margin_ppt, left_column_y, left_width, Cm(0.8))
            add_text_to_shape(keywords_title, '主要検索キーワード', font_size=Pt(12), is_bold=True, 
                             font_name='Meiryo UI', fill_color=RGBColor(200, 230, 200))
            left_column_y += Cm(1)
            
            # 検索ボリューム順にソートして上位10件を取得
            all_keywords = timeline_analysis['keywords']
            sorted_keywords = sorted(all_keywords, 
                                   key=lambda k: k.get('estimated_volume', k.get('search_volume', 0)), 
                                   reverse=True)
            keywords = sorted_keywords[:10]
            keywords_text = ""
            for i, kw in enumerate(keywords, 1):
                keyword_text = f"{i}. {kw['keyword']} ("
                if kw['time_diff_days'] < 0:
                    keyword_text += f"{abs(kw['time_diff_days']):.1f}日前)"
                else:
                    keyword_text += f"{kw['time_diff_days']:.1f}日後)"
                keywords_text += keyword_text + "\n"
            
            keywords_shape = slide.shapes.add_textbox(left_margin_ppt, left_column_y, left_width, Cm(5))
            add_text_to_shape(keywords_shape, keywords_text.strip(), font_size=Pt(9), font_name='Meiryo UI')
        
        # 右カラム：AI分析レポート（グラフの下に配置）
        right_column_y = Cm(9.5)  # 左カラムと同じ高さから開始
        analysis_title = slide.shapes.add_textbox(right_column_x, right_column_y, right_width, Cm(0.8))
        add_text_to_shape(analysis_title, 'AI分析レポート', font_size=Pt(12), is_bold=True, 
                         font_name='Meiryo UI', fill_color=RGBColor(200, 230, 200))
        
        ai_analysis_text = timeline_analysis.get('ai_analysis', '')
        if ai_analysis_text:
            analysis_shape = slide.shapes.add_textbox(right_column_x, right_column_y + Cm(1), right_width, Cm(10))
            add_text_to_shape(analysis_shape, ai_analysis_text, font_size=Pt(9), font_name='Meiryo UI')
    
    # Save to memory stream
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

# Google Maps Static API用のプロキシエンドポイント（セキュア）
@app.get("/api/google-maps-static")
async def get_google_maps_static(
    center: str,
    zoom: int = 14,
    size: str = "600x400",
    markers: str = None,
    username: str = Depends(verify_admin_credentials)
):
    """Google Maps Static APIのプロキシ（APIキーを隠蔽）"""
    api_key = os.getenv("GOOGLE_MAPS_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="Google Maps API key not configured")
    
    # Google Maps Static API URL
    base_url = "https://maps.googleapis.com/maps/api/staticmap"
    params = {
        "center": center,
        "zoom": zoom,
        "size": size,
        "key": api_key
    }
    
    if markers:
        params["markers"] = markers
    
    async with aiohttp.ClientSession() as session:
        async with session.get(base_url, params=params) as response:
            if response.status == 200:
                content = await response.read()
                return Response(content=content, media_type="image/png")
            else:
                raise HTTPException(status_code=response.status, detail="Failed to fetch map")

# 診療科リストを取得するエンドポイント
@app.get("/api/departments/{category}")
async def get_departments_list(category: str):
    """指定カテゴリの診療科リストを取得"""
    try:
        chief_complaints = load_chief_complaints_data()
        if category not in chief_complaints:
            raise HTTPException(status_code=404, detail=f"Category {category} not found")
        
        # 診療科名のリストを返す
        departments = list(chief_complaints[category].keys())
        return {"departments": departments}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# デバッグ用エンドポイント（本番環境では削除すること）
@app.get("/api/debug/google-maps-status")
async def check_google_maps_status(username: str = Depends(verify_admin_credentials)):
    """Google Maps API の状態を確認"""
    api_key_set = bool(os.getenv("GOOGLE_MAPS_API_KEY"))
    
    return {
        "api_key_configured": api_key_set,
        "api_key_first_chars": os.getenv("GOOGLE_MAPS_API_KEY", "")[:4] + "..." if api_key_set else None,
        "message": "API key is configured" if api_key_set else "API key is NOT configured"
    }

@app.post("/api/debug/test-geocoding")
async def test_geocoding(request: Request, username: str = Depends(verify_admin_credentials)):
    """Test geocoding with a specific address"""
    try:
        data = await request.json()
        address = data.get("address", "東京都千代田区丸の内1-1-1")
        
        google_maps = GoogleMapsService()
        result = await google_maps._geocode_address(address)
        
        return {
            "address": address,
            "geocoding_result": result,
            "api_key_set": bool(google_maps.api_key),
            "last_status": getattr(google_maps, '_last_geocoding_status', None)
        }
    except Exception as e:
        return {
            "error": str(e),
            "address": address,
            "api_key_set": bool(os.getenv("GOOGLE_MAPS_API_KEY"))
        }

# 競合分析API
@app.post("/api/competitive-analysis")
async def analyze_competitors(request: Request, username: str = Depends(verify_admin_credentials)):
    """競合分析を実行"""
    # レート制限チェック
    from backend.utils.rate_limiter import competitive_analysis_limiter, check_rate_limit
    check_rate_limit(competitive_analysis_limiter, request, username)
    
    try:
        data = await request.json()
        
        # 必須パラメータの確認
        clinic_info = data.get('clinic_info')
        search_radius = data.get('search_radius', 3000)
        additional_info = data.get('additional_info', '')
        
        if not clinic_info or not clinic_info.get('address'):
            return JSONResponse(
                status_code=400,
                content={"error": "clinic_info with address is required"}
            )
        
        # 入力検証
        if not isinstance(search_radius, int) or search_radius < 100 or search_radius > 50000:
            return JSONResponse(
                status_code=400,
                content={"error": "search_radius must be between 100 and 50000 meters"}
            )
        
        if 'departments' in clinic_info and not isinstance(clinic_info['departments'], list):
            return JSONResponse(
                status_code=400,
                content={"error": "departments must be a list"}
            )
        
        # 住所の基本的な検証
        address = clinic_info.get('address', '').strip()
        if len(address) < 5 or len(address) > 200:
            return JSONResponse(
                status_code=400,
                content={"error": "Invalid address length"}
            )
        
        # 競合分析サービスのインスタンスを作成
        competitive_service = CompetitiveAnalysisService()
        
        # 分析を実行（正しいメソッド名を使用: analyze_competition）
        result = await competitive_service.analyze_competition({
            "clinic_info": clinic_info,
            "search_radius": search_radius,
            "additional_info": additional_info
        })
        
        if result.get("error"):
            return JSONResponse(
                status_code=500,
                content={"error": result["error"]}
            )
        
        return JSONResponse(content=result)
        
    except Exception as e:
        print(f"Error in competitive analysis: {e}")
        traceback.print_exc()
        return JSONResponse(
            status_code=500,
            content={"error": f"競合分析中にエラーが発生しました: {str(e)}"}
        )

# Google Maps APIキー取得エンドポイント（認証必須）
@app.get("/api/google-maps-key")
async def get_google_maps_api_key(username: str = Depends(verify_admin_credentials)):
    """Google Maps APIキーを取得（認証ユーザーのみ）"""
    api_key = os.getenv("GOOGLE_MAPS_API_KEY")
    if not api_key:
        return JSONResponse(
            status_code=500,
            content={"error": "Google Maps API key not configured"}
        )
    return {"api_key": api_key}

@app.post("/api/competitive-analysis/search-clinics")
async def search_nearby_clinics(request: Request, username: str = Depends(verify_admin_credentials)):
    """近隣の医療機関を検索"""
    try:
        data = await request.json()
        
        # パラメータの取得
        address = data.get('address')
        radius = data.get('radius', 3000)
        department_types = data.get('department_types', [])
        
        if not address:
            return JSONResponse(
                status_code=400,
                content={"error": "address is required"}
            )
        
        # Google Maps サービスのインスタンスを作成
        from .services.google_maps_service import GoogleMapsService
        google_maps = GoogleMapsService()
        
        # 検索を実行
        result = await google_maps.search_nearby_clinics(
            location=address,
            radius=radius,
            department_types=department_types,
            limit=20
        )
        
        if result.get("error"):
            return JSONResponse(
                status_code=500,
                content={"error": result["error"]}
            )
        
        return JSONResponse(content=result)
        
    except Exception as e:
        print(f"Error searching nearby clinics: {e}")
        traceback.print_exc()
        return JSONResponse(
            status_code=500,
            content={"error": f"検索中にエラーが発生しました: {str(e)}"}
        )
